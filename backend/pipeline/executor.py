"""
非同步子 process 執行器。

使用 asyncio.create_subprocess_shell，即時串流輸出到 logger，
支援 timeout 強制終止。

Skill 模式：LLM 解讀自然語言任務描述，自主撰寫並執行程式碼完成任務。
"""
import asyncio
import base64
import hashlib
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage

# 模組層 logger，給沒有 per-step logger 的輔助函式使用（例如沙盒路由）
log = logging.getLogger(__name__)

from config import GROQ_API_KEY, GROQ_MODEL_MAIN

SKILL_TOOL_TIMEOUT = 60          # 預設值：給沒透過 execute_step_with_skill 流程的呼叫者用
SKILL_TOOL_TIMEOUT_MAX = 300     # 動態 tool timeout 上限(2026-05-24 從 180→300:node build_pptx 等任務常超 180s、放寬避免 SIGTERM)
SKILL_MAX_ITERATIONS = 20  # 互動式 skill(python-cli-extractor:問 A-B/subcommand/參數 + 掃描大專案找函式簽名 + GUI 解耦 + retry 補套件 + 重跑)耗 iter 多;較不聰明的模型(Gemma 等)還會空回覆 / 多繞、15 仍會撞牆。20 留足容錯。純運算 skill 照樣 4-5 輪結束

# 連續 N 輪 LLM 口頭說「完成」但沒下 <tool>done</tool> → 強制 done 收尾
# 解決 Sonnet / GPT 等強模型常見的「死循環式 verify」— 寫好檔案後一直確認、不打標準 tag
_DONE_KEYWORDS = (
    "任務已完成", "已完成任務", "已成功完成", "✅ 任務", "✅ 已完成", "✅ 完成",
    "task complete", "task is complete", "successfully completed", "all done",
    "已成功", "成功完成", "已順利完成", "完整完成",
)


def _looks_like_done(reply: str) -> bool:
    """偵測 LLM 口頭表示完成但沒打 <tool>done</tool> 標準 tag。"""
    if not reply:
        return False
    r = reply.lower()
    return any(k.lower() in r for k in _DONE_KEYWORDS)


# 第 N 輪起把更早的 tool 結果摺成短摘要、防 context 雪崩
# 保留最近 KEEP_RECENT_FULL 輪的完整結果、更早的截首尾各 PREVIEW_CHARS 字
SKILL_CONTEXT_KEEP_RECENT_FULL = 3
SKILL_CONTEXT_PREVIEW_CHARS = 200

# Prose-before-tool 守門(backport 自 subagent_runner、task #160)。
# native FC 下「prose before tool」= AIMessage.content（伴隨 tool_calls 的長篇文字）。
# 共用 subagent 的環境變數 SUBAGENT_PROSE_CAP_MODE(soft/enforce/off、預設 soft),
# 不另開旗標。實測 Sonnet 4.6 在 skill loop 單輪寫 5 萬字 parser → 拖斷連線 → 燒額度。
_SKILL_PROSE_BEFORE_TOOL_THRESHOLD = 3000   # 跟 subagent 同閾值
_SKILL_PROSE_BEFORE_TOOL_LIMIT = 2          # 連 N 輪違規才中止(enforce)
_SKILL_PROSE_CAP_MODE = (os.environ.get("SUBAGENT_PROSE_CAP_MODE", "soft").strip().lower())
if _SKILL_PROSE_CAP_MODE not in ("soft", "enforce", "off"):
    _SKILL_PROSE_CAP_MODE = "soft"


# ── 錯誤分類 + 對症提示 ───────────────────────────────────────────────────────
# 連續 2 次同類錯誤時、注入具體的恢復策略給 LLM(取代「換策略」這種空話)
# 每條 (kind, regex, hint),kind 是用來 dedup 計數的、regex 比對 stderr/tool_result
import re as _re
_ERROR_HINTS: list[tuple[str, "_re.Pattern[str]", str]] = [
    ("syntax_triple_quote",
     _re.compile(r"unterminated.*triple-quoted|EOL while scanning string|unterminated string literal", _re.IGNORECASE),
     '你用 Python `"""..."""` 包大段 JS/HTML/SQL 容易破。改成 `Path(p).write_text(content, encoding="utf-8")` '
     '把要寫的內容當一般字串、別當 Python source code 嵌入。或者把長字串拆成多個 .write() append。'),
    ("syntax_general",
     _re.compile(r"SyntaxError|IndentationError", _re.IGNORECASE),
     'Python 語法錯誤。常見原因:(1) f-string 內含 `{` `}` 字面值要寫 `{{` `}}`;(2) 中文標點 `,` `:` `:` `;` 被誤用;'
     '(3) 縮排混 tab/space。建議把這次 run_python 拆 < 100 行分批跑、好定位問題。'),
    ("attr_str_no_method",
     _re.compile(r"AttributeError:\s*['\"]?str['\"]?\s*object has no attribute"),
     '「str object has no attribute」幾乎都是 pickle 反序列化壞了(class 找不到變成 str)。**改用 JSON**:'
     '`json.dump(obj, f, default=str)` 寫、`json.load(f)` 讀。不要用 pickle 保存 LLM/第三方套件的物件。'),
    ("missing_module",
     _re.compile(r"(ModuleNotFoundError|ImportError):\s*No module named\s+['\"]?([\w\.]+)"),
     '套件未安裝。**不要 pip install、也不要 try/except 繞**,直接呼叫 '
     '`done(success=false, missing_packages=["套件名"])`。系統會幫你安裝後自動重跑這步。'),
    # docker_chdir 必須在 file_not_found 前面、否則「No such file or directory」會被 file_not_found 先匹配
    ("docker_chdir",
     _re.compile(r"OCI runtime exec failed|chdir.*no such file or directory", _re.IGNORECASE),
     'Docker 容器看不到這條路徑、不在 bind mount 範圍。容器只掛了 PROJECT_DIR + .agents。'
     '改寫到 `PROJECT_DIR/ai_output/` 內、不要寫 `~/ai_output/` 或其他絕對路徑。'),
    ("file_not_found",
     _re.compile(r"FileNotFoundError|No such file or directory", _re.IGNORECASE),
     '找不到檔案。檢查:(1) 路徑含 `~` 沒展開 → `Path(p).expanduser()`;'
     '(2) 寫檔前 `Path(p).parent.mkdir(parents=True, exist_ok=True)`;'
     '(3) 上一步是否真的有輸出該檔(先 `Path(p).exists()` 驗證)。'),
    ("permission_denied",
     _re.compile(r"PermissionError|Permission denied", _re.IGNORECASE),
     '檔案權限被拒。Windows 上多半是檔案被 Excel/PDF Reader/編輯器鎖住。'
     '建議:(1) 寫到不同檔名(加 timestamp 後綴);(2) 或在 done 訊息提示使用者關閉該檔。'),
    ("json_decode",
     _re.compile(r"json\.(decoder\.)?JSONDecodeError|Expecting value"),
     'JSON 解析失敗。先 `print(repr(s[:300]))` 看實際拿到的字串長怎樣 — 可能含 BOM、HTML、'
     '純文字或裝飾字元(```json 開頭)。確認是合法 JSON 再 parse。'),
    ("timeout",
     _re.compile(r"asyncio.*TimeoutError|timeout after \d+s", _re.IGNORECASE),
     '單次工具超時。把任務拆小:(1) 先用 3-5 筆資料測試邏輯;'
     '(2) 拆成多次 run_python、每次 <30s;(3) 用 async/concurrency 平行(若是 IO bound)。'),
]


def _classify_tool_error(tool_result: str) -> Optional[str]:
    """回傳錯誤 kind(用於 dedup 計數);找不到對應規則回 None。"""
    if not tool_result:
        return None
    for kind, pattern, _hint in _ERROR_HINTS:
        if pattern.search(tool_result):
            return kind
    return None


def _hint_for_error_kind(kind: str) -> Optional[str]:
    for k, _p, hint in _ERROR_HINTS:
        if k == kind:
            return hint
    return None


# LLM agent 自己寫的程式碼壓縮(對 AIMessage 用)— 防 pptx skill 那種「寫 15KB JS wrap」case
# 跟 _compact_old_tool_results 對稱:tool 結果第 4 輪起壓、LLM 自己寫的程式碼也該第 4 輪起壓
SKILL_AI_INPUT_COMPRESS_THRESHOLD = 3000  # AIMessage 內 <input>...</input> 段超過此長度才壓
SKILL_AI_INPUT_PREVIEW_CHARS = 300        # 壓縮後保留 head + tail 各 N 字


def _compact_old_llm_input_blocks(messages: list) -> int:
    """壓縮舊 AIMessage 內的 `<input>...</input>` 程式碼段。回傳壓縮的訊息數。

    對應場景:pptx skill 需要 LLM 寫 15KB+ Python wrap JS;這種「LLM 自己寫的巨碼」
    在歷史輪數會一直留著、context 5 輪後爆。

    規則:
    - 找所有 AIMessage(LLM 回覆)
    - 保留最近 SKILL_CONTEXT_KEEP_RECENT_FULL 個完整
    - 更早的:用 regex 抓 `<input>...</input>` 段,內容 > THRESHOLD 就截成 head+tail
    - 已壓過(含 `(舊輪程式摘要)` 標記)的不重複壓
    """
    from langchain_core.messages import AIMessage as _AI
    ai_indices = [i for i, m in enumerate(messages) if isinstance(m, _AI) and isinstance(m.content, str)]
    if len(ai_indices) <= SKILL_CONTEXT_KEEP_RECENT_FULL:
        return 0
    # 要壓縮的:除了最後 KEEP_RECENT_FULL 個之外的
    to_compact = ai_indices[:-SKILL_CONTEXT_KEEP_RECENT_FULL]
    compacted = 0
    # 抓 <input>...</input>(非貪婪、跨行)、找 long 段壓縮
    _pat = _re.compile(r"(<input>)(.*?)(</input>)", _re.DOTALL)

    def _shrink(m: "_re.Match[str]") -> str:
        body = m.group(2)
        if len(body) <= SKILL_AI_INPUT_COMPRESS_THRESHOLD:
            return m.group(0)
        if "(舊輪程式摘要)" in body:
            return m.group(0)
        head = body[:SKILL_AI_INPUT_PREVIEW_CHARS]
        tail = body[-SKILL_AI_INPUT_PREVIEW_CHARS:]
        new_body = (
            f"\n(舊輪程式摘要、原長 {len(body)} 字)\n"
            f"--- 前 {SKILL_AI_INPUT_PREVIEW_CHARS} 字 ---\n{head}\n"
            f"--- 後 {SKILL_AI_INPUT_PREVIEW_CHARS} 字 ---\n{tail}\n"
        )
        return m.group(1) + new_body + m.group(3)

    for idx in to_compact:
        content = messages[idx].content
        if "<input>" not in content:
            continue
        new_content, n = _pat.subn(_shrink, content)
        if new_content != content:
            messages[idx] = _AI(content=new_content)
            compacted += 1
    return compacted


def _compact_old_tool_results(messages: list) -> int:
    """壓縮舊 tool 結果防 context 膨脹。回傳壓縮的訊息數。

    規則:
    - 找所有 `[工具結果 —` 開頭的 HumanMessage
    - 保留最後 KEEP_RECENT_FULL 個完整、更早的截成 head + tail preview
    - 已經是 [摘要] 的不重複壓
    """
    from langchain_core.messages import HumanMessage as _HM
    # 找出所有「工具結果」訊息的 index
    tool_msg_indices = []
    for i, m in enumerate(messages):
        if isinstance(m, _HM) and isinstance(m.content, str) and m.content.startswith("[工具結果 —"):
            tool_msg_indices.append(i)

    if len(tool_msg_indices) <= SKILL_CONTEXT_KEEP_RECENT_FULL:
        return 0  # 還不需要壓縮

    # 要壓縮的:除了最後 KEEP_RECENT_FULL 個之外的、且還沒壓過的
    to_compact = tool_msg_indices[:-SKILL_CONTEXT_KEEP_RECENT_FULL]
    compacted = 0
    for idx in to_compact:
        content = messages[idx].content
        if "(舊輪摘要)" in content:
            continue  # 已壓過
        # 抽 tool name(第一行的 [工具結果 — X])
        first_nl = content.find("\n")
        header = content[:first_nl] if first_nl > 0 else content[:80]
        body = content[first_nl + 1:] if first_nl > 0 else ""
        if len(body) <= SKILL_CONTEXT_PREVIEW_CHARS * 2 + 50:
            continue  # 短的不必壓
        head = body[:SKILL_CONTEXT_PREVIEW_CHARS]
        tail = body[-SKILL_CONTEXT_PREVIEW_CHARS:]
        new_content = (
            f"{header} (舊輪摘要、原長度 {len(body)} 字)\n"
            f"--- 前 {SKILL_CONTEXT_PREVIEW_CHARS} 字 ---\n{head}\n"
            f"--- 後 {SKILL_CONTEXT_PREVIEW_CHARS} 字 ---\n{tail}"
        )
        messages[idx] = _HM(content=new_content)
        compacted += 1
    return compacted


def _compute_tool_timeout(step_timeout: int) -> int:
    """從 step.timeout 推導「單次 run_python / run_shell 上限秒數」。
    使用者已用 step.timeout 標註過該步驟大概要多久、tool 上限自然該跟著放寬。

    公式(2026-05-24 寬鬆化):min(300, max(90, step.timeout // 3))
    - step.timeout=270(短任務)  → 90s
    - step.timeout=300(預設)    → 100s
    - step.timeout=600          → 200s
    - step.timeout=900+         → 300s(封頂)

    寬鬆化原因(取代 //5 60s 公式):
    - node build_pptx.js 生 9 slide 含複雜 shape 超 60s → SIGTERM、step fail
    - npm install 偶發 > 60s
    - 60s 上限對「複雜但合理」的單次 tool 太緊、user 體感卡住
    - 改 //3 仍留 2/3 budget 給 LLM 思考 + 多次 tool retry,不會吃光 step.timeout
    """
    return min(SKILL_TOOL_TIMEOUT_MAX, max(90, int(step_timeout) // 3))
ASK_USER_MAX = 6          # 一個 skill 節點最多 ask_user 次數（ask_mode ON 時取消）。互動式 skill（python-cli-extractor 要問 模式/A-B/subcommand/參數 ≥4 次）3 太緊跑不完、6 留容錯
ASK_USER_TIMEOUT = 3600   # 單次等待使用者回答的逾時（秒）

# ── web_search 成本 / context 保護 ─────────────────────────────────────────
# 兩段式設計（簡化自原本的 3-tier）：
#   OFF：輕量 — answer + URL 清單（~500 字）
#   ON： 完整 — answer + URL + 每則文章完整原文（~15000 字）
#        由 Tavily 端直接回完整內容（include_raw_content=True），Agent 不用自己寫爬蟲
WEB_SEARCH_MAX_PER_STEP = 8             # 單一 skill step 最多呼叫次數(深度研究要拆多子題各搜、5→8)
WEB_SEARCH_OUTPUT_CHAR_CAP_LIGHT = 2000 # OFF 模式：輕量硬上限
WEB_SEARCH_OUTPUT_CHAR_CAP_FULL = 30000 # ON 模式：完整內容硬上限(放寬 20K→30K 容納更多來源)
WEB_SEARCH_PER_RESULT_FULL_CHARS = 4000 # ON 模式：每則原文截斷長度(3000→4000)
WEB_SEARCH_TITLE_CHARS = 100            # Title 顯示最大長度


# ── ask_user 進行中的問題：run_id -> {question, options, context, event, answer} ──
# In-memory：後端重啟會清空，使用者需重新觸發
_pending_questions: dict[str, dict] = {}


def deliver_ask_user_answer(run_id: str, answer: str) -> bool:
    """外部（resume_pipeline）呼叫：把答案送給正在等待的 skill agent。"""
    pending = _pending_questions.get(run_id)
    if not pending:
        return False
    pending["answer"] = answer
    pending["event"].set()
    return True


def get_pending_question(run_id: str) -> Optional[dict]:
    """查詢某 run 目前是否正在等 ask_user 答案。
    回 {question, options, context}(不含 event);沒在等則回 None。
    (2026-05-31 修:原本 pending 存在時 fall through 隱式回 None、害 GET /ask-user 永遠
     pending=False、前端/API 查不到問題、自動化 resume 全卡死。)"""
    pending = _pending_questions.get(run_id)
    if not pending:
        return None
    return {
        "question": pending.get("question"),
        "options": pending.get("options") or [],
        "context": pending.get("context") or "",
    }


# ── Phase B: ask_mode 命令分類攔截（in-memory pending state）──
# 正在等待的命令授權：run_id -> {category, label, preview, event, decision}
_pending_command_approvals: dict[str, dict] = {}


def deliver_command_approval(run_id: str, decision: str) -> bool:
    """外部（resume_pipeline）呼叫：把命令授權結果送給正在等待的 skill agent。
    decision: "allow" | "deny" | "hint"
    """
    pending = _pending_command_approvals.get(run_id)
    if not pending:
        return False
    pending["decision"] = decision
    pending["event"].set()
    return True


def get_pending_command_approval(run_id: str) -> Optional[dict]:
    """查詢某 run 目前是否在等命令授權。"""
    pending = _pending_command_approvals.get(run_id)
    if not pending:
        return None
    return {k: v for k, v in pending.items() if k != "event"}


# ── 命令分類規則（順序：先具體、後通用；first-match-wins）──
# 安全原則：判斷「是不是安裝行為」而非「是不是缺依賴」— LLM 主動裝套件也要攔
import re as _cmd_cls_re
_COMMAND_CLASSIFY_RULES = [
    # privileged: 提升權限（最高優先）
    (_cmd_cls_re.compile(r'\bsudo\b'), 'privileged', 'sudo 提升權限'),
    (_cmd_cls_re.compile(r'\brunas\b'), 'privileged', 'runas 提升權限'),

    # remote-exec: 從網路下載執行（次高優先）
    (_cmd_cls_re.compile(r'curl\s+[^|;]*\|\s*(?:sh|bash)\b'), 'remote-exec', '從網路下載並執行（curl | sh）'),
    (_cmd_cls_re.compile(r'wget\s+[^|;]*\|\s*(?:sh|bash)\b'), 'remote-exec', '從網路下載並執行（wget | sh）'),

    # install: 安裝/移除套件 — 須先於 subprocess / destructive 規則,優先匹配
    (_cmd_cls_re.compile(r'\bpip3?\s+install\b'), 'install', 'pip 安裝套件'),
    (_cmd_cls_re.compile(r'\bpip3?\s+uninstall\b'), 'install', 'pip 移除套件'),
    (_cmd_cls_re.compile(r'\bpython3?\s+-m\s+pip\s+install\b'), 'install', 'pip 安裝套件（python -m pip）'),
    (_cmd_cls_re.compile(r'\bnpm\s+(?:install|i)\b'), 'install', 'npm 安裝套件'),
    (_cmd_cls_re.compile(r'\bpnpm\s+(?:add|install)\b'), 'install', 'pnpm 安裝套件'),
    (_cmd_cls_re.compile(r'\byarn\s+add\b'), 'install', 'yarn 安裝套件'),
    (_cmd_cls_re.compile(r'\bpoetry\s+add\b'), 'install', 'poetry 新增套件'),
    (_cmd_cls_re.compile(r'\bconda\s+install\b'), 'install', 'conda 安裝套件'),
    (_cmd_cls_re.compile(r'\bapt(?:-get)?\s+install\b'), 'install', 'apt 安裝系統套件'),
    (_cmd_cls_re.compile(r'\bbrew\s+install\b'), 'install', 'brew 安裝套件'),
    (_cmd_cls_re.compile(r'\bInstall-Module\b'), 'install', 'PowerShell 安裝模組'),
    (_cmd_cls_re.compile(r'\bwinget\s+install\b'), 'install', 'winget 安裝程式'),
    # subprocess.run(["pip", "install", ...]) 等用 subprocess 包裝的安裝命令(優先列為 install)
    (_cmd_cls_re.compile(r"""(?:subprocess\.(?:run|call|Popen|check_output|check_call)|os\.system)\s*\(\s*['"\[]+\s*['"]?(?:pip3?|npm|pnpm|yarn|poetry|conda|apt|apt-get|brew|winget)['"]?\s*[,'"\]\s]+(?:install|add|i\b)"""), 'install', 'subprocess/os.system 安裝套件'),

    # destructive: 刪除檔案/目錄
    (_cmd_cls_re.compile(r'\brm\s+(?:-[rRfF]+\s+|--recursive\s+|--force\s+)'), 'destructive', '刪除目錄/檔案（rm -r/-f）'),
    (_cmd_cls_re.compile(r'shutil\.rmtree\b'), 'destructive', '刪除目錄（shutil.rmtree）'),
    (_cmd_cls_re.compile(r'os\.remove\b|os\.unlink\b'), 'destructive', '刪除檔案（os.remove / os.unlink）'),
    (_cmd_cls_re.compile(r'\b\w+\.unlink\(\)|\bPath\([^)]+\)\.unlink'), 'destructive', '刪除檔案（Path.unlink）'),

    # subprocess: 執行外部命令(排除 python 自呼)
    (_cmd_cls_re.compile(r"subprocess\.(?:run|call|Popen|check_output|check_call)\s*\(\s*\[\s*['\"](?!python)[^'\"]+['\"]"), 'subprocess', '透過 subprocess 執行外部命令'),
]


def classify_command(code: str) -> Optional[tuple[str, str, str]]:
    """偵測 ask_mode 下需要強制詢問的敏感命令。
    回 (category, label, preview) 或 None。
    """
    if not code:
        return None
    for pattern, cat, label in _COMMAND_CLASSIFY_RULES:
        m = pattern.search(code)
        if m:
            line_start = code.rfind('\n', 0, m.start()) + 1
            line_end = code.find('\n', m.end())
            line = code[line_start:line_end if line_end > 0 else len(code)].strip()
            preview = line[:140] + ('...' if len(line) > 140 else '')
            return (cat, label, preview)
    return None
    return {
        "question": pending["question"],
        "options": pending["options"],
        "context": pending["context"],
    }

# ── Per-run subprocess tracking（for immediate abort）─────────────────────────
import threading

_proc_lock = threading.Lock()
_running_procs: dict[str, list] = {}  # run_id → list of (subprocess.Popen | asyncio.subprocess.Process)


def register_proc(run_id: str, proc):
    """註冊一個正在執行的子進程，供 abort 時立即 kill"""
    with _proc_lock:
        _running_procs.setdefault(run_id, []).append(proc)


def unregister_proc(run_id: str, proc):
    """反註冊子進程"""
    with _proc_lock:
        if run_id in _running_procs:
            try:
                _running_procs[run_id].remove(proc)
            except ValueError:
                pass
            if not _running_procs[run_id]:
                del _running_procs[run_id]


def kill_run_processes(run_id: str):
    """立即終止指定 run 的所有子進程,連同 process tree(防 cmd.exe wrapper 殺掉、Python 變孤兒)。
    Windows 經典問題:create_subprocess_shell 透過 cmd.exe 開 Python,proc.kill() 只殺 cmd.exe,
    Python 子進程繼續活著(尤其有 GUI 的時候 GUI window 留著)。用 psutil 走完整 process tree。
    """
    with _proc_lock:
        procs = _running_procs.pop(run_id, [])
    if not procs:
        return
    try:
        import psutil as _psutil
    except ImportError:
        _psutil = None
    killed = 0
    for proc in procs:
        pid = getattr(proc, "pid", None)
        if pid is None:
            continue
        # 用 psutil 找 children + kill 整棵樹
        if _psutil is not None:
            try:
                parent = _psutil.Process(pid)
                children = parent.children(recursive=True)
                for child in children:
                    try:
                        child.kill()
                        killed += 1
                    except (_psutil.NoSuchProcess, _psutil.AccessDenied):
                        pass
                try:
                    parent.kill()
                    killed += 1
                except (_psutil.NoSuchProcess, _psutil.AccessDenied):
                    pass
                continue
            except (_psutil.NoSuchProcess, _psutil.AccessDenied):
                pass
            except Exception:
                pass
        # psutil 不在 / 失敗時的 fallback:只 kill 直接 proc
        try:
            proc.kill()
            killed += 1
        except (ProcessLookupError, OSError):
            pass
    if killed > 0:
        try:
            import logging as _logging
            _logging.getLogger("pipeline").info(f"🧹 kill_run_processes({run_id}):終結 {killed} 個進程(含子樹)")
        except Exception:
            pass

def _build_clean_success_stdout(all_stdout: list[str], done_marker_prefix: str) -> str:
    """Skill / Outlook agent 成功時，從 all_stdout 中只取最後一個工具結果 + 完成訊息。

    試錯的 Traceback 留在 .log 檔即可，validator 看到就會誤判 failed
    （LLM 看到 stderr 推論「步驟失敗」，即使 exit_code=0、檔案已正確產出）。
    這個誤判導致 retry → 整個 agent 重跑 5 個 iter → 工作流變慢數倍。"""
    if not all_stdout:
        return ""
    done_idx = None
    for i in range(len(all_stdout) - 1, -1, -1):
        if all_stdout[i].startswith(done_marker_prefix):
            done_idx = i
            break
    if done_idx is None:
        return "\n".join(all_stdout)
    start = max(0, done_idx - 1)
    return "\n".join(all_stdout[start:])


# Skill 模式需要的核心套件（探測用，從 skill_packages.txt 讀取）
def _load_skill_required_pkgs() -> tuple[str, ...]:
    pkg_file = Path(__file__).parent.parent / "skill_packages.txt"
    if pkg_file.exists():
        lines = pkg_file.read_text(encoding="utf-8").splitlines()
        pkgs = [l.strip() for l in lines if l.strip() and not l.strip().startswith("#")]
        if pkgs:
            return tuple(pkgs[:3])  # 取前 3 個作為探測用
    return ("matplotlib", "pandas", "openpyxl")

_SKILL_REQUIRED_PKGS = _load_skill_required_pkgs()


def _detect_python_interpreter() -> str:
    """
    跨平台偵測最適合 Skill 模式的 Python 直譯器：
    1. 優先鎖定專案目錄下的 .venv (確保 AI 能看到 UI 安裝的套件)
    2. 其次使用環境變數 SKILL_PYTHON 指定的路徑
    3.Fallback 到系統路徑或其他位置
    """
    import sys
    from pathlib import Path
    
    # 強制優先檢查專案內的 .venv (backend/.venv)
    proj_venv = Path(__file__).parent.parent / ".venv"
    if os.name == "nt":
        venv_exe = proj_venv / "Scripts" / "python.exe"
    else:
        venv_exe = proj_venv / "bin" / "python"
        
    if venv_exe.exists():
        return str(venv_exe.absolute())

    override = os.getenv("SKILL_PYTHON")
    if override and Path(override).exists():
        return override

    candidates: list[str] = []
    is_windows = os.name == "nt"
    # Windows: python, py.exe, python.exe；Unix: python3, python
    probe_names = ("python", "py", "python.exe", "py.exe") if is_windows else ("python3", "python")
    for name in probe_names:
        p = shutil.which(name)
        if p and p not in candidates:
            candidates.append(p)
    # Unix 常見路徑（Windows 會自動 skip 因為 os.path.exists 回 False）
    if not is_windows:
        for p in ("/usr/bin/python3", "/usr/local/bin/python3", "/opt/homebrew/bin/python3"):
            if os.path.exists(p) and p not in candidates:
                candidates.append(p)
    if sys.executable and sys.executable not in candidates:
        candidates.append(sys.executable)

    test_code = "import " + ", ".join(_SKILL_REQUIRED_PKGS)
    for py in candidates:
        try:
            r = subprocess.run(
                [py, "-c", test_code],
                capture_output=True, text=True, timeout=5,
            )
            if r.returncode == 0:
                return py
        except Exception:
            continue

    # 都不完整 → 退回第一個可用的；使用者會看到 ModuleNotFoundError，可自行 pip install
    if candidates:
        return candidates[0]
    return "python" if is_windows else "python3"


_SKILL_PYTHON = _detect_python_interpreter()
# Groq Free tier: 30 RPM → 每次 LLM 呼叫間隔至少 2 秒
SKILL_REQUEST_INTERVAL = 2.0
# 每 N 次 LLM 呼叫後強制冷卻（避免撞 TPM 上限）
SKILL_COOLDOWN_EVERY = 14
SKILL_COOLDOWN_SECONDS = 60


def _clean_env() -> dict:
    """移除 venv 對 PATH 的影響，並把 _SKILL_PYTHON 的目錄插到 PATH 最前面，
    確保 subprocess 內的 `python`/`python3` 都解析到有必要套件的 interpreter。
    """
    env = os.environ.copy()
    venv = env.pop("VIRTUAL_ENV", None)
    env.pop("PYTHONHOME", None)
    paths = env.get("PATH", "").split(os.pathsep)
    if venv:
        venv_bin = os.path.join(venv, "Scripts" if os.name == "nt" else "bin")
        paths = [p for p in paths if p != venv_bin]
    # 把 _SKILL_PYTHON 所在目錄放到 PATH 最前面
    global _SKILL_PYTHON
    skill_py = globals().get("_SKILL_PYTHON")
    if skill_py:
        skill_dir = os.path.dirname(skill_py)
        if skill_dir:
            paths = [p for p in paths if p != skill_dir]
            paths.insert(0, skill_dir)
    env["PATH"] = os.pathsep.join(paths)
    # 強制 stdout/stderr 用 UTF-8 編碼，避免 Windows cp1252/cp950 遇到中文 print() 炸出 UnicodeEncodeError
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"  # Python 3.7+ UTF-8 mode 全域啟用
    return env


def _script_env() -> dict:
    """script 節點專用環境(刻意與 V5 venv 脫鉤)。

    設計決策(2026-06-01):script 節點未勾「使用虛擬環境」時,裸 `python` 走
    **系統全域** Python,而非 V5 後端 venv —— 避免使用者腳本的依賴污染編排器自己的
    執行環境。缺依賴就 loud fail(由使用者改用專案自帶 venv 解決)。
    與 `_clean_env()`(skill/AI code 用、會導向 _SKILL_PYTHON)相反:這裡**防禦性移除**
    V5 venv 的 bin 目錄與 active VIRTUAL_ENV 對 PATH 的影響,確保裸 `python` 不會落到 V5 venv。
    """
    env = os.environ.copy()
    venv = env.pop("VIRTUAL_ENV", None)
    env.pop("PYTHONHOME", None)
    paths = env.get("PATH", "").split(os.pathsep)
    drop = []
    if venv:
        drop.append(os.path.join(venv, "Scripts" if os.name == "nt" else "bin"))
    skill_py = globals().get("_SKILL_PYTHON")
    if skill_py:
        d = os.path.dirname(skill_py)
        if d:
            drop.append(d)
    if drop:
        dropn = {os.path.normcase(os.path.normpath(p)) for p in drop}
        paths = [p for p in paths if os.path.normcase(os.path.normpath(p)) not in dropn]
    env["PATH"] = os.pathsep.join(paths)
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"
    return env


import re as _re


_PY_CMD_RE = _re.compile(
    r'^(\s*)(python3|python|py)(\.exe)?(\s|$)',
    _re.IGNORECASE if os.name == "nt" else 0,
)


def _quote_path(path: str) -> str:
    """跨平台為含空格的路徑加引號。"""
    if os.name == "nt":
        return f'"{path}"' if (" " in path or "\t" in path) else path
    import shlex as _shlex
    return _shlex.quote(path)


def _rewrite_python_cmd(command: str) -> str:
    """把指令開頭的 python / python3 / py 換成 _SKILL_PYTHON（驗證過套件可用的 interpreter）。

    - 跨平台：Windows 用 py.exe/python.exe、Unix 用 python3/python
    - 只改最前面那顆，不動 pipe、&&、; 之後的
    - 不 re-tokenize 整個命令（避免反斜線路徑被破壞）
    """
    if not _SKILL_PYTHON:
        return command
    m = _PY_CMD_RE.match(command)
    if not m:
        return command
    prefix = m.group(1)
    trailing = m.group(4)
    rest = command[m.end():]
    return f"{prefix}{_quote_path(_SKILL_PYTHON)}{trailing}{rest}"


# 整條命令 = `<python> -c "<code>"`(允許前後空白、code 可跨行)
_PY_DASH_C_RE = _re.compile(r'^\s*(\S.*?)\s+-c\s+(["\'])(.*)\2\s*$', _re.DOTALL)
_PY_INTERP_RE = _re.compile(
    r'(?:^|[\\/ "])(?:python3?|py)(?:\.exe)?"?$', _re.IGNORECASE
)


def _maybe_extract_multiline_python_c(command: str):
    """多行 `python -c "..."` 在 Windows cmd.exe 會被換行切斷 → 只剩第一行送進
    cmd、其餘靜默丟掉,結果是 exit 0 卻什麼都沒做(過夜測試 BUG 5)。

    偵測到「整條命令就是一個跨行的 python -c」就把程式碼抽出寫進暫存 .py、
    改成 `<python> "<暫存檔>"`。單行 -c 不受影響、照舊。

    回傳 (改寫後命令, 暫存檔路徑 or None)。暫存檔由 caller 跑完負責刪。
    """
    m = _PY_DASH_C_RE.match(command)
    if not m:
        return command, None
    interpreter, _q, code = m.group(1), m.group(2), m.group(3)
    if "\n" not in code:
        return command, None  # 單行 -c shell 接得住、不動
    if not _PY_INTERP_RE.search(interpreter.strip()):
        return command, None  # 開頭不是 python、避免誤判其他含 -c 的命令
    import tempfile
    fd, path = tempfile.mkstemp(suffix=".py", prefix="step_inline_")
    with os.fdopen(fd, "w", encoding="utf-8") as f:
        f.write(code)
    return f'{interpreter} {_quote_path(path)}', path


@dataclass
class ExecResult:
    exit_code: int
    stdout: str
    stderr: str
    pending_recipe: Optional[dict] = None  # 延遲儲存的 recipe 資料
    missing_packages: list = None          # LLM 回報缺少的套件（供 runner 產生安裝建議）
    # skill agent 主動呼叫 done(success=false) — 這是「明確結論」非 crash / exception。
    # runner 看到此旗標時跳過 step retry、直接進 awaiting_human 把 summary 給使用者
    # （重試只會重複同樣結論、浪費一輪 LLM、延後使用者看到結論）。
    agent_concluded_fail: bool = False
    # Trace / token tracking — skill / outlook agent loop 累計填入；shell 路徑保持 default 空。
    # 結構: {"input_tokens": int, "output_tokens": int, "total_tokens": int}
    token_usage: dict = field(default_factory=dict)
    # 結構: list[{"name": str, "input_preview": str, "result_preview": str}]
    tool_calls: list = field(default_factory=list)


async def execute_step(
    command: str,
    timeout: int,
    logger: logging.Logger,
    step_name: str,
    run_id: str = "",
    working_dir: Optional[str] = None,
    background: bool = False,
    ready_after_seconds: int = 0,
) -> ExecResult:
    """
    執行 shell 命令，串流輸出到 logger，回傳完整結果。

    Args:
        command:     shell 命令字串
        timeout:     最大執行秒數
        logger:      file logger（記錄完整輸出）
        step_name:   用於 log 標籤
        run_id:      pipeline run id（用於立即中止追蹤）
        working_dir: 當前工作目錄（會注入 PIPELINE_OUTPUT_DIR）

    Returns:
        ExecResult(exit_code, stdout, stderr)
    """
    # 設計決策(2026-06-01):script 節點**不**把裸 `python` 改寫成 V5 venv interpreter。
    # 裸 `python` 走系統全域(見 _script_env);要用特定 venv 請在 batch 用該 venv 的
    # 絕對路徑(UI「使用虛擬環境」勾選會自動填)。避免污染編排器自己的 venv。
    # （skill / AI 生成 code 仍走 _SKILL_PYTHON,不受此影響。）

    # 多行 `python -c "..."` 在 Windows cmd 會被換行切斷 → 改寫成暫存 .py 執行
    command, _inline_tmp = _maybe_extract_multiline_python_c(command)
    if _inline_tmp:
        logger.info(f"[{step_name}] 偵測到多行 python -c、已改寫成暫存腳本 {_inline_tmp}")

    logger.info(f"[{step_name}] ▶ 開始執行：{command}")

    stdout_lines: list[str] = []
    stderr_lines: list[str] = []

    # 準備環境變數(script 節點專用:與 V5 venv 脫鉤、裸 python 走系統全域)
    env = _script_env()
    cwd_arg: Optional[str] = None
    if working_dir:
        # 強制將工作目錄注入環境變數,供 stage 系列腳本主動讀取
        env["PIPELINE_OUTPUT_DIR"] = str(Path(working_dir).absolute())
        # 把 subprocess CWD 設成 workflow dir
        # → 一般使用者寫的 Python 工具就算用 open("x.csv") / Path("x.csv").write_text(...)
        # 也會落在 workflow 資料夾、snapshot diff 抓得到、下游 {{ steps.X.output.path }} 自動代入
        cwd_arg = str(Path(working_dir).absolute())

    try:
        proc = await asyncio.create_subprocess_shell(
            command,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            env=env,
            cwd=cwd_arg,
        )
        if run_id:
            register_proc(run_id, proc)

        # ── 背景模式:不等 exit、給 daemon 一段時間 boot up 後直接回 success ──
        if background:
            if ready_after_seconds > 0:
                logger.info(f"[{step_name}] 🚀 背景模式:等 {ready_after_seconds}s 讓 daemon 啟動完成…")
                await asyncio.sleep(ready_after_seconds)
                # 啟動期間若 proc 已 exit 了、應該當作失敗
                if proc.returncode is not None:
                    rc = proc.returncode
                    logger.warning(f"[{step_name}] ⚠ 背景進程在 {ready_after_seconds}s 內已退出(exit={rc}),非預期 daemon 行為")
                    return ExecResult(
                        exit_code=rc if rc is not None else -1,
                        stdout="(背景進程提早退出)",
                        stderr=f"進程在 ready_after_seconds={ready_after_seconds} 內就 exit、應該設成非背景或檢查啟動失敗",
                    )
            else:
                logger.info(f"[{step_name}] 🚀 背景模式啟動、不等 exit、立即下一步")
            # 子程序留著、由 run 結束時統一 kill(register_proc 註冊過了)
            return ExecResult(
                exit_code=0,
                stdout=f"(背景啟動 OK pid={proc.pid})",
                stderr="",
            )

        async def _drain(stream: asyncio.StreamReader, buf: list[str], tag: str):
            while True:
                raw = await stream.readline()
                if not raw:
                    break
                line = raw.decode("utf-8", errors="replace").rstrip()
                buf.append(line)
                logger.debug(f"[{step_name}][{tag}] {line}")

        try:
            await asyncio.wait_for(
                asyncio.gather(
                    _drain(proc.stdout, stdout_lines, "out"),
                    _drain(proc.stderr, stderr_lines, "err"),
                    proc.wait(),
                ),
                timeout=timeout,
            )
        except asyncio.TimeoutError:
            try:
                proc.kill()
            except ProcessLookupError:
                pass
            logger.error(f"[{step_name}] ⏱ 執行超時（>{timeout}s），已強制終止")
            if run_id:
                unregister_proc(run_id, proc)
            return ExecResult(
                exit_code=-1,
                stdout="\n".join(stdout_lines),
                stderr=f"執行超時（>{timeout}s）",
            )

        if run_id:
            unregister_proc(run_id, proc)

        exit_code = proc.returncode if proc.returncode is not None else -99
        level = logging.INFO if exit_code == 0 else logging.WARNING
        logger.log(level, f"[{step_name}] ■ 結束，exit code: {exit_code}")

        return ExecResult(
            exit_code=exit_code,
            stdout="\n".join(stdout_lines),
            stderr="\n".join(stderr_lines),
        )

    except FileNotFoundError as e:
        logger.error(f"[{step_name}] 命令找不到：{e}")
        return ExecResult(exit_code=-2, stdout="", stderr=f"命令找不到：{e}")

    except Exception as e:
        logger.error(f"[{step_name}] 執行異常：{e}")
        return ExecResult(exit_code=-3, stdout="", stderr=str(e))

    finally:
        if _inline_tmp:
            try:
                os.unlink(_inline_tmp)
            except OSError:
                pass


# ── Skill 模式執行器 ─────────────────────────────────────────────────────────

_DANGEROUS_COMMANDS = {'rm', 'rmdir', 'del', 'format', 'mkfs', 'dd', 'kill', 'shutdown', 'reboot'}

# 雙模型 cache:key = (role, settings_signature)
_skill_llm_cache: dict[tuple[str, str], Any] = {}


def _get_skill_llm(role: str = "primary"):
    """依 role(primary/secondary)拿對應的 LLM。設定變動時 cache 失效。"""
    from settings import settings_signature
    from llm_factory import build_llm
    sig = settings_signature()
    key = (role, sig)
    if key not in _skill_llm_cache:
        _skill_llm_cache[key] = build_llm(temperature=0, role=role)
    return _skill_llm_cache[key]


def _skill_run_python(code: str, cwd: Optional[str] = None, run_id: str = "",
                      tool_timeout: int = SKILL_TOOL_TIMEOUT) -> str:
    """在 subprocess 中執行 Python 程式碼。"""
    # 截斷混入程式碼中的 <tool> 標籤（LLM 有時在 run_python 輸入末尾附加 <tool>done</tool>）
    tool_tag_pos = code.find('<tool>')
    if tool_tag_pos > 0:
        code = code[:tool_tag_pos].rstrip()
    # 注入 done / view_image / read_file 的 no-op stub，避免 LLM 把工具名當 Python 函式呼叫而崩潰
    # 另外抑制所有 warnings，避免 pandas FutureWarning 等雜訊污染 stderr 害 LLM 誤以為失敗
    # 第一行必須是 UTF-8 encoding 宣告（PEP 263）：即使我們用 UTF-8 寫檔，也保險讓 Python 明確識別
    preamble = (
        "# -*- coding: utf-8 -*-\n"
        "import warnings\n"
        "warnings.filterwarnings('ignore')\n"
        "def done(*args, **kwargs):\n"
        "    print('[info] done() is a tool, not a Python function - ignored in script context')\n"
        "def view_image(*args, **kwargs):\n"
        "    print('[info] view_image() is a tool, not a Python function - ignored')\n"
        "def read_file(*args, **kwargs):\n"
        "    print('[info] read_file() is a tool, not a Python function - ignored')\n"
    )
    code = preamble + code
    tmp_path = None
    proc = None
    try:
        # 必須明確指定 UTF-8 編碼，否則 Windows 會用系統 locale（cp950/cp1252）寫檔，
        # LLM 的程式碼只要含任何非該 locale 的字元（中文註解、em dash 等）就會產生
        # "Non-UTF-8 code starting with '\\xXX'" SyntaxError
        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
            f.write(code)
            tmp_path = f.name
        # 子程序強制用 UTF-8 I/O，避免 Windows cp950/cp1252 locale 把含中文的 Traceback
        # 解不出來 → stderr 被吃光 → LLM 收到 [exit code: 1] 卻沒錯誤訊息可改，無限重試
        child_env = _clean_env()
        child_env["PYTHONIOENCODING"] = "utf-8"
        child_env["PYTHONUTF8"] = "1"  # Python 3.7+ 強制 UTF-8 模式
        # skill_llm helper:host 模式也注入 SKILL_LLM_* + PYTHONPATH(Windows 路徑),
        # 讓 `from skill_llm import llm` 走系統當前模型、不必 import openai。
        try:
            from pipeline.skill_llm_env import runtime_env, HELPERS_DIR_WIN
            _renv = runtime_env()
            if _renv:
                child_env.update(_renv)
                _pp = str(HELPERS_DIR_WIN)
                child_env["PYTHONPATH"] = (
                    _pp + os.pathsep + child_env["PYTHONPATH"]
                    if child_env.get("PYTHONPATH") else _pp
                )
        except Exception:
            pass
        proc = subprocess.Popen(
            [_SKILL_PYTHON, tmp_path],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",  # 出現無法解碼的 byte 就用 U+FFFD 代替，不讓解碼錯誤吃掉訊息
            env=child_env,
            cwd=cwd,
        )
        if run_id:
            register_proc(run_id, proc)
        try:
            stdout, stderr = proc.communicate(timeout=tool_timeout)
        except subprocess.TimeoutExpired:
            proc.kill()
            proc.communicate()
            return f"[錯誤] Python 執行超時（>{tool_timeout}秒）"
        finally:
            if run_id and proc:
                unregister_proc(run_id, proc)
        output = ""
        if stdout:
            output += stdout
        if stderr:
            # 區分錯誤 vs 警告：exit code 0 + stderr 只有警告不該讓 LLM 以為失敗
            tag = "stderr" if proc.returncode != 0 else "warnings"
            output += f"\n[{tag}]\n{stderr}"
        if proc.returncode != 0:
            output += f"\n[exit code: {proc.returncode}]"
            # 保險：若 exit code 非零但 stdout / stderr 都空 → 明確告訴 LLM 捕捉不到錯誤，
            # 讓它改變策略（例如加 try/except 印 traceback、改寫 log 到檔案）而不是重送同一份程式
            if not stdout and not stderr:
                output += (
                    "\n[提示] 子程序非正常結束但沒捕捉到任何 stdout / stderr。"
                    "可能是非 UTF-8 位元組、C-level crash 或進程被殺。"
                    "請在程式碼外層包 try/except 印出 traceback 到 stdout，或改寫 log 到檔案排查。"
                )
        elif not stdout:
            # 成功執行但沒輸出 → 明確告訴 LLM 任務已完成，避免誤以為失敗
            output += "\n[執行成功，程式無 stdout 輸出]"
        return output.strip() or "(無輸出)"
    except Exception as e:
        return f"[錯誤] Python 執行失敗：{e}"
    finally:
        if tmp_path:
            Path(tmp_path).unlink(missing_ok=True)


def _is_pip_install_cmd(cmd: str) -> bool:
    """偵測 run_shell 命令是否在 pip install 套件(pip/pip3/python -m pip/uv pip + install)。
    擋下後強制走 missing_dependency 使用者確認 —— skill 不准自己裝套件繞過確認。
    只攔 install、不攔 pip list / pip show / pip --version 等查詢命令。"""
    import re
    if not cmd:
        return False
    return bool(re.search(
        r'(^|[\s;&|(])(pip3?|python3?\s+-m\s+pip|uv\s+pip)\s+install\b',
        cmd, re.IGNORECASE))


# ── 共用 helper:native skill / subagent loop 都用,避免保護 drift ───────────────
# (2026-06-01 重構:把「pip 安裝偵測 + 缺套件偵測」集中。涵蓋 run_shell command
#  與 run_python code 內(subprocess/os.system/pip.main/['pip','install'])兩種繞法。)
_PIP_IN_CODE_PATTERNS = [
    r'\bpip3?\s+install\b',                          # os.system("pip install x") / !pip install
    r'-m\s+pip\s+install\b',                         # python -m pip install
    r'["\']pip3?["\']\s*,\s*["\']install["\']',      # subprocess([... 'pip','install' ...])
    r'\bpip\.main\s*\(',                             # pip.main(['install',...])
    r'\bensurepip\b',                                # ensurepip
]
_PKG_STOPWORDS = {
    "pip", "pip3", "install", "python", "python3", "u", "upgrade", "user",
    "q", "quiet", "m", "no", "cache", "dir", "r", "import", "subprocess",
    "os", "system", "run", "check_call", "call", "popen", "main", "uv",
}


def _extract_pip_pkgs(text: str) -> list:
    """從 pip install 命令 / code 片段抽套件名(去 flags / 版本 / 路徑)。
    遇到不像套件名的 token(含 . ( 等 code 邊界)就停 —— 避免把 install 後整段 code 都當套件。"""
    import re
    if not text or "install" not in text:
        return []
    after = text.split("install", 1)[1]
    out = []
    for raw in re.split(r"[\s,]+", after):
        tok = raw.strip().strip("'\"[](){};:")
        if not tok:
            continue
        if tok.startswith("-"):   # flag(-q / -U / --no-cache-dir)→ 跳過、不中斷
            continue
        base = tok.split("==")[0].split(">")[0].split("<")[0].split("~")[0]
        if (re.fullmatch(r"[A-Za-z][\w\-]*", base)
                and base.lower() not in _PKG_STOPWORDS
                and not base.endswith((".txt", ".cfg", ".toml", ".py"))):
            if base not in out:
                out.append(base)
            # 原始 token 帶字串/命令結束標記(' " ) ;)→ run_python code 內的 pip、
            # 套件清單到此為止(後面是程式碼、別把 result / print 等 identifier 當套件)
            if any(c in raw for c in (")", "'", '"', ";")):
                break
        else:
            break   # 遇到不像套件的 token(含 . ( 等 → code 邊界)→ 套件清單結束
    return out[:5]


def detect_pip_install(tc_name: str, tc_args) -> list:
    """共用:偵測 run_shell command 或 run_python code 內的 pip install 企圖。
    回偵測到的套件名 list(caller 應攔截、轉 missing_dependency / steering);空 list = 放行。
    tc_args 接 dict 或 JSON 字串(native FC 兩種格式都處理)。"""
    import re
    if not isinstance(tc_args, dict):
        try:
            import json as _j
            tc_args = _j.loads(tc_args) if (isinstance(tc_args, str) and tc_args.strip().startswith("{")) else {}
        except Exception:
            tc_args = {}
    if tc_name == "run_shell":
        cmd = str(tc_args.get("command") or tc_args.get("input") or "")
        if _is_pip_install_cmd(cmd):
            return _extract_pip_pkgs(cmd) or ["(未知套件)"]
    elif tc_name == "run_python":
        code = str(tc_args.get("code") or tc_args.get("input") or "")
        if any(re.search(p, code, re.IGNORECASE) for p in _PIP_IN_CODE_PATTERNS):
            return _extract_pip_pkgs(code) or ["(未知套件)"]
    return []


# import 名 ≠ PyPI 安裝名 的知名套件。key=程式碼裡 import 的名字、value=pip install 的名字。
# 缺套件偵測抓到的是 import 名(ModuleNotFoundError 給的),直接拿去裝會裝錯/裝不到
# (例:pip install whisper 裝到的是另一個不相干的套件,真正要的是 openai-whisper)。
_IMPORT_TO_PKG = {
    "whisper": "openai-whisper",
    "cv2": "opencv-python",
    "PIL": "Pillow",
    "bs4": "beautifulsoup4",
    "sklearn": "scikit-learn",
    "skimage": "scikit-image",
    "yaml": "PyYAML",
    "docx": "python-docx",
    "pptx": "python-pptx",
    "fitz": "PyMuPDF",
    "dotenv": "python-dotenv",
    "dateutil": "python-dateutil",
    "git": "GitPython",
    "Crypto": "pycryptodome",
    "OpenSSL": "pyOpenSSL",
    "serial": "pyserial",
    "usb": "pyusb",
    "win32com": "pywin32",
    "win32api": "pywin32",
    "win32con": "pywin32",
    "speech_recognition": "SpeechRecognition",
    "telebot": "pyTelegramBotAPI",
    "Xlib": "python-xlib",
    "attr": "attrs",
    "jose": "python-jose",
    "magic": "python-magic",
    "slugify": "python-slugify",
}


def _resolve_pkg_name(import_name: str, task_description: str = "") -> str:
    """把 import 名解析成真實 pip 安裝名。優先序:
    1) 靜態別名表(知名 import≠package、零誤判)
    2) 任務描述線索(保守):任務文字裡有帶連字號/底線的 token、且其「最後一段」
       剛好等於 import 名 → 採用(openai-whisper 末段=whisper ✓;
       pandas-profiling 末段=profiling≠pandas → 不採用、避免亂猜)
    3) 退回 import 名本身
    """
    import re as _re
    base = (import_name or "").split(".")[0]
    if not base:
        return import_name
    # 1) 靜態表優先(最確定)
    if base in _IMPORT_TO_PKG:
        return _IMPORT_TO_PKG[base]
    # 2) 任務描述線索(嚴格末段比對、低誤判)
    if task_description:
        for tok in _re.findall(r"[A-Za-z][A-Za-z0-9_.-]+", task_description):
            t = tok.strip(" .,，。、)(」「")
            if ("-" in t or "_" in t) and t.lower() != base.lower():
                last = _re.split(r"[-_]", t)[-1]
                if last.lower() == base.lower():
                    return t
    # 3) 退回 import 名
    return base


def detect_missing_module(result: str, task_description: str = "") -> list:
    """共用:從 tool result 偵測 ModuleNotFoundError、回缺的套件名 list(或空)。
    回傳的是「真實 pip 安裝名」(已經過 import 名→package 名解析)、不是 import 名。"""
    import re
    if not result:
        return []
    m = re.search(r"ModuleNotFoundError: No module named ['\"]([\w.]+)['\"]", result)
    if not m:
        return []
    return [_resolve_pkg_name(m.group(1), task_description)]


def _skill_run_shell(cmd: str, cwd: Optional[str] = None, run_id: str = "",
                     tool_timeout: int = SKILL_TOOL_TIMEOUT) -> str:
    """執行 shell 命令。"""
    first_word = cmd.strip().split()[0] if cmd.strip() else ""
    if first_word in _DANGEROUS_COMMANDS:
        return f"[拒絕] 命令 '{first_word}' 被安全策略封鎖"
    # 把 python/python3/py 開頭的指令改用 _SKILL_PYTHON（有 pandas 等套件的 interpreter）
    cmd = _rewrite_python_cmd(cmd)
    proc = None
    try:
        # 同 run_python：強制 UTF-8 避免 Windows locale 吃掉含中文的 stderr
        shell_env = _clean_env()
        shell_env["PYTHONIOENCODING"] = "utf-8"
        shell_env["PYTHONUTF8"] = "1"
        proc = subprocess.Popen(
            cmd, shell=True,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            env=shell_env,
            cwd=cwd,
        )
        if run_id:
            register_proc(run_id, proc)
        try:
            stdout, stderr = proc.communicate(timeout=tool_timeout)
        except subprocess.TimeoutExpired:
            proc.kill()
            proc.communicate()
            return f"[錯誤] 命令執行超時（>{tool_timeout}秒）"
        finally:
            if run_id and proc:
                unregister_proc(run_id, proc)
        output = ""
        if stdout:
            output += stdout
        if stderr:
            output += f"\n[stderr]\n{stderr}"
        if proc.returncode != 0:
            output += f"\n[exit code: {proc.returncode}]"
        return output.strip()[:5000] or "(無輸出)"
    except Exception as e:
        return f"[錯誤] 命令執行失敗：{e}"


def _skill_read_file(path: str, max_lines: int = 100, offset: int = 0) -> str:
    """讀取檔案內容。

    offset：跳過開頭 N 行(用於分段讀大檔)。
    max_lines：本次最多讀幾行。
    兩者搭配 24KB byte 上限、讓大檔可被「分段讀完整」而非只看開頭。"""
    try:
        # 清理 LLM 常見的錯誤格式：read_file("path"), 引號, 空白
        cleaned = path.strip()
        import re as _re
        m = _re.match(r'read_file\(["\']?(.+?)["\']?\)\s*$', cleaned)
        if m:
            cleaned = m.group(1)
        cleaned = cleaned.strip().strip('"').strip("'")
        # 沙盒路徑 → Windows 路徑（同 view_image 的修補）：LLM 在沙盒裡跑時會給 /mnt/c/...，
        # 但 read_file 在 host 上跑、需要 Windows 路徑
        cleaned = _wsl_to_windows_path(cleaned)
        p = Path(cleaned).expanduser()
        if not p.exists():
            return f"[錯誤] 檔案不存在：{path}（解析後：{p}）"
        if p.is_dir():
            files = sorted(p.iterdir())[:30]
            listing = "\n".join(f"  {'📁' if f.is_dir() else '📄'} {f.name} ({f.stat().st_size:,} B)" for f in files)
            return f"目錄內容：\n{listing}"
        # 偵測二進制檔案，避免汙染 LLM context
        binary_exts = {'.xlsx', '.xls', '.docx', '.pptx', '.pdf', '.png', '.jpg', '.jpeg', '.gif', '.zip', '.gz', '.tar', '.pkl', '.npy', '.parquet'}
        if p.suffix.lower() in binary_exts:
            size = p.stat().st_size
            return (f"[提示] {p.name} 是二進制檔案（{size:,} bytes），無法用 read_file 讀取。\n"
                    f"請改用 run_python 搭配適當的套件讀取：\n"
                    f"- .xlsx/.xls → pandas.read_excel() 或 openpyxl\n"
                    f"- .docx → python-docx\n"
                    f"- .png/.jpg → PIL 或 view_image 工具\n"
                    f"- .pdf → PyPDF2")
        # byte 上限:防 100 行內但全是長行(minified / 資料檔)的檔案撐爆 context。
        # 實測 retry round LLM 讀 187KB 大檔 → 單輪 input 衝到 273K 字 → LLM 異常。
        _CHAR_BUDGET = 24000  # ≈ 8K token,單次讀取上限
        offset = max(0, int(offset or 0))
        with open(p, 'r', encoding='utf-8', errors='replace') as f:
            lines = []
            char_count = 0
            truncated_reason = ""
            next_offset = offset
            read_count = 0
            for i, line in enumerate(f):
                if i < offset:
                    continue  # 跳過 offset 之前的行
                if read_count >= max_lines:
                    truncated_reason = f"已達本次 {max_lines} 行上限"
                    break
                stripped = line.rstrip()
                char_count += len(stripped) + 1
                if char_count > _CHAR_BUDGET:
                    truncated_reason = f"已達 {_CHAR_BUDGET:,} 字元上限"
                    break
                lines.append(stripped)
                read_count += 1
                next_offset = i + 1
            header = f"(檔案 {p.name}、共 {p.stat().st_size:,} bytes" + (
                f"、本次從第 {offset + 1} 行起" if offset else "") + ")"
            if truncated_reason:
                lines.append(
                    f"... ({truncated_reason}、檔案還有後續內容。"
                    f"要接著讀下一段、用 read_file 帶 offset={next_offset}:"
                    f' {{"path": "{cleaned}", "offset": {next_offset}}};'
                    f"或用 grep 直接搜尋關鍵字)"
                )
        body = "\n".join(lines) or "(空檔案 / offset 已超過檔尾)"
        return f"{header}\n{body}"
    except Exception as e:
        return f"[錯誤] 讀取失敗：{e}"


# ── write_file / edit_file / grep / glob tools ─────────────────────────────
# 這 4 個 tool 解掉「LLM 必須用 Python 包外國語言」的結構性問題:
# - write_file:LLM 直接傳 content 寫檔、不用 Python `"""..."""` 包(解 triple-quote 雷)
# - edit_file:局部替換、不用整檔 read+rewrite
# - grep:搜尋字串、不用 subprocess 包 grep
# - glob:列檔、不用 Path.glob() in Python
# 跟 run_python 互補:run_python 仍是邏輯處理主力,新 tool 接管「LLM 已知做什麼」場景

def _skill_write_file(input_str: str) -> str:
    """write_file(path, content) — 兩種 input 格式:
    1. JSON(短內容、< 5KB 推薦):{"path": "...", "content": "..."}
    2. RAW MULTILINE(長內容、推薦):
       path: /absolute/path/file.js
       ---
       <raw content here, no escaping needed, can contain \\u and " freely>

    raw 格式 detection:input 第一行是 `path: ...` 且第二行(或之後)有 `---` 分隔線。
    避免 30K+ JS code 包 JSON 踩 escape 雷(\\u / 引號 / 換行)。"""
    try:
        stripped = (input_str or "").lstrip()

        # 偵測 raw multiline 格式:第一行 `path: ...` + `---` 分隔
        if stripped.lower().startswith("path:") and "\n---" in stripped:
            first_newline = stripped.find("\n")
            path_line = stripped[:first_newline].strip()
            path = path_line.split(":", 1)[1].strip()
            rest = stripped[first_newline + 1:]
            # 找第一個獨佔行的 `---`
            sep_match = re.search(r"^---\s*$", rest, flags=re.MULTILINE)
            if not sep_match:
                return ("[錯誤] raw 格式需要獨佔行 '---' 分隔 path 跟 content\n"
                        "範例:\n  path: /abs/path.js\n  ---\n  <content>")
            content = rest[sep_match.end():].lstrip("\n")
        else:
            # JSON 格式
            import json as _json
            data = _json.loads(input_str)
            path = data.get("path") or data.get("file_path") or ""
            content = data.get("content")
            if content is None:
                return "[錯誤] write_file 需要 'content' 欄位"

        if not path:
            return "[錯誤] write_file 需要 'path' 欄位"
        # WSL→Windows 路徑轉換
        path = _wsl_to_windows_path(path)
        p = Path(path).expanduser()
        # 自動建 parent 目錄
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding="utf-8")
        size = p.stat().st_size
        return f"✓ 已寫入 {p}({size:,} bytes、{len(content.splitlines()) if isinstance(content, str) else '?'} 行)"
    except _json.JSONDecodeError as e:
        # 30K+ content 包 JSON 太容易踩 escape 雷、引導 LLM 改 raw multiline 或 run_python
        _len = len(input_str)
        if _len > 5000:
            return (
                f"[錯誤] JSON 解析失敗(content {_len:,} 字、太長易踩 escape 雷):{e}\n\n"
                f"⚠ 長內容(> 5KB)請改用 RAW MULTILINE 格式(無需 escape):\n"
                f"```\n"
                f"path: /絕對路徑/檔名.js\n"
                f"---\n"
                f"const x = `中文 \\u OK`;  // content 可自由用 \" \\ \\n \\u 任何字元、不必 escape\n"
                f"const y = \"也可以\";\n"
                f"```\n"
                f"或用 run_python + r''' heredoc 寫檔。"
            )
        return (f"[錯誤] JSON 格式錯誤:{e}\n"
                f"短內容 JSON:{{\"path\":\"x.js\",\"content\":\"...\"}}(content 內 \" 要 escape 為 \\\")\n"
                f"長內容請改 RAW MULTILINE:\n  path: /abs/path.js\n  ---\n  <content>")
    except Exception as e:
        return f"[錯誤] 寫檔失敗:{e.__class__.__name__}: {e}"


def _skill_edit_file(input_str: str) -> str:
    """edit_file(path, old, new) — JSON: {"path":"...", "old_text":"...", "new_text":"..."}.
    局部替換、不用整檔重寫。要求 old_text 唯一出現一次(避免歧義)。"""
    try:
        import json as _json
        data = _json.loads(input_str)
        path = data.get("path") or data.get("file_path") or ""
        old_text = data.get("old_text") or data.get("old") or ""
        new_text = data.get("new_text") or data.get("new") or ""
        replace_all = bool(data.get("replace_all", False))
        if not path:
            return "[錯誤] edit_file 需要 'path' 欄位"
        if not old_text:
            return "[錯誤] edit_file 需要 'old_text' 欄位(要被替換的文字)"
        path = _wsl_to_windows_path(path)
        p = Path(path).expanduser()
        if not p.exists():
            return f"[錯誤] 檔案不存在:{p}"
        content = p.read_text(encoding="utf-8", errors="replace")
        count = content.count(old_text)
        if count == 0:
            return (f"[錯誤] 在 {p.name} 找不到 old_text。檢查:(1) 完全相符的字串(含空格/換行);"
                    f"(2) 用 read_file 確認檔案目前內容")
        if count > 1 and not replace_all:
            return (f"[錯誤] old_text 在檔案內出現 {count} 次、不唯一。"
                    f"加更多上下文讓它唯一、或加 \"replace_all\": true 全部替換")
        new_content = content.replace(old_text, new_text) if replace_all else content.replace(old_text, new_text, 1)
        p.write_text(new_content, encoding="utf-8")
        return f"✓ 已替換 {p.name}({count if replace_all else 1} 處)、檔案 {p.stat().st_size:,} bytes"
    except _json.JSONDecodeError as e:
        return f"[錯誤] JSON 格式錯誤:{e}"
    except Exception as e:
        return f"[錯誤] 編輯失敗:{e.__class__.__name__}: {e}"


def _skill_export_var(input_str: str, cwd: Optional[str] = None) -> str:
    """export_var(name, value) — 把一個算好的值傳給下游節點(尤其 condition 條件節點)。

    JSON input: {"name": "score", "value": 380}
    實作:寫進 workflow 資料夾的 _step_export.json(扁平 dict);runner 在這一步結束後
    讀進該步的 step_vars,下游用 {{ steps.<本步驟名>.output.<name> }} 引用。
    多次呼叫會累積(各自的 name 都留著)。skill 完全不用知道檔名 / 格式。"""
    import json as _json
    try:
        d = _json.loads(input_str)
    except Exception:
        return ('[錯誤] export_var 需要 JSON 輸入,格式 '
                '{"name": "變數名", "value": 值}')
    if not isinstance(d, dict) or "name" not in d:
        return ('[錯誤] export_var 需要 "name" 欄位,格式 '
                '{"name": "變數名", "value": 值}')
    name = str(d["name"]).strip()
    if not name:
        return "[錯誤] export_var 的 name 不能為空"
    value = d.get("value")
    try:
        # skill 在 WSL 沙盒跑時 cwd 是 /mnt/c/... → host(Windows)上 Path("/mnt/c/..") 會解成
        # 磁碟根的 \mnt\c\.. bogus 路徑、寫入 FileNotFoundError(其它工具都有做這轉換、唯獨這裡漏了)。
        # 轉成 Windows 路徑後才寫,runner 才讀得到 _step_export.json → 變數才進得了 step_vars。
        base = Path(_wsl_to_windows_path(cwd)) if cwd else Path(".")
        export_f = base / "_step_export.json"
        existing: dict = {}
        if export_f.is_file():
            try:
                _loaded = _json.loads(export_f.read_text(encoding="utf-8"))
                if isinstance(_loaded, dict):
                    existing = _loaded
            except Exception:
                existing = {}
        existing[name] = value
        export_f.write_text(
            _json.dumps(existing, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        return (f"[export_var] 已匯出 {name} = {value!r}。"
                f"下游節點可用 steps.<本步驟名>.output.{name} 引用。")
    except Exception as e:
        return f"[錯誤] export_var 寫入失敗:{e.__class__.__name__}: {e}"


def _skill_grep(input_str: str) -> str:
    """grep(pattern, path) — JSON: {"pattern":"...", "path":"...", "glob":"*.py"(選填), "max_results":50(選填)}.
    用 ripgrep / fallback re。回 matching lines with file:line prefix。"""
    try:
        import json as _json
        import re as _re
        data = _json.loads(input_str)
        pattern = data.get("pattern") or ""
        search_path = data.get("path") or "."
        glob_filter = data.get("glob") or ""
        max_results = int(data.get("max_results") or 50)
        case_sensitive = bool(data.get("case_sensitive", True))
        if not pattern:
            return "[錯誤] grep 需要 'pattern' 欄位"
        search_path = _wsl_to_windows_path(search_path)
        root = Path(search_path).expanduser()
        if not root.exists():
            return f"[錯誤] 路徑不存在:{root}"
        # 用 Python re scan(避免依賴外部 ripgrep);大量檔案搜尋仍 OK
        flags = 0 if case_sensitive else _re.IGNORECASE
        try:
            regex = _re.compile(pattern, flags)
        except _re.error as e:
            return f"[錯誤] regex 語法錯:{e}。試試簡化 pattern 或加 \\ 跳脫特殊字元"
        results: list[str] = []
        files_iter = [root] if root.is_file() else root.rglob(glob_filter or "*")
        for fp in files_iter:
            if not fp.is_file() or fp.suffix.lower() in {'.pkl', '.zip', '.gz', '.png', '.jpg', '.jpeg', '.pdf', '.xlsx', '.docx', '.pptx'}:
                continue
            try:
                with open(fp, "r", encoding="utf-8", errors="replace") as f:
                    for lineno, line in enumerate(f, 1):
                        if regex.search(line):
                            rel = fp.relative_to(root) if root.is_dir() else fp.name
                            results.append(f"{rel}:{lineno}:{line.rstrip()[:200]}")
                            if len(results) >= max_results:
                                break
            except (OSError, UnicodeDecodeError):
                continue
            if len(results) >= max_results:
                break
        if not results:
            return f"(no match for `{pattern}` in {root})"
        suffix = f"\n... (截斷,只顯示前 {max_results} 筆)" if len(results) >= max_results else ""
        return "\n".join(results) + suffix
    except _json.JSONDecodeError as e:
        return f"[錯誤] JSON 格式錯誤:{e}"
    except Exception as e:
        return f"[錯誤] grep 失敗:{e.__class__.__name__}: {e}"


def _skill_glob(input_str: str) -> str:
    """glob(pattern, path) — JSON: {"pattern":"*.py", "path":"."(選填), "max_results":100(選填)}.
    用 pathlib.rglob 列檔。回 sorted file paths。"""
    try:
        import json as _json
        data = _json.loads(input_str)
        pattern = data.get("pattern") or ""
        search_path = data.get("path") or "."
        max_results = int(data.get("max_results") or 100)
        if not pattern:
            return "[錯誤] glob 需要 'pattern' 欄位(例 *.py、**/*.csv)"
        search_path = _wsl_to_windows_path(search_path)
        root = Path(search_path).expanduser()
        if not root.exists():
            return f"[錯誤] 路徑不存在:{root}"
        if not root.is_dir():
            return f"[錯誤] glob 的 path 必須是目錄:{root}"
        matches = []
        for fp in root.rglob(pattern):
            if fp.is_file():
                rel = fp.relative_to(root)
                size = fp.stat().st_size
                matches.append(f"{rel} ({size:,} B)")
                if len(matches) >= max_results:
                    break
        if not matches:
            return f"(no file matches `{pattern}` in {root})"
        matches.sort()
        suffix = f"\n... (截斷,只顯示前 {max_results} 筆)" if len(matches) >= max_results else ""
        return f"{len(matches)} files matching `{pattern}` in {root}:\n" + "\n".join(matches) + suffix
    except _json.JSONDecodeError as e:
        return f"[錯誤] JSON 格式錯誤:{e}"
    except Exception as e:
        return f"[錯誤] glob 失敗:{e.__class__.__name__}: {e}"


IMAGE_EXTS_SKILL = {'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                    '.gif': 'image/gif', '.webp': 'image/webp', '.bmp': 'image/bmp'}


def _wsl_to_windows_path(path: str) -> str:
    """LLM 在沙盒裡跑時常吐 `/mnt/c/Users/...` 路徑，但 view_image 在 host（Windows）上
    讀檔，要把它轉回 `C:\\Users\\...`。已是 Windows 路徑或非 /mnt/ 開頭就原樣回。
    沒這層翻譯 V3 log 看到 view_image 永遠回「圖片不存在」+ LLM 幻覺出假結果。"""
    import re as _re
    m = _re.match(r"^/mnt/([a-z])/(.*)$", path.strip())
    if not m:
        return path
    drive = m.group(1).upper()
    rest = m.group(2).replace("/", "\\")
    return f"{drive}:\\{rest}"


def _skill_view_image(path: str) -> dict:
    """讀圖片並回 base64 給 agent loop 注入多模態訊息。
    回 {"text": ..., "image_b64": str|None, "image_mime": str|None}
    上限 20 MB；超過或不是圖片就回錯誤訊息（image_b64=None）。"""
    try:
        cleaned = path.strip().strip('"').strip("'")
        # 沙盒路徑 → Windows 路徑（V3 view_image bug：LLM 跑沙盒給 /mnt/c/... 結果讀不到）
        cleaned = _wsl_to_windows_path(cleaned)
        p = Path(cleaned).expanduser()
        if not p.exists():
            return {"text": f"[錯誤] 圖片不存在：{path}（解析後：{p}）", "image_b64": None, "image_mime": None}
        ext = p.suffix.lower()
        if ext not in IMAGE_EXTS_SKILL:
            return {"text": f"[錯誤] 不支援的圖片格式：{ext}，支援 {list(IMAGE_EXTS_SKILL.keys())}",
                    "image_b64": None, "image_mime": None}
        data = p.read_bytes()
        if len(data) > 20 * 1024 * 1024:
            return {"text": f"[錯誤] 圖片過大（{len(data):,} bytes，上限 20MB）",
                    "image_b64": None, "image_mime": None}
        b64 = base64.b64encode(data).decode()
        mime = IMAGE_EXTS_SKILL[ext]
        return {"text": f"圖片 {p.name}（{len(data):,} bytes），已載入供視覺分析",
                "image_b64": b64, "image_mime": mime}
    except Exception as e:
        return {"text": f"[錯誤] 圖片讀取失敗：{e}", "image_b64": None, "image_mime": None}


def _skill_web_search(tool_input: str, call_count: int = 0,
                      logger: Optional[logging.Logger] = None) -> str:
    """用 Tavily API 搜網。兩段式輸出：
      OFF（include_full_content=false）= answer + URL 清單（~500 字元）
      ON （include_full_content=true ）= answer + URL + 每則完整原文（~15000 字元）
    ON 模式由 Tavily 端直接回完整文章內容（include_raw_content=True），
    Agent 不用自己寫 requests.get / newspaper 爬蟲（省失敗率）。
    """
    _lg = logger if logger is not None else log
    # ── 1. 設定檢查 ──
    try:
        import sys as _sys
        _backend_dir = str(Path(__file__).resolve().parent.parent)
        if _backend_dir not in _sys.path:
            _sys.path.insert(0, _backend_dir)
        from settings import get_settings as _gs
    except Exception as e:
        return f"[web_search 錯誤] 無法載入 settings：{e}"
    s = _gs()
    if not s.get("web_search_enabled"):
        return "[web_search 錯誤] 網路搜尋未啟用（Settings → 網路搜尋 → 啟用）"
    key = (s.get("tavily_api_key") or "").strip()
    if not key:
        return "[web_search 錯誤] Tavily API key 未設定（Settings → 網路搜尋 → API Key）"
    # ── 2. 呼叫次數上限 ──
    if call_count > WEB_SEARCH_MAX_PER_STEP:
        return (f"[web_search 錯誤] 本步驟已達搜尋次數上限（{WEB_SEARCH_MAX_PER_STEP} 次）。"
                "請整合前面搜尋結果回答，或呼叫 done(success=false) 說明需要更多搜尋。")
    # ── 3. 參數解析 ──
    params: dict = {}
    tool_input = tool_input.strip()
    if tool_input.startswith("{"):
        try:
            params = json.loads(tool_input)
        except json.JSONDecodeError as e:
            return f"[web_search 錯誤] input 不是合法 JSON：{e}（或直接傳純字串當 query）"
    else:
        params = {"query": tool_input}
    query = (params.get("query") or "").strip()
    if not query:
        return "[web_search 錯誤] query 不可為空"
    max_results = max(1, min(int(params.get("max_results", 5)), 8))
    # search_depth 預設讀 settings.web_search_deep_default(預設 True、advanced)、
    # LLM 可以在 input 內傳 search_depth 個別 override(基本上不必、預設就好)。
    _deep_default = bool(s.get("web_search_deep_default", True))
    _depth_input = params.get("search_depth")
    if _depth_input is None:
        search_depth = "advanced" if _deep_default else "basic"
    else:
        search_depth = "advanced" if str(_depth_input).lower() == "advanced" else "basic"
    # 完整內容模式：預設從 settings 取、agent 可 per-call 覆寫
    full_content = bool(params.get("include_full_content",
                                   s.get("web_search_full_content_default", False)))
    # ── 4. 呼叫 Tavily ──
    import requests as _requests
    try:
        resp = _requests.post(
            "https://api.tavily.com/search",
            json={
                "api_key": key,
                "query": query,
                "max_results": max_results,
                "search_depth": search_depth,
                "include_answer": True,
                # ON 模式：讓 Tavily 回完整文章原文（他們處理 CF / JS 渲染等）
                # OFF 模式：不要原文，輕量模式節省 context
                "include_raw_content": full_content,
            },
            timeout=45 if full_content else 20,  # full content 回傳較慢，給長一點 timeout
        )
        if resp.status_code == 401:
            return "[web_search 錯誤] Tavily API key 無效（401）"
        if resp.status_code == 429:
            return "[web_search 錯誤] Tavily 配額用盡或速率受限（429），請稍後再試"
        resp.raise_for_status()
        data = resp.json()
    except _requests.Timeout:
        return "[web_search 錯誤] Tavily 連線逾時"
    except _requests.HTTPError as e:
        return f"[web_search 錯誤] Tavily HTTP {resp.status_code}：{resp.text[:300]}"
    except Exception as e:
        return f"[web_search 錯誤] Tavily 呼叫失敗：{type(e).__name__}: {e}"
    # ── 5. 組裝輸出 ──
    answer = (data.get("answer") or "").strip()
    results = data.get("results") or []
    mode_tag = "full" if full_content else "light"
    lines = [f"[web_search] query=\"{query[:80]}\" (depth={search_depth}, mode={mode_tag})"]
    if answer:
        lines.append(f"answer: {answer}")
    lines.append("")
    lines.append(f"來源 (共 {len(results)} 項)：")
    for i, r in enumerate(results, start=1):
        title = (r.get("title") or "").strip()[:WEB_SEARCH_TITLE_CHARS]
        url = (r.get("url") or "").strip()
        lines.append(f"[{i}] {title} — {url}")
        if full_content:
            # include_raw_content 會回 raw_content；回 None 時退到 content（短摘要）
            raw = (r.get("raw_content") or r.get("content") or "").strip()
            if raw:
                # 正規化換行空白，避免 agent 吃到一堆 \n\n\n
                raw = re.sub(r"\n{3,}", "\n\n", raw)
                if len(raw) > WEB_SEARCH_PER_RESULT_FULL_CHARS:
                    raw = raw[:WEB_SEARCH_PER_RESULT_FULL_CHARS] + "…（本篇截斷）"
                lines.append("--- 內文 ---")
                lines.append(raw)
                lines.append("--- /內文 ---")
    output = "\n".join(lines)
    cap = WEB_SEARCH_OUTPUT_CHAR_CAP_FULL if full_content else WEB_SEARCH_OUTPUT_CHAR_CAP_LIGHT
    truncated = False
    if len(output) > cap:
        output = output[:cap] + f"\n…（總輸出已截斷，完整 {len(output)} 字；下次縮小 max_results 或關閉 include_full_content）"
        truncated = True
    _lg.info(
        f"[web_search] query={query[:60]!r} → 回傳 {len(output)} 字元 "
        f"(mode={mode_tag}, depth={search_depth}, results={len(results)}"
        f"{', truncated' if truncated else ''})"
    )
    return output


# 抽 code block 的兩條 regex(稽查 B / bug 家族 1:行內 ``` 腰斬):
# OWNLINE 要求「閉合 ``` 自成一行(前有換行)」→ LLM 程式碼字串內的行內 ``` 不會被誤判成結束
# (與 main.py YAML 抽取同一招、已驗證)。找不到才退回 LEGACY 非貪婪(保證不比修前差)。
_CODE_FENCE_OWNLINE = re.compile(r'```(?:python|json|bash|sh)?[ \t]*\n([\s\S]*?)\n[ \t]*```')
_CODE_FENCE_LEGACY = re.compile(r'```(?:python|json|bash|sh)?\s*\n(.*?)```', re.DOTALL)


def _find_code_blocks(text: str) -> list:
    """回傳所有 code block 的 match;優先用「閉合自成一行」、抓不到才退回舊式。"""
    blocks = list(_CODE_FENCE_OWNLINE.finditer(text))
    return blocks if blocks else list(_CODE_FENCE_LEGACY.finditer(text))


def _balanced_json(s: str) -> Optional[str]:
    """從字串第一個 '{' 起,做「字串感知的大括號平衡掃描」抽出完整 JSON 物件。
    取代 re.search(r'{.*?}') —— 非貪婪會在第一個 '}' 截斷、巢狀 JSON 必壞(稽查 B)。"""
    start = s.find('{')
    if start < 0:
        return None
    depth = 0
    in_str = False
    esc = False
    for j in range(start, len(s)):
        c = s[j]
        if in_str:
            if esc:
                esc = False
            elif c == '\\':
                esc = True
            elif c == '"':
                in_str = False
        else:
            if c == '"':
                in_str = True
            elif c == '{':
                depth += 1
            elif c == '}':
                depth -= 1
                if depth == 0:
                    return s[start:j + 1]
    return None


def _extract_code_block(text: str) -> Optional[str]:
    """從 markdown code block 中提取程式碼內容。"""
    blocks = _find_code_blocks(text)
    return blocks[0].group(1).strip() if blocks else None


def _extract_text(content) -> str:
    """把 LLM 回應的 content 正規化成純文字(稽查 D / bug 家族 5)。
    - str → 原樣
    - list-of-blocks(Gemini-3 / Anthropic 回 [{'type':'text','text':..},{'type':'thinking',..}])
      → 只串接 text 區塊(thinking / 其他類型不收)
    - 其他 → ''
    取代散落各處的 `response.content if isinstance(.,str) else ""` —— 那會把 list 內容整個丟成空字串,
    害 gemini-3 的文字回答被當『模型什麼都沒說』(實測 gemini-3-flash content 就是 list)。"""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for b in content:
            if isinstance(b, dict):
                if b.get("type") == "text" and isinstance(b.get("text"), str):
                    parts.append(b["text"])
            elif isinstance(b, str):
                parts.append(b)
        return "".join(parts)
    return ""


def _is_malformed_empty(response) -> bool:
    """稽查 D / bug 家族 5:Gemini 偶發 MALFORMED_FUNCTION_CALL ——
    回應的 content 與 tool_calls 都空、finish_reason 標 MALFORMED。
    這是『模型這次生成壞了、重試常會過』的暫時性錯誤,不是『模型不肯呼叫工具』。
    用來在 native FC 迴圈裡把它當可重試、而非當成 no-tool 直接放棄。"""
    try:
        content = _extract_text(response.content)
        tc = len(getattr(response, "tool_calls", []) or [])
        if content or tc:
            return False
        fr = str((getattr(response, "response_metadata", None) or {}).get("finish_reason") or "").upper()
        return "MALFORMED" in fr
    except Exception:
        return False


def _sanitize_code(code: str) -> str:
    """清除混入程式碼中的 LLM 解釋文字（非 Python/Shell 語法的行）。"""
    lines = code.split('\n')
    # 找到第一行有效程式碼（import, from, def, class, #, 變量賦值, 函式呼叫等）
    code_pattern = re.compile(
        r'^(\s*(import |from |def |class |if |for |while |with |try:|except |'
        r'return |print|#|[a-zA-Z_]\w*\s*[=(]|plt\.|df\.|pd\.|np\.|sns\.|'
        r'\[|{|}|\]|\)|"|\'|$))'
    )
    start_idx = 0
    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue
        if code_pattern.match(stripped):
            start_idx = i
            break
    # 從第一行有效程式碼開始，過濾掉純中文解釋行（不在字串內的非 ASCII 開頭行）
    result = []
    for line in lines[start_idx:]:
        stripped = line.strip()
        if not stripped:
            result.append(line)
            continue
        # 如果整行以中文/全形字元開頭且不是 Python 字串或註解
        first_char = stripped[0]
        if ord(first_char) > 0x2E00 and not stripped.startswith('#') and not stripped.startswith(("'", '"')):
            continue  # 跳過純中文解釋行
        result.append(line)
    return '\n'.join(result).strip()


def _parse_skill_tool_calls(text: str) -> list[dict]:
    """
    解析 LLM 回覆中的工具呼叫。

    LLM 常見輸出格式：
    1. <tool>name</tool> <input>content</input>                     （標準）
    2. <tool>name</tool> ```python\ncode```                          （code block）
    3. <tool>name</tool> ```json\n{"key":"val"}```                   （json block）
    4. <tool>name</tool>\n直接跟隨程式碼或JSON                        （無標籤無block）
    5. ```python\n<tool>name</tool>\n<input>content</input>\n```     （整體在block內）

    關鍵：run_python/run_shell 的 input 只應包含可執行程式碼，
    不能混入 LLM 的解釋文字（會導致 SyntaxError）。
    """
    calls = []

    # ── Step 1：嘗試標準 <input>...</input> 格式 ──
    pattern_std = re.compile(r'<tool>(.*?)</tool>\s*<input>(.*?)</input>', re.DOTALL)
    for m in pattern_std.finditer(text):
        calls.append({"tool": m.group(1).strip(), "input": m.group(2).strip()})
    if calls:
        return calls

    # ── Step 2：找所有 code blocks，再找離 <tool> 最近的那個 ──
    # 先提取所有 code blocks 及其位置(用「閉合自成一行」抽取、避免行內 ``` 腰斬)
    code_blocks = _find_code_blocks(text)
    # 找所有 <tool> 標籤
    tool_tags = list(re.finditer(r'<tool>(.*?)</tool>', text))

    for tag in tool_tags:
        tool_name = tag.group(1).strip()
        tag_start = tag.start()
        tag_end = tag.end()

        # 先找 tag 之後最近的 code block
        best_block = None
        for block in code_blocks:
            if block.start() >= tag_end:
                best_block = block
                break

        # 如果 tag 之後沒有 code block，往前找最近的（LLM 先放 code 再放 tag）
        if not best_block:
            for block in reversed(code_blocks):
                if block.end() <= tag_start:
                    best_block = block
                    break

        if best_block:
            content = best_block.group(1).strip()
            # 對 run_python 清洗混入的中文解釋
            if tool_name in ('run_python', 'run_shell'):
                content = _sanitize_code(content)
            if content and len(content) > 2:
                calls.append({"tool": tool_name, "input": content})
                return calls  # 一次只處理一個工具呼叫

    # ── Step 3：done 工具 — 找 JSON ──
    done_match = re.search(r'<tool>done</tool>', text)
    if done_match:
        # 在 done 標籤後找 JSON(大括號平衡掃描、避免巢狀 JSON 被第一個 '}' 截斷)
        after_done = text[done_match.end():]
        json_str = _balanced_json(after_done)
        if json_str:
            return [{"tool": "done", "input": json_str.strip()}]

    # ── Step 4：沒有 <tool> 標籤，但有 code block（LLM 忘記加標籤）──
    if not tool_tags and code_blocks:
        content = code_blocks[-1].group(1).strip()  # 取最後一個 code block
        # 猜測工具類型
        if content.startswith('{') and 'success' in content:
            return [{"tool": "done", "input": content}]
        elif content.startswith('{') and 'status' in content:
            return [{"tool": "done", "input": content}]

    # ── Step 5：fallback — 清除 code block 標記後找 raw content ──
    cleaned = re.sub(r'```(?:python|json|bash|sh)?\s*\n?', '', text)
    cleaned = cleaned.replace('```', '')

    pattern_raw = re.compile(r'<tool>(.*?)</tool>\s*(.+?)(?=<tool>|$)', re.DOTALL)
    for m in pattern_raw.finditer(cleaned):
        tool_name = m.group(1).strip()
        content = m.group(2).strip()
        if tool_name in ('run_python', 'run_shell'):
            content = _sanitize_code(content)
        if content and len(content) > 2:
            calls.append({"tool": tool_name, "input": content})
            break

    return calls


def _execute_skill_tool(tool_name: str, tool_input: str, cwd: Optional[str] = None, run_id: str = "",
                        logger: Optional[logging.Logger] = None, force_host: bool = False,
                        tool_timeout: int = SKILL_TOOL_TIMEOUT) -> str:
    """執行單一工具。
    若 settings.skill_sandbox_mode='wsl_docker' 且沙盒可用，run_python / run_shell
    會走沙盒容器；其餘情況走原本 host subprocess。
    force_host=True：跳過沙盒檢查直接走 host（使用者透過 ask_user 同意 fallback 時 caller 會傳）。
    logger: per-step 的 pipeline logger（有寫到 .log 檔）；None 的話沙盒標記只會印到 backend stdout。
    tool_timeout：單次 run_python / run_shell 上限秒數（從 step.timeout 推導，見 _compute_tool_timeout）。"""
    # sandbox fallback 警告:設成 wsl_docker 但 sandbox 跑不起來時 LLM 不知情、
    # 繼續用 /mnt/d/ Linux 路徑撞 Windows host、燒 token 瘋狂 retry(2026-05-24 root cause)。
    # 在 host result 前綴明確告知 LLM 路徑與 shell 已切換。
    # ── 硬攔 run_shell 的 pip install:skill 不准自己裝套件、強制走 missing_dependency 使用者確認 ──
    # (2026-05-31 發現:gemma 會 run_shell pip install 繞過確認;executor 的 missing_module steering
    #  只是軟提示、擋不住會自作主張的模型。這裡硬攔、回 steering 逼它改 done(missing_packages)。)
    if tool_name == "run_shell" and _is_pip_install_cmd(tool_input):
        if logger:
            logger.warning(f"[run_shell] 攔截 pip install、要求改走 done(missing_packages)：{(tool_input or '')[:80]}")
        return (
            "[系統攔截] 不允許用 run_shell 直接 pip install 裝套件。\n"
            "缺套件時請改呼叫 done(success=false, missing_packages=[\"套件名\"])、"
            "系統會跳出『允許安裝』確認給使用者、同意後才裝到容器、並自動重跑這步。\n"
            "(這是系統層級的安全規定:裝任何依賴都必須經使用者確認、不能由 skill 自行安裝。)"
        )

    _sandbox_warn_prefix = ""
    if tool_name in ("run_python", "run_shell") and not force_host:
        _sandbox_meta: dict = {}
        sandbox_out = _try_sandbox_exec(
            tool_name, tool_input, cwd, run_id, logger,
            tool_timeout=tool_timeout, out_meta=_sandbox_meta,
        )
        if sandbox_out is not None:
            return sandbox_out
        # 真的 fallback 到 host 才加 warning(不是 host 模式 normal 跑)
        if _sandbox_meta.get("mode_was") == "wsl_docker":
            _reason = _sandbox_meta.get("fallback_reason") or "未知"
            _sandbox_warn_prefix = (
                f"[⚠️ 沙盒不可用、本次 fallback 到 Windows host 執行 — 原因:{_reason}]\n"
                f"⚠ 路徑請用 Windows 格式:D:\\... 或 C:\\...(不要用 /mnt/d/ 等 Linux 路徑)\n"
                f"⚠ shell 是 cmd.exe、不認 ls / cat / grep / which:用 dir / type / findstr / where\n"
                f"⚠ 設定請使用者去 Settings 切換沙盒開關恢復 — 之前已正常跑、現在突然不行通常是 Docker 或 WSL 暫斷\n"
                f"---\n"
            )
    if tool_name == "run_python":
        _r = _skill_run_python(tool_input, cwd=cwd, run_id=run_id, tool_timeout=tool_timeout)
        return _sandbox_warn_prefix + _r if _sandbox_warn_prefix else _r
    elif tool_name == "run_shell":
        _r = _skill_run_shell(tool_input, cwd=cwd, run_id=run_id, tool_timeout=tool_timeout)
        return _sandbox_warn_prefix + _r if _sandbox_warn_prefix else _r
    elif tool_name == "read_file":
        # 兩種輸入皆支援:
        #   1) 裸路徑字串(舊行為、向後相容)
        #   2) JSON {"path": "...", "offset": N, "limit": N}(分段讀大檔)
        _rf_path, _rf_offset, _rf_limit = tool_input, 0, 100
        _ti_stripped = (tool_input or "").strip()
        if _ti_stripped.startswith("{"):
            try:
                import json as _json_rf
                _d = _json_rf.loads(_ti_stripped)
                if isinstance(_d, dict) and _d.get("path"):
                    _rf_path = str(_d["path"])
                    _rf_offset = int(_d.get("offset", 0) or 0)
                    _rf_limit = int(_d.get("limit", 100) or 100)
            except Exception:
                pass  # 解析失敗 → 當裸路徑處理
        return _skill_read_file(_rf_path, max_lines=_rf_limit, offset=_rf_offset)
    elif tool_name == "write_file":
        return _skill_write_file(tool_input)
    elif tool_name == "edit_file":
        return _skill_edit_file(tool_input)
    elif tool_name == "export_var":
        return _skill_export_var(tool_input, cwd=cwd)
    elif tool_name == "grep":
        return _skill_grep(tool_input)
    elif tool_name == "glob":
        return _skill_glob(tool_input)
    elif tool_name == "web_search":
        # call_count 由呼叫方維護（每個 skill step 獨立計數）— 這邊拿不到，交由 agent loop 處理呼叫前計數
        return _skill_web_search(tool_input, logger=logger)
    elif tool_name == "view_image":
        # 特殊標記，agent loop 會看到後改走多模態 HumanMessage 路徑（注入 image_url）
        return "__VIEW_IMAGE__"
    elif tool_name == "done":
        return "__DONE__"
    else:
        return f"[錯誤] 未知工具：{tool_name}"


# ── 沙盒路由（V3） ────────────────────────────────────────────────
# 避免每次呼叫都 log「沙盒不可用」洗頻、但又要避免「第一次失敗後永久靜音」
# 導致 LLM 一直用 sandbox 路徑撞 host(2026-05-24 ai_coding_market_research 燒
# $1.5-2 token 的 root cause)。改用 dict[reason, last_warn_ts],過 5 分鐘可再 log。
_SANDBOX_WARNED: dict[str, float] = {}
_SANDBOX_WARN_COOLDOWN_SEC = 300.0  # 同 reason 每 5 分鐘最多 log 一次


async def _preflight_sandbox(
    ask_mode: bool,
    fallback_state: dict,
    run_id: str,
    step_name: str,
    logger: logging.Logger,
) -> str:
    """在 run_python / run_shell 被真的執行前，先判斷沙盒可不可用。
    回傳 'sandbox' | 'host' | 'abort' 三種決策，交給 agent loop 處理。

    fallback_state: 跨 iteration 的可變 dict，用來記「使用者這一步已經同意 fallback」
                    的決定，同一步內後續 tool 呼叫不會再被問一次。
                    格式：{'allowed': bool}

    ask_mode=False：維持舊行為（靜默 fallback），回傳 'host' 或 'sandbox'（看狀態）
    ask_mode=True ：沙盒不可用時呼叫 ask_user 問使用者，選項：重試 / 退 host / 中止
    """
    try:
        import sys as _sys
        _backend_dir = str(Path(__file__).resolve().parent.parent)
        if _backend_dir not in _sys.path:
            _sys.path.insert(0, _backend_dir)
        from settings import get_settings
        from pipeline import sandbox as _sandbox
    except Exception:
        # 設定 / 沙盒模組無法載入 → 當作 host 模式處理
        return "host"

    mode = (get_settings().get("skill_sandbox_mode") or "host").strip()
    if mode != "wsl_docker":
        return "sandbox"  # 不用沙盒；交給 _execute_skill_tool 走 host（不會觸發 sandbox 路徑）

    # 使用者這一步已同意 fallback 了，不要再問
    if fallback_state.get("allowed"):
        return "host"

    ok, reason = _sandbox.ensure_running()
    if ok:
        return "sandbox"

    # 沙盒不可用，ask_mode OFF → 靜默 fallback（維持舊行為，不中斷 pipeline）
    if not ask_mode:
        fallback_state["allowed"] = True
        return "host"

    # ask_mode ON → 問使用者怎麼處理；最多問 5 輪「重試」避免無限迴圈
    for attempt in range(5):
        answer = await _wait_for_ask_user(
            run_id=run_id,
            question=(
                f"⚠️ 沙盒容器不可用 ── {reason}\n\n"
                "請選擇如何繼續：\n"
                "• 重試沙盒：再試一次（WSL 冷啟動通常一兩次就通）\n"
                "• 退回 host 模式：直接在 Windows host 跑（本次步驟的後續 tool 也都走 host）\n"
                "• 中止步驟：放棄這個 skill step"
            ),
            options=["重試沙盒", "退回 host 模式", "中止步驟"],
            context=f"ask_mode 已啟用，沙盒狀態異常。若沙盒只是短暫忙（VM 冷啟）選「重試沙盒」。",
            logger=logger,
            step_name=step_name,
        )
        if answer is None:
            logger.warning(f"[{step_name}] 沙盒 ask_user 取消或逾時 → 中止")
            return "abort"
        if "中止" in answer:
            return "abort"
        if "重試" in answer:
            logger.info(f"[{step_name}] 使用者選擇重試沙盒（第 {attempt + 1} 次）")
            ok, reason = _sandbox.ensure_running()
            if ok:
                logger.info(f"[{step_name}] 重試後沙盒已恢復")
                return "sandbox"
            # 繼續下一輪問
            continue
        if "host" in answer.lower() or "退" in answer:
            logger.info(f"[{step_name}] 使用者同意此步驟 fallback 到 host")
            fallback_state["allowed"] = True
            return "host"
        # 非預期的回答 → 當作 host 比較安全（至少工作能繼續）
        logger.warning(f"[{step_name}] 無法解析沙盒 ask_user 回答：{answer!r} → 預設 host")
        fallback_state["allowed"] = True
        return "host"
    # 重試 5 次還是不行
    logger.warning(f"[{step_name}] 沙盒連續 5 次重試失敗 → 中止")
    return "abort"


def _try_sandbox_exec(tool_name: str, tool_input: str, cwd: Optional[str], run_id: str,
                      logger: Optional[logging.Logger] = None,
                      tool_timeout: int = SKILL_TOOL_TIMEOUT,
                      out_meta: Optional[dict] = None) -> Optional[str]:
    """若 settings.skill_sandbox_mode='wsl_docker' 且沙盒可用，就把 run_python/run_shell
    送進 pipeline-sandbox-v4 容器執行。回傳組好的 output 字串（格式對齊 host 版本）；
    若 mode=host 或沙盒不可用則回傳 None 讓 caller fallback 到 host subprocess。
    logger: per-step pipeline logger；若提供則沙盒標記會出現在 .log 檔，否則只出現在 backend stdout。
    tool_timeout：單次 tool 執行上限秒數（從 step.timeout 推導，見 _compute_tool_timeout）。
    out_meta: mutable dict, 失敗時填 {"fallback_reason": str, "mode_was": "wsl_docker"}
              caller 用來判斷要不要在 host result 前綴警告給 LLM（防 LLM 還以為在 sandbox)。"""
    _lg = logger if logger is not None else log
    try:
        import sys as _sys
        _backend_dir = str(Path(__file__).resolve().parent.parent)
        if _backend_dir not in _sys.path:
            _sys.path.insert(0, _backend_dir)
        from settings import get_settings
        from pipeline import sandbox as _sandbox
    except Exception as e:
        _lg.warning(f"[sandbox] import 失敗（fallback 到 host）：{e}")
        if out_meta is not None:
            out_meta["fallback_reason"] = f"import failed: {e}"
            out_meta["mode_was"] = "wsl_docker"  # 不知道、保守假設
        return None

    settings_dict = get_settings()
    mode = (settings_dict.get("skill_sandbox_mode") or "host").strip()
    # 每次 skill tool 呼叫都 log 一下目前讀到什麼模式，方便追蹤使用者看到的 UI
    # 跟後端實際決策有沒有差距（之前出現過 UI 顯示藍色但實際走 host 的懸案）
    _lg.info(f"[sandbox] 檢查：skill_sandbox_mode={mode!r}（來自 settings cache）")
    if mode != "wsl_docker":
        return None

    # ensure_running 失敗時 retry 一次,可能只是 Docker 暫時 hiccup
    ok, reason = _sandbox.ensure_running()
    if not ok:
        _lg.info(f"[sandbox] ensure_running 失敗({reason})、1s 後 retry...")
        import time as _t
        _t.sleep(1.0)
        ok, reason = _sandbox.ensure_running()
    if not ok:
        # cooldown:同 reason 5 分鐘最多 log warning 一次,但不會永久靜音
        key = reason or "unknown"
        import time as _t
        now = _t.time()
        last = _SANDBOX_WARNED.get(key, 0)
        if now - last >= _SANDBOX_WARN_COOLDOWN_SEC:
            _lg.warning(
                f"[sandbox] ⚠ 沙盒不可用、此次 fallback 到 Windows host:{reason}"
                f"(同錯誤 {_SANDBOX_WARN_COOLDOWN_SEC:.0f}s 內不重複 log)"
            )
            _SANDBOX_WARNED[key] = now
        if out_meta is not None:
            out_meta["fallback_reason"] = reason or "sandbox not running"
            out_meta["mode_was"] = "wsl_docker"
        return None
    # 沙盒恢復健康後，清掉之前的告警記錄下次若又壞可再提醒
    if _SANDBOX_WARNED:
        _SANDBOX_WARNED.clear()

    _lg.info(f"[sandbox] 🛡 在容器內執行 {tool_name}（{len(tool_input)} 字元、timeout={tool_timeout}s）")
    if tool_name == "run_python":
        res = _sandbox.run_python(
            tool_input, cwd=cwd,
            timeout=tool_timeout,
            run_id=run_id,
            register_cb=register_proc,
            unregister_cb=unregister_proc,
        )
    else:  # run_shell
        res = _sandbox.run_shell(
            tool_input, cwd=cwd,
            timeout=tool_timeout,
            run_id=run_id,
            register_cb=register_proc,
            unregister_cb=unregister_proc,
        )
    _lg.info(f"[sandbox] ✓ 容器執行完畢 rc={res.returncode}"
             + (" (timed out)" if res.timed_out else ""))

    # 組裝輸出 — 格式刻意與 host 版本一致，LLM 分不出差別
    output = ""
    if res.stdout:
        output += res.stdout
    if res.stderr:
        tag = "stderr" if res.returncode != 0 else "warnings"
        output += f"\n[{tag}]\n{res.stderr}"
    if res.returncode != 0:
        output += f"\n[exit code: {res.returncode}]"
        if not res.stdout and not res.stderr:
            output += (
                "\n[提示] 子程序非正常結束但沒捕捉到任何 stdout / stderr。"
                "請把整段程式用 try/except 包起來，except 裡 "
                "`import traceback; traceback.print_exc()` 再 `sys.exit(0)`。"
            )
    elif not res.stdout and tool_name == "run_python":
        output += "\n[執行成功，程式無 stdout 輸出]"
    return output.strip() or "(無輸出)"


async def _wait_for_ask_user(
    run_id: str,
    question: str,
    options: list,
    context: str,
    logger: logging.Logger,
    step_name: str,
) -> Optional[str]:
    """
    把問題送出去（Pipeline 進 awaiting_human + Telegram/前端 推問題），
    in-memory 等待答案送達（asyncio.Event），或 timeout 回 None。
    """
    from pipeline.store import get_store
    store = get_store()
    run = store.load(run_id)
    if not run:
        logger.warning(f"[{step_name}] ask_user 失敗：找不到 run {run_id}")
        return None

    # gemma 常把範例路徑寫成 `C:\\Users\\X`(雙反斜線、像在寫程式字串字面值)、顯示給人看很怪。
    # 問句 / context 純粹是給人讀的提示文字、不影響使用者實際輸入的答案 → 把 2+ 連續反斜線收斂成
    # 單一(純顯示修正、log 與前端對話框一起乾淨)。
    _bs = __import__("re")
    if question:
        question = _bs.sub(r"\\{2,}", lambda _m: "\\", question)
    if context:
        context = _bs.sub(r"\\{2,}", lambda _m: "\\", context)

    event = asyncio.Event()
    _pending_questions[run_id] = {
        "question": question,
        "options": options,
        "context": context,
        "event": event,
        "answer": None,
    }

    # 更新 run 狀態：進入 awaiting_human
    run.status = "awaiting_human"
    run.awaiting_type = "ask_user"
    run.awaiting_message = question
    run.awaiting_suggestion = json.dumps({"options": options, "context": context}, ensure_ascii=False)
    store.save(run)

    # 發通知（Telegram + 前端會 poll 到狀態變化）
    try:
        from pipeline.runner import _send_ask_user_notification
        await _send_ask_user_notification(run, question, options, context, step_name)
    except Exception as e:
        logger.warning(f"[{step_name}] ask_user 通知發送失敗：{e}")

    # 等待答案或 timeout
    logger.info(f"[{step_name}] ⏸ ask_user 等待中：{question}")
    try:
        await asyncio.wait_for(event.wait(), timeout=ASK_USER_TIMEOUT)
    except asyncio.TimeoutError:
        logger.warning(f"[{step_name}] ask_user 逾時（{ASK_USER_TIMEOUT}s）")
        _pending_questions.pop(run_id, None)
        # 恢復狀態（注意：可能其他邏輯已接手改狀態，這裡僅在仍為 ask_user 時清除）
        run2 = store.load(run_id)
        if run2 and run2.awaiting_type == "ask_user":
            run2.status = "running"
            run2.awaiting_type = ""
            run2.awaiting_message = ""
            run2.awaiting_suggestion = ""
            store.save(run2)
        return None

    answer = _pending_questions[run_id]["answer"]
    _pending_questions.pop(run_id, None)

    # 恢復 running 狀態
    run3 = store.load(run_id)
    if run3:
        run3.status = "running"
        run3.awaiting_type = ""
        run3.awaiting_message = ""
        run3.awaiting_suggestion = ""
        store.save(run3)

    logger.info(f"[{step_name}] ▶ ask_user 收到答案：{answer}")
    return answer


# ── Phase B: 命令授權等待（同 _wait_for_ask_user pattern）──
COMMAND_APPROVAL_TIMEOUT = 1800  # 30 分鐘等用戶按按鈕

async def _wait_for_command_approval(
    run_id: str,
    category: str,
    label: str,
    preview: str,
    tool_name: str,
    logger: logging.Logger,
    step_name: str,
) -> str:
    """ask_mode ON 時、敏感命令呼叫前的授權等待。
    回 "allow" / "deny" / "hint"（hint 是用戶要改任務）/ "timeout"。
    """
    from pipeline.store import get_store
    store = get_store()
    run = store.load(run_id)
    if not run:
        logger.warning(f"[{step_name}] command_approval 失敗：找不到 run {run_id}")
        return "deny"  # 找不到 run 安全起見拒絕

    event = asyncio.Event()
    _pending_command_approvals[run_id] = {
        "category": category,
        "label": label,
        "preview": preview,
        "tool_name": tool_name,
        "event": event,
        "decision": None,
    }

    # 寫 awaiting_human 狀態
    run.status = "awaiting_human"
    run.awaiting_type = "command_approval"
    run.awaiting_message = f"敏感操作需授權:{label}"
    run.awaiting_suggestion = json.dumps({
        "category": category,
        "label": label,
        "preview": preview,
        "tool_name": tool_name,
        "step_name": step_name,
    }, ensure_ascii=False)
    store.save(run)

    # 通知 TG / 前端
    try:
        from pipeline.runner import _send_command_approval_notification
        await _send_command_approval_notification(run, category, label, preview, step_name)
    except Exception as e:
        logger.warning(f"[{step_name}] command_approval 通知發送失敗：{e}")

    logger.warning(
        f"[{step_name}] ⏸ ask_mode 攔截敏感命令 ({category})、等用戶授權:{preview[:80]}"
    )
    try:
        await asyncio.wait_for(event.wait(), timeout=COMMAND_APPROVAL_TIMEOUT)
    except asyncio.TimeoutError:
        logger.warning(f"[{step_name}] command_approval 逾時 ({COMMAND_APPROVAL_TIMEOUT}s)")
        _pending_command_approvals.pop(run_id, None)
        run2 = store.load(run_id)
        if run2 and run2.awaiting_type == "command_approval":
            run2.status = "running"
            run2.awaiting_type = ""
            run2.awaiting_message = ""
            run2.awaiting_suggestion = ""
            store.save(run2)
        return "timeout"

    decision = _pending_command_approvals[run_id]["decision"]
    _pending_command_approvals.pop(run_id, None)

    run3 = store.load(run_id)
    if run3:
        run3.status = "running"
        run3.awaiting_type = ""
        run3.awaiting_message = ""
        run3.awaiting_suggestion = ""
        store.save(run3)

    logger.info(f"[{step_name}] ▶ command_approval 收到決定:{decision}")
    return decision or "deny"


def _read_file_sample(path: str, max_chars: int = 700) -> str:
    """讀前步驟輸出檔的一小段樣本,給 skill LLM 看「實際格式」——
    欄位名、值的大小寫、是字串還是數字、檔案結構 —— 讓它不用憑空猜。

    回傳空字串 = 不取樣(檔案不存在 / 二進位 / 太大 / 出錯)。"""
    try:
        p = Path(path)
        if not p.is_file():
            return ""
        size = p.stat().st_size
        ext = p.suffix.lower()
        if ext == ".json":
            if size > 3_000_000:
                return f"(JSON 檔較大 ~{size // 1024} KB、未取樣;動手前請自己讀檔頭確認格式)"
            import json as _j
            raw = p.read_text(encoding="utf-8", errors="replace")
            try:
                data = _j.loads(raw)
            except Exception:
                return raw[:max_chars]
            if isinstance(data, list):
                head = data[:2]
                s = _j.dumps(head, ensure_ascii=False, indent=2)
                return f"(JSON 陣列、共 {len(data)} 筆,前 {len(head)} 筆:)\n{s[:max_chars]}"
            if isinstance(data, dict):
                s = _j.dumps(data, ensure_ascii=False, indent=2)
                return f"(JSON 物件:)\n{s[:max_chars]}"
            return raw[:max_chars]
        if ext in (".csv", ".tsv", ".txt", ".md", ".log"):
            with p.open("r", encoding="utf-8", errors="replace") as f:
                buf = f.read(8000)
            lines = buf.splitlines()[:15]
            return "(前幾行:)\n" + "\n".join(lines)[:max_chars]
        if ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
                parts = [f"(Excel、工作表:{', '.join(wb.sheetnames)};第一張表前幾列:)"]
                ws = wb[wb.sheetnames[0]]
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= 4:
                        break
                    parts.append(str(list(row)))
                wb.close()
                return "\n".join(parts)[:max_chars]
            except Exception:
                return "(Excel 檔、無法取樣)"
        return ""  # 二進位 / 不認得的副檔名:不取樣
    except Exception:
        return ""


import re as _re_runts
_RUN_TS_RE = _re_runts.compile(r"run_\d{8}_\d{6}")


def _rebase_recipe_run_dir(code: str, output_path: Optional[str]) -> tuple[str, int]:
    """B′:快取 recipe 的程式碼裡寫死了「錄製當時」的 run_<時間戳> 子夾路徑,
    但每次執行的 run 夾時間戳不同。重播前把碼內所有 run_<舊ts> 換成「本次」的 run_<ts>。

    本次 run_<ts> 從 output_path 取(output_path 一定是 .../run_<ts>/... 的形式),
    比用 working_dir.basename 穩 —— 當 output.path 帶子目錄時(如 sub/report.md),
    working_dir 會是 .../run_<ts>/sub、basename='sub' 不是 run_ts,但 output_path 仍含 run_ts。

    回傳 (改寫後的碼, 替換處數)。下列情況回 (原碼, 0):
      - output_path 為空 / 不含 run_<ts>(例:output.path 是絕對路徑、本就不在 run 夾)
      - 快取碼裡根本沒有 run_<ts>(例:LLM 自己寫死了不含時間戳的固定路徑)→ 無從替換
    """
    if not code or not output_path:
        return code, 0
    m = _RUN_TS_RE.search(str(output_path))
    if not m:
        return code, 0
    return _RUN_TS_RE.subn(m.group(0), code)


# ── 自動探查輸入檔結構 → 餵給 code-gen，杜絕模型腦補欄位/標籤(grounding、模型無關)──
# batch 裡引用到的「已存在」資料檔(輸出檔還沒生成、會自動跳過),在 skill 生 code 前
# 把真實 schema(Excel 欄位 / CSV 表頭 / PPT/Word 佔位標籤 / JSON keys)攤給模型看,
# 連弱模型(Gemma)也不會把欄位腦補成「產品/價格」「姓名/照片」。輕量:最多探 3 檔、
# 大檔跳過、read_only 只讀表頭+前幾列。
_INPUT_PEEK_RE = re.compile(r"[^\s\"'`\n,;()\[\]<>]+\.(?:xlsx|xls|csv|json|pptx|docx)", re.IGNORECASE)
_INPUT_PEEK_MAX_FILES = 3
_INPUT_PEEK_MAX_BYTES = 30 * 1024 * 1024  # 單檔 > 30MB 跳過(避免拖慢)
_TAG_RE = re.compile(r"\{\{[^{}\n]{1,40}?\}\}")


def _summarize_input_file(p: "Path") -> str:
    ext = p.suffix.lower()
    nm = p.name
    try:
        if ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
            parts = []
            for sn in wb.sheetnames[:3]:
                ws = wb[sn]
                rows = list(ws.iter_rows(max_row=4, values_only=True))
                if not rows:
                    parts.append(f"  工作表「{sn}」: (空)"); continue
                hdr = [str(c) if c is not None else "" for c in rows[0]]
                samp = ["、".join(str(c) if c is not None else "" for c in r) for r in rows[1:4]]
                parts.append(f"  工作表「{sn}」欄位: {', '.join(hdr)}" +
                             (("\n    前幾列: " + " | ".join(samp)) if samp else ""))
            wb.close()
            return f"Excel {nm}:\n" + "\n".join(parts)
        if ext == ".csv":
            import csv as _csv
            with open(p, encoding="utf-8", errors="replace") as fh:
                rows = []
                for i, r in enumerate(_csv.reader(fh)):
                    rows.append(r)
                    if i >= 3:
                        break
            if not rows:
                return f"CSV {nm}: (空)"
            return (f"CSV {nm} 欄位: {', '.join(rows[0])}" +
                    (("\n  前幾列: " + " | ".join("、".join(r) for r in rows[1:4])) if len(rows) > 1 else ""))
        if ext == ".json":
            import json as _json
            d = _json.loads(p.read_text(encoding="utf-8", errors="replace"))
            if isinstance(d, dict):
                return f"JSON {nm} keys: {', '.join(list(d.keys())[:20])}"
            if isinstance(d, list) and d and isinstance(d[0], dict):
                return f"JSON {nm}(陣列、每筆 keys): {', '.join(list(d[0].keys())[:20])}"
            return f"JSON {nm}: {type(d).__name__}"
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(p))
            slides = list(prs.slides)
            tags = set()
            for s in slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        tags.update(_TAG_RE.findall(sh.text_frame.text))
            return (f"PPT {nm}: {len(slides)} 頁、"
                    f"佔位標籤: {', '.join(sorted(tags)) if tags else '(無 {{}} 標籤)'}")
        if ext == ".docx":
            from docx import Document
            doc = Document(str(p))
            tags = set()
            for par in doc.paragraphs:
                tags.update(_TAG_RE.findall(par.text))
            return f"Word {nm}: 佔位標籤: {', '.join(sorted(tags)) if tags else '(無 {{}} 標籤)'}"
    except Exception as e:
        return f"{nm}: (讀取結構失敗、{e.__class__.__name__})"
    return ""


def _peek_input_files(task_description: str, logger) -> str:
    seen: set = set()
    summaries: list = []
    for raw in _INPUT_PEEK_RE.findall(task_description or ""):
        cand = raw.strip()
        wm = re.match(r"^/mnt/([a-z])/(.*)$", cand)  # 沙盒 linux 路徑 → Windows
        if wm:
            cand = f"{wm.group(1).upper()}:\\" + wm.group(2).replace("/", "\\")
        try:
            p = Path(cand).expanduser()
        except Exception:
            continue
        key = str(p).lower()
        if key in seen:
            continue
        seen.add(key)
        try:
            if not (p.exists() and p.is_file()):
                continue  # 不存在 = 多半是「待產出的輸出檔」→ 跳過
            if p.stat().st_size > _INPUT_PEEK_MAX_BYTES:
                continue
        except Exception:
            continue
        s = _summarize_input_file(p)
        if s:
            summaries.append(s)
        if len(summaries) >= _INPUT_PEEK_MAX_FILES:
            break
    if not summaries:
        return ""
    blob = "\n".join(summaries)
    if len(blob) > 1800:
        blob = blob[:1800] + "…(截斷)"
    try:
        logger.info(f"[input-peek] 已注入 {len(summaries)} 個輸入檔的實際結構")
    except Exception:
        pass
    return ("\n\n【📂 輸入檔的實際結構(以下是程式實際讀到的、是真的)】\n" + blob +
            "\n→ 請**完全依照上面的真實欄位 / 標籤名稱**填寫,絕對不可自行假設或改寫"
            "(例:不可把欄位腦補成「產品/價格/規格」或「姓名/職稱/照片」等常見模板)。")


async def execute_step_with_skill(
    task_description: str,
    timeout: int,
    logger: logging.Logger,
    step_name: str,
    output_path: Optional[str] = None,
    working_dir: Optional[str] = None,
    prev_outputs: Optional[list] = None,
    pipeline_id: Optional[str] = None,
    use_recipe: bool = True,
    no_save_recipe: bool = False,
    readonly: bool = False,
    run_id: str = "",
    previous_failures: Optional[list] = None,
    recipe_step_key: Optional[str] = None,
    skill_name: str = "",
    ask_mode: bool = False,
    silent_recipe: bool = False,
    llm_role: str = "primary",
    has_external_validator: bool = False,  # runner 傳:此 step 之後有沒有 AI validator 跑
) -> ExecResult:
    """
    Skill 模式執行器：LLM 解讀自然語言任務描述，自主撰寫並執行程式碼。

    Args:
        task_description: 自然語言任務描述（取代 shell 命令）
        timeout:          最大執行秒數（整體 agent 迴圈）
        logger:           file logger
        step_name:        步驟名稱
        output_path:      預期輸出路徑（可選，讓 agent 知道要把結果存在哪）
        prev_outputs:     前幾步的輸出檔案資訊列表，格式 [{"path": "...", "schema": "..."}]
    """
    # 從 step.timeout 推導「單次 run_python / run_shell 上限」
    # — step.timeout=600 (skill 預設) → 120s tool_timeout
    # — step.timeout=1200 (爬蟲類) → 180s tool_timeout（封頂）
    # 寫進 prompt 給 LLM 看，讓它一開始就知道單次工具呼叫的限制、會自己用 asyncio 併發或分批
    tool_timeout = _compute_tool_timeout(timeout)
    logger.info(f"[{step_name}] 工具呼叫單次上限：{tool_timeout}s（從 step.timeout={timeout}s 推導）")

    # 展開 ~ 為完整路徑
    if output_path:
        output_path = str(Path(output_path).expanduser())
    # 自動建立輸出路徑的父目錄和工作目錄
    if output_path:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        # 刪除舊的 output 檔案，避免 done guard 被上次執行的殘留檔案騙過
        # Windows 上使用者若開著 Excel 檢視輸出，unlink 會 PermissionError
        # → 改成 rename 成 .stale-<timestamp>.bak 繞過鎖，並用 try/except 確保整個 step 不會因此卡死
        _out = Path(output_path)
        if _out.exists():
            logger.info(f"[{step_name}] 刪除舊輸出檔案：{output_path}")
            try:
                _out.unlink()
            except PermissionError as _e:
                import time as _t
                _bak = _out.with_suffix(_out.suffix + f".stale-{int(_t.time())}.bak")
                try:
                    _out.rename(_bak)
                    logger.warning(
                        f"[{step_name}] 舊輸出檔案被佔用（可能你在 Excel 開著），"
                        f"已改名為 {_bak.name} 讓這次執行繼續。請關閉 Excel 後手動清掉 .bak 檔。"
                    )
                except Exception as _e2:
                    # rename 也失敗（極少見，通常是檔案被獨佔）→ 讓使用者知道但不中斷
                    logger.warning(
                        f"[{step_name}] 無法刪除或改名舊輸出檔案（{_out.name}）：{_e2.__class__.__name__}。"
                        f"通常是 Excel / 其他程式正打開此檔。LLM 寫入時可能也會失敗，請先關閉該檔案再重跑。"
                    )
            except Exception as _e:
                logger.warning(f"[{step_name}] 刪除舊輸出檔案時發生錯誤：{_e}")
    if working_dir:
        Path(working_dir).mkdir(parents=True, exist_ok=True)

    logger.info(f"[{step_name}] 🔬 Skill 模式啟動：{task_description}")

    # ── Recipe Book：檢查是否有可重用的成功快取 ────────────────────────────
    _rkey = recipe_step_key or step_name  # recipe DB key（含索引，避免同名覆蓋）
    input_paths = [po["path"] for po in (prev_outputs or []) if po.get("path")]
    if pipeline_id and use_recipe:
        try:
            from db import get_recipe, match_recipe, save_recipe, mark_recipe_failed
            from pipeline.recipe import _sha1 as _recipe_sha1, _fingerprint_input as _recipe_fp
            # Debug: 先載入 recipe 看詳細匹配狀況
            _raw = get_recipe(pipeline_id, _rkey)
            if _raw:
                _cur_hash = _recipe_sha1(task_description)
                _cur_fps = {p: _recipe_fp(p) for p in input_paths}
                saved_fps = _raw["input_fingerprints"]
                if isinstance(saved_fps, str):
                    import json as _json
                    saved_fps = _json.loads(saved_fps)
                if _raw["disabled"]:
                    logger.info(f"[{step_name}] 📖 Recipe 存在但已停用")
                elif _raw["task_hash"] != _cur_hash:
                    logger.info(f"[{step_name}] 📖 Recipe 存在但 task_hash 不符（saved={_raw['task_hash']}, current={_cur_hash}）")
                elif _cur_fps != saved_fps:
                    logger.info(f"[{step_name}] 📖 Recipe 存在但輸入指紋不符")
                    for k in set(list(_cur_fps.keys()) + list(saved_fps.keys())):
                        sv = saved_fps.get(k, '(無)')
                        cv = _cur_fps.get(k, '(無)')
                        if sv != cv:
                            logger.info(f"[{step_name}]   {k}: saved={sv} → current={cv}")
            else:
                logger.debug(f"[{step_name}] 📖 無 Recipe 紀錄")
            _fp = {p: _recipe_fp(p) for p in input_paths}
            cached = match_recipe(pipeline_id, _rkey, _recipe_sha1(task_description), _fp)
            # 互動式 recipe(過程用過 ask_user)不直接重播:快取碼把上次的人工選擇寫死在裡面,
            # 直接 replay 等於拿舊答案、永遠不再問使用者 → 使用者這次的選擇「沒被記錄/沒被使用」。
            # 退回 LLM 重跑,讓 ask_user 重新觸發、收這次的回答。was_interactive 一直有存、
            # 但 replay 從來沒檢查它(稽查 ask_user/recipe 問題的 root cause)。
            if cached and cached.get("was_interactive"):
                logger.info(
                    f"[{step_name}] 📖 命中互動式 recipe(含 ask_user 回答),不直接重播 → "
                    f"退回 LLM 重跑以重新詢問使用者"
                )
                cached = None
            if cached:
                logger.info(
                    f"[{step_name}] 📖 找到快取 recipe (成功 {cached['success_count']} 次, "
                    f"平均 {cached['avg_runtime_sec']:.1f}s)，跳過 LLM 直接執行"
                )
                import time as _time
                t0 = _time.time()
                loop = asyncio.get_event_loop()
                # 修：原本直呼叫 _skill_run_python（host subprocess）會把 /mnt/c/... 沙盒路徑
                # 丟到 Windows 跑 → 找不到檔 → recipe 一直 fail 重學。
                # 改走 _try_sandbox_exec 先判斷沙盒可用性，可用就用沙盒（原本錄製就是在沙盒）；
                # 沙盒不可用才退 host（這時 LLM 重學會學到 host 路徑、新 recipe 自洽）
                def _replay_recipe():
                    # B′ per-run 夾相容:快取碼寫死了錄製當時的 run_<ts> 子夾,重播前換成本次的
                    # run_<ts>(從 output_path 取),否則快取碼寫到舊夾、本次夾空 → output_path
                    # 不存在 → recipe 被誤判失敗、退回 LLM、0 成本失效。詳見 _rebase_recipe_run_dir。
                    code_to_run, _n = _rebase_recipe_run_dir(cached["code"], output_path)
                    if _n:
                        logger.info(f"[{step_name}] 🔁 Recipe 重播:快取碼內 {_n} 處舊 run 夾 → 本次 run 夾")
                    sandbox_out = _try_sandbox_exec(
                        "run_python", code_to_run, working_dir, run_id, logger,
                        tool_timeout=tool_timeout,
                    )
                    if sandbox_out is not None:
                        return sandbox_out
                    return _skill_run_python(code_to_run, cwd=working_dir, run_id=run_id,
                                              tool_timeout=tool_timeout)
                tool_result = await loop.run_in_executor(None, _replay_recipe)
                runtime = _time.time() - t0
                # 成功條件：輸出檔存在（若有指定）且無 [exit code: X]
                ok = "[exit code:" not in tool_result
                if ok and output_path:
                    ok = Path(output_path).exists()
                if ok:
                    import sys as _sys
                    save_recipe(pipeline_id, _rkey, _recipe_sha1(task_description),
                                _fp, output_path, cached["code"],
                                f"{_sys.version_info.major}.{_sys.version_info.minor}", runtime)
                    logger.info(f"[{step_name}] ✅ Recipe 重跑成功（{runtime:.1f}s）")
                    return ExecResult(exit_code=0, stdout=tool_result, stderr="__RECIPE_HIT__")
                else:
                    logger.warning(f"[{step_name}] Recipe 重跑失敗，改用 LLM 重新學習。輸出：{tool_result[:300]}")
                    mark_recipe_failed(pipeline_id, _rkey)
        except Exception as e:
            logger.warning(f"[{step_name}] Recipe 檢查失敗：{e}")
    # ───────────────────────────────────────────────────────────────────────

    # 注入當前日期/時間 — 避免 LLM 的 training cutoff 造成「2026 年還沒到」之類誤判
    # 用 host 本地時間（TZ=Asia/Taipei 之類由系統決定）；skill 任務都是跟使用者同時區
    from datetime import datetime as _dt
    _now = _dt.now()
    _date_block = (
        f"【當前日期時間（host system 時鐘）】\n"
        f"  {_now.strftime('%Y-%m-%d %H:%M:%S')}（週{'一二三四五六日'[_now.weekday()]}）\n"
        "  使用者提到「今天」、「最新」、「本月」、「Q1」等相對時間時，以上面這個日期為準，"
        "不要以你訓練資料的時間為準。若使用者指定的年份早於或等於目前年份，那是真實存在可查的時間，"
        "不要回覆「尚未到達」、「無法獲取」等錯誤判斷。\n"
    )

    system_prompt = _date_block + """你是 pipeline Skill 執行 agent。根據任務描述自主寫程式並跑、回繁體中文。

工具(每次回覆只能呼叫一個、用 <tool>名稱</tool><input>...</input> 格式):

1. run_python — Python 程式碼(在工作目錄執行)
   <input>print("hello")</input>
2. run_shell — Linux 命令(優先 run_python、Python 跨平台較穩)
   <input>wc -l output.csv</input>
3. read_file — 讀檔(路徑不加引號)
   <input>path/to/file.txt</input>
   單次最多 100 行 / 24KB,超過會截斷。**大檔分段讀**:用 JSON 帶 offset
   <input>{"path": "path/to/big.py", "offset": 100, "limit": 100}</input>
   截斷訊息會告訴你下一段的 offset、照著帶就能接著讀完整檔。
   ✓ 用於:讀上一步輸出 / 樣本檔 / 設定檔(餵推理用)
   ✗ 不要用於:「驗證自己剛 write 的檔存在」(Python exit 0 已證、用 Path.exists() 在 run_python 內就行)
   ✗ 不要用於:「再確認一次自己剛寫的內容對不對」(交給外部 validator、不要重複工作浪費 iter)
4. write_file — **整檔寫入**(LLM 已知 content、直接寫,**避免用 Python 三引號包大段外語碼**)
   **短內容(< 5KB)用 JSON**:
   <input>{"path":"output.js","content":"const x = ...整段內容..."}</input>
   注意 JSON content 內 " 要 escape 為 \\"、換行用 \\n
   **長內容(> 5KB、含特殊字元、不想 escape)用 RAW MULTILINE**(推薦):
   <input>
   path: /abs/path.js
   ---
   const x = `中文直接寫`;  // 完全免 escape、可自由用引號 / 反斜線 / 換行 / unicode
   const y = "也可以放雙引號";
   </input>
   ✓ 用於:寫 JS/HTML/CSS/SQL/markdown/設定檔等「LLM 知道全部內容、要原樣寫到檔」
   ✗ 不要用於:寫 pandas DataFrame、計算結果(用 run_python + to_csv 才對)
5. edit_file — **局部替換**(部分修改既有檔案、不整檔重寫)
   <input>{"path":"x.py","old_text":"return None","new_text":"return result"}</input>
   ✓ 用於:換一行、改一個函式、修一個 bug
   要求:old_text 在檔內**唯一**(否則加上下文讓它唯一、或加 "replace_all": true)
6. grep — **搜尋字串 in files**
   <input>{"pattern":"import\\\\s+pandas","path":".","glob":"*.py"}</input>
   ✓ 用於:找 keyword、找 imports、找 TODO、找錯誤訊息
   選填:glob(*.py 之類過濾)、max_results(預設 50)、case_sensitive(預設 true)
7. glob — **列檔案**(找符合 pattern 的檔)
   <input>{"pattern":"*.csv","path":"data/"}</input>
   ✓ 用於:看資料夾有哪些檔、找特定副檔名
   pattern 用 pathlib 規則:`*.py` / `**/*.csv` / `report_*.md`
8. view_image — 看圖(png/jpg/gif/webp/bmp、上限 20MB)。驗證圖表 / 從圖擷取資訊用。
   模型不支援視覺時直接 done(success=false) 並在 error 註記。
   <input>path/to/chart.png</input>
9. ask_user — 不確定 / 模糊 / 高風險時**優先用、不要硬猜**:任務歧義(欄位 / 格式 / 路徑)、
   覆寫 / 刪除使用者檔、外部 API、多種合理做法、環境狀態不確定。
   <input>{"question":"輸出哪種格式?","options":["PDF","Word","MD"],"context":"資料 120 筆"}</input>
   question 必填(中文、可一次多題)、options 選填(陣列 → UI 顯示按鈕)、context 選填。
   使用者回 → 工具回傳「使用者回答：<答案>」、再依答案繼續。逾時 / 取消 → 用合理預設或 done(success=false)。
10. export_var — 把這步算好的一個值傳給下游節點(尤其 condition 條件節點)。
   <input>{"name":"score","value":380}</input>
   ✓ 用於:這步算出一個數字 / 字串、下游 condition 要拿它判斷分支時。
   呼叫後下游就能用 {{ steps.<本步驟名>.output.score }} 引用。
   name 用英文、value 數字或字串皆可;可多次呼叫匯出多個值。
   ✗ 不用於:傳整份資料(那用輸出檔 + {{ steps.X.output.path }})。
11. done — 完成回報。
   成功:<input>{"success":true,"summary":"完成了什麼"}</input>
   失敗(已窮盡所有可用工具與方法後):<input>{"success":false,"error":"已試 X/Y/Z 各自失敗原因","missing_packages":["套件A"]}</input>
   - 必須先試完已安裝套件的所有替代方案才呼叫 success=false
   - missing_packages 只填**未安裝**、裝後有合理機率解的套件
   - 不要第一方案失敗就放棄、要切策略再試

【可用 Python 套件】
標準庫:csv, json, random, os, pathlib, re, math, datetime, io, collections, itertools, functools, glob, shutil, hashlib, urllib
已安裝第三方:{installed_packages}
不在清單上的也可 import、系統會偵測缺失提示安裝。

【規則(必守)】
- 嚴格遵守任務指定的欄位名 / 路徑 / 數值範圍、不得自改
- 一律絕對路徑(根據工作目錄 + 輸出路徑 hint)、用 pathlib.Path 或 os.path.join、不要字串拼 `/`
- 只用上方套件、絕對不要 sudo / pip install / apt install
- 跑其他 Python script 用 sys.executable、不要寫死 python3 / python(PATH 不一定對)
  正確:`subprocess.run([sys.executable, "script.py"], ...)`
- 隨機資料須唯一(姓名等):先用集合生成不重複組合再 random.sample、不要迴圈內 random.choice 累積
- ❌ 禁止 hardcode 結果(摘要 / 情緒 / 分類):必須讀實際資料 + 程式邏輯處理、不可在原始碼寫死答案 dict / list、會被驗證階段抓
- 🧠 **你自己就是 LLM**:任務說「用 LLM 校對 / 翻譯 / 分類 / 摘要 / 改寫」時,那個 LLM 就是你。
  **小批就自己一次做完**(數十筆內、預估輸出 < 約 8000 字 = 一次塞得下 → read_file 整批讀進來 →
  在同一次回應裡產出全部處理後內容 → write_file 一次寫出;**別拆成逐筆**,5 行也逐筆呼叫純浪費)。
  **絕不要** `import openai` / `import anthropic` / `from google ... genai` 呼叫外部 LLM API
  (沒金鑰、會 rc=1 卡死、繞過防呆)。**只有大批**(一次塞不下 / 輸出超過約 8000 字 / 上看百筆)自己做不完,
  才用內建 helper(走系統現在這顆模型、免金鑰、已在 PYTHONPATH):`from skill_llm import llm` →
  `out = llm("提示詞" + text)`(可加 `system=...`);失敗會 raise。**仍然不要 import openai**。
- 嚴禁 input() / getpass() / sys.stdin.read() — pipeline 非互動環境、會永久卡死
- 任務需選擇:優先用任務指定;無指定 → 最合理預設值 + summary 註明假設;只有「會嚴重影響結果(覆蓋重要檔、無法回復)」才用 done(success=false)
- 讀別的步驟產出的檔(csv / xlsx / json / md 等)第一步先 run_python 看前幾行,確認實際的欄位名、值的格式(大小寫、引號、數字或字串),或文字檔的結構(例:第一行是不是標題)、不要憑空假設
- 重試:絕不重複同一方法、回顧歷史、用尚未嘗試的不同套件 / 策略;耗盡才 done(success=false) + missing_packages

【工具呼叫格式(最重要)】
- 每次回覆只一個 tool call、所有程式碼放 <tool> + <input> 標籤內
- 禁止 markdown ``` 區塊展示程式碼:回覆中不應出現 ``` 符號
- 主處理邏輯集中在一個 run_python、別無謂拆步;但讀「別的步驟產出的檔」時,先 peek 看前幾行確認格式、再寫處理 —— 這個 peek 是必要的(見上)、不算拆步
- 絕不在 Python 程式碼裡呼叫 done(...) / view_image(...) / read_file(...) — 那是 tool 名、不是函式
- 程式跑成功 → 下回覆直接 done"""

    # 網路搜尋工具：僅在 settings.web_search_enabled AND 有 tavily_api_key 時對 agent 揭露
    # 沒啟用就完全不提（agent 連這工具名都看不到，不會誤呼叫）
    try:
        import sys as _sys3
        _backend_dir3 = str(Path(__file__).parent.parent.absolute())
        if _backend_dir3 not in _sys3.path:
            _sys3.path.insert(0, _backend_dir3)
        from settings import get_settings as _gs_ws
        _ws_settings = _gs_ws()
        if _ws_settings.get("web_search_enabled") and (_ws_settings.get("tavily_api_key") or "").strip():
            # FIRST_ITER 區塊:LLM 第一輪看過後內化、後續輪數會被 strip 省 token
            system_prompt += r"""

<!--FIRST_ITER_BEGIN-->
【🔍 工具 7：web_search】
Tavily 搜網、結果回對話。**不是每個任務都要搜**:
- ✅ 即時資訊(股價 / 新聞 / 匯率)、使用者提「查」「最新」、缺背景知識、確認套件 / API 最新做法
- ❌ 純資料處理、任務已給完整資料、為「驗證想法」亂搜(先動手)

<input>{"query":"今天美國科技新聞","max_results":5,"include_full_content":true}</input>
- max_results: 1-5(預設 5)
- search_depth: 不必傳、系統會用設定頁的預設(預設 advanced、深度好、適合研究類任務)
- include_full_content=true → 拿每則完整原文(~3000 字/篇)

⭐ 任務要「擷取內文 / 分析全文」時直接 `include_full_content=true`、別寫 requests/newspaper 自己爬
   (Tavily 已處理 CF / JS 渲染 / 反爬、自爬幾乎一定 403 或拿不到 SPA 內容)

⚠️ 每步驟最多搜 5 次($0.01-0.025/次)、整合後再 query
⚠️ include_full_content=true 回傳 ~15000 字、只在確需才開
<!--FIRST_ITER_END-->"""
            logger.info(f"[{step_name}] 🔍 web_search 工具已啟用（Tavily）")
    except Exception as _e:
        logger.debug(f"[{step_name}] web_search 工具注入失敗（略過）：{_e}")

    # 唯讀模式：注入禁止修改的約束
    if readonly:
        system_prompt += """

【🔒 唯讀驗證模式】
此步驟為「唯讀深度驗證」，你的職責是：
- **只能讀取、分析、檢查檔案內容**
- **嚴禁修改、覆寫、重新命名任何檔案或欄位**
- **嚴禁用程式碼「修正」資料來通過驗證**
- 如果檢查結果不符合預期 → 直接用 done 回報 success=false 並說明哪裡不符
- 如果檢查結果符合預期 → 用 done 回報 success=true 並說明驗證通過的理由
- **你只是驗證者，不是修復者**"""
        logger.info(f"[{step_name}] 🔒 唯讀驗證模式已啟用")

    # 詢問模式：鼓勵 LLM 遇到任何模糊處就用 ask_user 主動問使用者
    # 預設（未勾選）：保守使用 ask_user，優先靠任務描述 + 合理預設完成任務
    # 勾選後：把「遇到不確定就問」的優先度拉到最高，減少 LLM 自己猜的情況
    if ask_mode:
        # 先覆寫 base 裡跟詢問模式相衝突的「優先用預設值」那行,避免 LLM 拿到兩條矛盾指令
        system_prompt = system_prompt.replace(
            "- 任務需選擇:優先用任務指定;無指定 → 最合理預設值 + summary 註明假設;只有「會嚴重影響結果(覆蓋重要檔、無法回復)」才用 done(success=false)",
            "- 任務需選擇:**一律優先用 ask_user 問使用者**(詢問模式 ask_user 無次數上限);只有使用者已明確指定 / 唯一明顯答案才用預設",
        )
        system_prompt += """

【❓ 詢問模式已啟用】
你**最優先**的工具是 `ask_user`，不是 `run_python`。下面任何一項吻合就必須用 ask_user，不得自行推論：
- 任務描述有模糊處（欄位名、輸出格式、數值範圍、是否覆蓋檔、要不要 dry-run…）
- 有多種合理做法 → 列成 options 讓使用者選
- 要動到關鍵檔案 / 覆蓋既有資料 / 呼叫外部 API / 花錢或耗時的操作
- 環境狀態不確定（例：沙盒是否可用、套件有無安裝、預期檔案是否存在）
- 第一次嘗試失敗、在選下一種做法前（先問「要繼續試其他套件還是放棄」）
**判斷原則反過來**：base prompt 預設「能推論就不問」，詢問模式下改成「**有任何疑慮就問**」。
**詢問模式下 ask_user 沒有次數上限**（原本限制 3 次，此模式下取消），請放心多問幾次。
每個 ask_user 可以同時包 1 題或多題相關問題（用換行或編號）一次收齊，減少往返。

【🚀 啟動既有 Python 專案的特別規則】
任務描述含「啟動 / 跑 / 執行」+ 一個 .py 檔案路徑(尤其是 main.py / run.py / app.py / start.py)時:
- 在 subprocess 啟動子程序**之前**,**先 read_file 看一遍源碼**(至少前 200 行)
- 重點找:
  - `input(...)` 互動輸入呼叫
  - `argparse` / `sys.argv` 命令列參數
  - 環境變數讀取(`os.environ` / `os.getenv`)
  - `print` 出來的 menu/選單(會跟 input 配對)
- **發現 input() 或選單** → 用 ask_user 問使用者該回什麼,再用 stdin pipe 餵進子程序:
  ```python
  result = subprocess.run(
      [sys.executable, "main.py"],
      input="user_answer_1\nuser_answer_2\n",  # 對應 input() 順序
      text=True, capture_output=True, timeout=120
  )
  ```
- **發現需要 argv 參數** → 用 ask_user 問參數值,組進 subprocess 命令
- **絕對不要**直接 `subprocess.run([sys.executable, "main.py"])` 然後等 timeout、再去猜為什麼卡住
原因:子程序 input() 會 block stdin 直到 timeout(60-120 秒),純粹浪費時間;先讀源碼 + 主動問用戶才是專業做法。"""
        logger.info(f"[{step_name}] ❓ 詢問模式已啟用（LLM 遇到模糊處會主動問使用者；ask_user 無上限）")

    # ── 沙盒環境提示（僅在 wsl_docker 模式注入）──
    # Host 模式在 Windows 上跑，agent 用 Windows 路徑 / win32com 都 OK；
    # wsl_docker 模式在 Linux 容器，需要告訴 agent「你不在 Windows」避免浪費迭代
    # （實測 agent 常犯：用 C:\ 路徑、呼叫 win32com、以為有 PowerShell 等）
    try:
        import sys as _sys2
        _backend_dir2 = str(Path(__file__).parent.parent.absolute())
        if _backend_dir2 not in _sys2.path:
            _sys2.path.insert(0, _backend_dir2)
        from settings import get_settings as _get_settings_for_sandbox
        if (_get_settings_for_sandbox().get("skill_sandbox_mode") or "host").strip() == "wsl_docker":
            try:
                _v5_root_win = Path(__file__).parent.parent.parent.absolute()
                _drive = str(_v5_root_win)[0].lower()
                _rest = str(_v5_root_win)[3:].replace("\\", "/")
                _v5_root_wsl = f"/mnt/{_drive}/{_rest}"
            except Exception:
                _v5_root_wsl = "/mnt/c/" + Path(__file__).resolve().parents[2].name
            # FIRST_ITER 區塊:沙盒環境規則第一輪講完、後續輪數 strip 省 ~70 行 / 輪
            system_prompt += rf"""

<!--FIRST_ITER_BEGIN-->
【🛡️ Sandbox 環境(Linux Docker 容器、python:3.13-slim、不是 Windows)】
- OS = Linux:沒有 win32com / pywin32 / PowerShell / cmd.exe — 用純 Python 或 Linux 工具
- 產 PPT 用 `python-pptx`(首選)或 Node.js + `pptxgenjs`(走 `.agents/skills/pptx`);**不要 import win32com.client**(永遠 ImportError)
- Node.js 全域套件已預裝(`pptxgenjs`、`docx`、可能還有其他 npm 包)、NODE_PATH 已自動對齊
  → **JS 檔內 require 用裸模組名**:`require('pptxgenjs')`、`require('docx')`(NODE_PATH 自動解析、不必寫絕對路徑)
  → **絕對禁止** `require('C:/...')`、`require('/mnt/c/.../node_modules/pptxgenjs')`(那是 Windows / 主機路徑、container 內找不到、必定 Cannot find module)
  → **絕對禁止** `find / -name "..."`(掃整碟 > 60s timeout、會 SIGTERM、且 NODE_PATH 已對齊不需要找)
- 路徑轉換:Windows `C:\...` / `C:/...` 在容器無效、pathlib 會當相對路徑導致找不到檔
  - `C:\Users\X\...` → `/mnt/c/Users/X/...`
  - `D:\data\...` → `/mnt/d/data/...`
  - 容器 `~`(`/root`)mount `.agents`、`Path.home()/".agents"` 跟 `/mnt/c/Users/X/.agents` 同一份
- PATH 只有 Linux 工具(node/npm/python3/bash/ls/grep/curl);沒 where/dir/type/copy
- 任務給 Windows 路徑 → 自動轉 `/mnt/<drive>/...` 再用

【⛔ 兩個 mount 不要混淆 — 寫產物搞錯位置 = step 找不到產物 fail】
容器有兩個獨立 mount 各管不同事:

A. **專案目錄**(寫 workflow 產物的地方、絕大多數情境):
   容器內:`{_v5_root_wsl}/`
   host 對應:`{_v5_root_win}`
   workflow 產物寫到 `{_v5_root_wsl}/ai_output/<workflow_name>/<檔名>`
   ✅ 範例:`/mnt/d/Atlas/pipeline-orchestratorV5/ai_output/sales_q1_analysis/report.md`

B. **Skill / Agent 目錄**(讀 only、skill 自身程式碼住的地方):
   容器內:`/root/.agents/`
   host 對應:`~\.agents\`
   裡面是 SKILL.md / scripts / references — **不是 workflow 產物存放區**

⛔ **絕對不要把 workflow 產物寫到 `/root/.agents/ai_output/<xxx>/`**!
   那個路徑你可能看得到資料夾(因為 .agents/ai_output 巧合也存在)、
   但**不對應到 host 的 ai_output**、Pipeline runner 看不到、step 標 fail。
   這是踩過的真實坑、必須記牢。

✅ **正確寫法**:`output_path` 提示給的是哪個容器內絕對路徑、就寫到那、不要自作主張改路徑。

【📁 專案根目錄】`{_v5_root_wsl}`
- 任務裡相對路徑(`external_projects/...` / `scripts/...` / `docs/...`)以**專案根**展開、不要拿 sandbox CWD
- 例:`external_projects/interactive_demo/main.py` → `{_v5_root_wsl}/external_projects/interactive_demo/main.py`
- 不確定位置直接 `list(Path("{_v5_root_wsl}").rglob("檔名"))`、不要一層層 ls

【🌐 網頁抓取(容器內)】
- **一律用 `crawl4ai`**(容器已預裝 + playwright + chromium、處理 JS / CF / cookies / 滾動 / 轉 md)
- ❌ 禁用 selenium(容器沒 chromedriver) / requests/urllib/httpx 直接 GET HTML(SPA 拿不到) / 自寫 playwright(繞過反爬)

最小樣板(直接抄改):
```python
import asyncio
from crawl4ai import AsyncWebCrawler, BrowserConfig, CrawlerRunConfig, CacheMode

async def fetch(url, scroll=True):
    js = ""
    if scroll:
        js = '''
        await new Promise(r => setTimeout(r, 2000));
        window.scrollTo(0, document.body.scrollHeight / 2);
        await new Promise(r => setTimeout(r, 1500));
        window.scrollTo(0, document.body.scrollHeight);
        await new Promise(r => setTimeout(r, 1500));
        '''
    async with AsyncWebCrawler(config=BrowserConfig(headless=True)) as c:
        r = await c.arun(url=url, config=CrawlerRunConfig(
            cache_mode=CacheMode.BYPASS,
            page_timeout=90000,
            wait_until="domcontentloaded",  # SPA 不能用 networkidle、會 timeout
            js_code=js if js else None,
        ))
        return r.markdown if r.success else None

md = asyncio.run(fetch("https://example.com/"))
```

SPA 站(Reddit/Twitter/X/Instagram/Threads/Bluesky):`wait_until="domcontentloaded"` + js_code 加滾動 + regex 過濾貼文 URL pattern(別拿頁首導覽連結)。
影片下載用 `yt-dlp`(容器預裝)、不要 pytube / youtube-dl。
<!--FIRST_ITER_END-->"""
            # 動態注入「單次 tool 上限秒數」— 從 step.timeout 推導出來、不是寫死 60s
            system_prompt += f"""

<!--FIRST_ITER_BEGIN-->
【⏱ 單次 tool call 上限 {tool_timeout} 秒(超過會被強制終止)】
- 多筆網路 I/O(抓 N 網頁) → asyncio.gather + asyncio.Semaphore(5) 限併發
- 大量 LLM 序列 → 拆多次 run_python(每次 batch 的 1/3)
- CPU 密集 → pandas / numpy 向量化、避免 Python for-loop
- ❌ 一個 run_python 裡 sequential 跑 10+ 慢操作(每個 5-10s、總和會超)
- ❌ run_python 內 sleep / wait / poll loop 等外部事件(直接超時)
- 看到「[錯誤] 執行超時」 → 立即改策略重寫、別重送同份程式
<!--FIRST_ITER_END-->"""
            logger.info(f"[{step_name}] 🛡 已注入 wsl_docker sandbox 環境資訊")
    except Exception as _e:
        logger.debug(f"[{step_name}] sandbox env 注入失敗（略過）：{_e}")

    # 掛載 skill：注入 SKILL.md 內容與子資源清單
    if skill_name:
        try:
            from skill_scanner import get_skill_prompt_injection
            skill_injection = get_skill_prompt_injection(skill_name)
            if skill_injection:
                system_prompt += skill_injection
                logger.info(f"[{step_name}] ✨ 已掛載 Skill：{skill_name}")
            else:
                logger.warning(f"[{step_name}] ⚠️ 找不到 Skill：{skill_name}（已略過）")
        except Exception as e:
            logger.warning(f"[{step_name}] ⚠️ 載入 Skill {skill_name} 失敗：{e}")

    # output_path / working_dir hint:在 sandbox 模式下、自動翻 Windows path → 容器內 /mnt/<drive>/...
    # 防 LLM 在沙盒內看到 D:\ 想破頭、或誤推「/root/.agents/ai_output/...」(那是 skill 目錄、不對應 host)
    def _to_wsl_path_for_hint(p: str) -> str:
        if not p:
            return p
        try:
            from settings import get_settings
            if (get_settings().get("skill_sandbox_mode") or "host").strip() != "wsl_docker":
                return p  # host 模式直接用原 path
            import re as _re
            m = _re.match(r"^([A-Za-z]):[\\/](.+)$", p)
            if m:
                drive = m.group(1).lower()
                rest = m.group(2).replace("\\", "/")
                return f"/mnt/{drive}/{rest}"
        except Exception:
            pass
        return p

    _hint_output = _to_wsl_path_for_hint(output_path) if output_path else None
    _hint_wd = _to_wsl_path_for_hint(working_dir) if working_dir else None
    output_hint = (
        f"\n輸出路徑提示:請將結果存到 `{_hint_output}`(這是**容器內絕對路徑**、寫到這 host 端會直接看到、不要改別處)"
        if _hint_output else ""
    )
    wd_hint = (
        f"\n工作目錄(容器內絕對路徑):`{_hint_wd}` (所有相對路徑相對於此、但建議直接用絕對路徑)"
        if _hint_wd else ""
    )

    # 組合前步驟的輸出資訊 —— 附上「實際內容樣本」,讓 LLM 一開始就看到真實的
    # 欄位名 / 值的格式(大小寫…),不用自己猜(猜錯又不當機就會靜默產出爛結果)。
    prev_hint = ""
    if prev_outputs:
        lines = ["\n【前步驟產生的檔案（可直接讀取使用）】"]
        for po in prev_outputs:
            lines.append(f"- {po['path']}")
            if po.get("schema"):
                lines.append(f"  欄位/結構：{po['schema']}")
            sample = _read_file_sample(po.get("path", ""))
            if sample:
                lines.append("  ↓ 實際內容樣本（欄位名、值的大小寫一律以此為準、不要自己猜）：")
                for sl in sample.splitlines():
                    lines.append(f"    {sl}")
        prev_hint = "\n".join(lines)

    # ── 動態注入已安裝的第三方套件清單 ──
    pkg_file = Path(__file__).parent.parent / "skill_packages.txt"
    if pkg_file.exists():
        pkg_lines = [l.strip() for l in pkg_file.read_text(encoding="utf-8").splitlines()
                     if l.strip() and not l.strip().startswith("#")]
        system_prompt = system_prompt.replace("{installed_packages}", ", ".join(pkg_lines))
    else:
        system_prompt = system_prompt.replace("{installed_packages}", "pandas, openpyxl, matplotlib, requests, beautifulsoup4, Pillow, python-docx")

    # ── matplotlib lazy injection:只有任務含繪圖關鍵字才注入(原本每次都送 ~15 行) ──
    # 90% 的 skill 任務不繪圖、塞這段純粹浪費 token
    _plot_keywords = ("圖", "繪", "chart", "plot", "視覺", "png", "jpg", "jpeg",
                      "折線", "長條", "直方", "柱狀", "圓餅", "散點", "pie", "scatter",
                      "boxplot", "heatmap", "dashboard", "視覺化", "matplotlib", "seaborn", "plotly")
    _td_lower = (task_description or "").lower()
    if any(kw.lower() in _td_lower for kw in _plot_keywords):
        system_prompt += """

<!--FIRST_ITER_BEGIN-->
【matplotlib 繪圖】
- 最前面加 `import matplotlib; matplotlib.use('Agg')` 避免 GUI 問題
- boxplot 的 `labels` 已棄用、改 `tick_labels`
- 中文字型:macOS 'PingFang HK';Windows 'Microsoft JhengHei' / 'SimHei'
  跨平台寫法:
  ```python
  import matplotlib
  for font in ['PingFang HK', 'Microsoft JhengHei', 'SimHei', 'Arial Unicode MS']:
      try:
          matplotlib.font_manager.findfont(font, fallback_to_default=False)
          matplotlib.rcParams['font.family'] = font
          break
      except: pass
  ```
- 完成 → `plt.savefig(路徑, dpi=150, bbox_inches='tight')` + `plt.close()`
<!--FIRST_ITER_END-->"""
        logger.info(f"[{step_name}] 📊 偵測到繪圖關鍵字、注入 matplotlib 提示")

    # ── 注入前次失敗歷史（重試時） ──
    failures_hint = ""
    if previous_failures:
        logger.info(f"[{step_name}] 🔄 重試：注入 {len(previous_failures)} 條失敗歷史到 user_prompt")
        lines = ["\n\n【⚠️ 前次嘗試失敗記錄 — 本次必須改用不同方法】"]
        for f in previous_failures:
            lines.append(f"\n第 {f['attempt']} 次嘗試失敗：")
            lines.append(f"  失敗原因：{f['reason']}")
            if f.get("suggestion"):
                lines.append(f"  驗證建議：{f['suggestion']}")
            if f.get("stdout_tail"):
                lines.append(f"  程式輸出（尾段）：{f['stdout_tail'][:400]}")
            if f.get("stderr_tail"):
                lines.append(f"  錯誤訊息：{f['stderr_tail'][:200]}")
        lines.append("\n→ 請分析上方失敗原因，改用已安裝套件中尚未嘗試過的不同方法或套件來完成任務。")
        failures_hint = "\n".join(lines)
    else:
        logger.debug(f"[{step_name}] 初次執行（無失敗歷史）")

    # ── 自動探查 batch 引用到的輸入檔結構、餵真實 schema 給模型(防腦補欄位/標籤)──
    input_peek_hint = _peek_input_files(task_description, logger)

    user_prompt = f"""請完成以下任務：

{task_description}{output_hint}{wd_hint}{prev_hint}{input_peek_hint}{failures_hint}

請直接使用 <tool>run_python</tool> 執行完整程式碼，不要用 markdown 展示。"""

    all_stdout: list[str] = []

    try:
        llm = _get_skill_llm(role=llm_role)
        # Prompt caching (#153):對 SystemMessage 加 ephemeral 1h cache_control。
        # SKILL loop 多輪、第 2 輪起 system_prompt 命中 cache、input cost 0.1x。
        # 非 Anthropic provider 會略過 cache_control(不影響功能)。
        _sys_cache_kwargs = {"cache_control": {"type": "ephemeral", "ttl": "1h"}}
        messages = [
            SystemMessage(content=system_prompt, additional_kwargs=_sys_cache_kwargs),
            HumanMessage(content=user_prompt),
        ]

        short_code_streak = 0  # 連續短程式碼計數器（偵測迴圈）
        last_error_sig = ""    # 上次錯誤簽名（偵測重複錯誤）
        same_error_count = 0   # 連續相同錯誤計數
        # 同一份 tool_input 重複偵測（不只比 stderr，連 exit_code-only 的失敗也能抓）
        last_tool_inputs: list[str] = []   # 最近幾次 (tool_name, tool_input_hash)
        # 連續失敗早停（任何類型的 tool failure 都算）
        consecutive_failures = 0
        # 錯誤類型追蹤:連續同類錯誤 2 次就注入對症提示(取代「換策略」空話)
        last_error_kind: Optional[str] = None
        same_error_kind_count = 0
        injected_hint_kinds: set[str] = set()  # 已注入過的 kind、避免每輪重複注入同一條
        # Output-driven done 提示:輸出檔已產生、提醒 LLM 早 done、避免「再驗證一輪」迴圈
        output_done_hint_injected = False
        import time as _time
        skill_start_time = _time.time()

        # Trace / token tracking — 累計 skill loop 全程 LLM 用量 + tool 呼叫時間軸；
        # 在每個 ExecResult return 前透過 _attach_trace() 接上、給前端 trace 視圖渲染。
        acc_usage: dict = {"input_tokens": 0, "output_tokens": 0, "total_tokens": 0}
        acc_tool_calls: list = []

        def _attach_trace(r: ExecResult) -> ExecResult:
            r.token_usage = dict(acc_usage)
            r.tool_calls = list(acc_tool_calls)
            return r

        # ── 檔案變化追蹤（C：iteration 之間的檔案 diff）────────────────────
        # 每次 run_python 後比對 working_dir 的 mtime 變化，log 哪些檔被新增/改了
        # 解決：LLM 多輪 rewrite 留下混雜的 PNG/XLSX，下游 validator 看到「不知是哪輪寫的」混亂狀態
        # 不在這做主動刪除（過於侵入），只 log；step retry 時 runner 已會清 output_path
        def _snapshot_dir_mtimes(d: str) -> dict:
            try:
                p = Path(d)
                if not p.is_dir():
                    return {}
                snap: dict[str, float] = {}
                # 只看頂層 + 一層子資料夾，避免 .venv 之類超大樹
                for entry in p.iterdir():
                    if entry.is_file():
                        try:
                            snap[entry.name] = entry.stat().st_mtime
                        except OSError:
                            pass
                return snap
            except Exception:
                return {}

        _wd_for_diff = working_dir or (str(Path(output_path).parent) if output_path else "")
        prev_mtimes = _snapshot_dir_mtimes(_wd_for_diff)
        if prev_mtimes:
            logger.info(f"[{step_name}] 📂 工作目錄初始 {len(prev_mtimes)} 個檔（{_wd_for_diff}）")

        last_successful_code: Optional[str] = None  # 供 Recipe Book 儲存：只記最後一段成功的 run_python
        # 防 LLM 在 run_python rc=1 後硬送 done(success=true)：done 來時要先看這個。
        # 沒跑過 run_python = None（容許純 read_file / web_search 後直接 done）；
        # 跑過 = True（成功）/ False（失敗，下次 done 拒絕並要 LLM 修錯）
        last_run_python_ok: Optional[bool] = None
        # 同理:防 LLM 在 run_shell rc=1(Cannot find module / no such file 等)後硬送 done。
        # 實測 Sonnet 會在 run_shell 失敗後直接 done(success=true)+ 寫假 summary 騙過 validator。
        last_run_shell_ok: Optional[bool] = None
        ask_user_count = 0               # ask_user 呼叫次數（上限 ASK_USER_MAX）
        web_search_count = 0             # web_search 呼叫次數（上限 WEB_SEARCH_MAX_PER_STEP）
        was_interactive = False          # 首次互動標記（給 recipe）
        # 沙盒 fallback 跨 iteration 狀態：使用者同意過就一路放行不再問
        # dict 是 mutable，傳進 helper 裡改 'allowed' 外層看得到
        sandbox_fallback_state: dict = {"allowed": False}

        # 「完成」字樣連續輪數計數 — 連 2 輪 LLM 口頭說完成但沒下 done tag 就強制收尾
        done_keyword_streak = 0
        prose_before_tool_violations = 0   # task #160:第一個 <tool> 前 prose 過長累計

        # ── Phase A.2 — Native function calling 切換點 ──────────────────
        # 共用 SUBAGENT_LOOP_MODE flag(SUBAGENT / SKILL 同個 .env 設定即可)。
        # native(**預設**,#157): 跑下方獨立 inner function、用 LangChain bind_tools()、消滅 <tool> 文字協議
        # text(opt-out): 繼續走原本 SKILL loop(向後相容、.env 設 SUBAGENT_LOOP_MODE=text 可退回)
        _skill_loop_mode = (os.environ.get("SUBAGENT_LOOP_MODE", "native") or "native").strip().lower()
        if _skill_loop_mode == "native":
            logger.info(f"[{step_name}] 🚀 SKILL loop 用 [NATIVE] function calling 模式")
            return _attach_trace(await _execute_skill_native_loop(
                llm=llm, messages=messages,
                output_path=output_path, working_dir=working_dir,
                run_id=run_id, step_name=step_name, logger=logger,
                tool_timeout=tool_timeout, ask_mode=ask_mode,
                acc_usage=acc_usage, acc_tool_calls=acc_tool_calls,
                all_stdout=all_stdout,
                pipeline_id=pipeline_id, task_description=task_description,
                input_paths=input_paths, no_save_recipe=no_save_recipe,
                silent_recipe=silent_recipe, recipe_step_key=_rkey,
            ))

        for iteration in range(SKILL_MAX_ITERATIONS):
            logger.info(f"[{step_name}] Skill 執行迭代 {iteration + 1}/{SKILL_MAX_ITERATIONS}")

            # Context 雪崩防護:第 4 輪起壓縮更早的 tool 結果(保留最近 3 輪完整、舊的摺成首尾預覽)
            # 解決 25 輪場景每輪 token 從 9k 漲到 21k 的問題、保持每輪 context size 大致穩定
            if iteration >= SKILL_CONTEXT_KEEP_RECENT_FULL + 1:
                _compacted = _compact_old_tool_results(messages)
                if _compacted:
                    logger.debug(f"[{step_name}] 🪶 壓縮 {_compacted} 條舊 tool 結果(防 context 膨脹)")
                # 同時壓縮 AIMessage 內的 <input>...</input>(LLM 自己寫的巨碼、pptx wrap JS 那種)
                # 跟 tool 結果壓縮對稱、舊輪保留 head+tail preview 即可、最近 3 輪保全
                _ai_compacted = _compact_old_llm_input_blocks(messages)
                if _ai_compacted:
                    logger.debug(f"[{step_name}] 🪶 壓縮 {_ai_compacted} 條舊 LLM 程式碼段(防 context 膨脹)")

            # 第一輪後 strip <!--FIRST_ITER_BEGIN/END--> 區塊(sandbox env / web_search / matplotlib /
            # tool_timeout 等規則第一輪 LLM 已內化、後續輪數重送純粹浪費 token)
            # 平均省 70-130 行 system prompt / 後續輪、約 1-2K tok / 輪
            if iteration == 1 and "<!--FIRST_ITER_BEGIN-->" in messages[0].content:
                import re as _re_strip
                _stripped = _re_strip.sub(
                    r"<!--FIRST_ITER_BEGIN-->.*?<!--FIRST_ITER_END-->\s*",
                    "",
                    messages[0].content,
                    flags=_re_strip.DOTALL,
                )
                messages[0] = SystemMessage(content=_stripped)
                logger.debug(f"[{step_name}] 🪶 第二輪起 system prompt 已 strip first-iter blocks")

            # 接近 max_iter 時、在送 LLM 前主動提醒收尾(比等沒下 tool 才提醒更早一步)
            # iter 剩 ≤ 2 時注入一條短 system reminder、不重複加(用 sentinel 防重)
            remaining_iters = SKILL_MAX_ITERATIONS - iteration
            if remaining_iters <= 2:
                _wrap_msg = (
                    f"⚠ 系統提醒:目前第 {iteration + 1}/{SKILL_MAX_ITERATIONS} 輪、只剩 {remaining_iters} 輪。"
                    "如果任務已可收尾、本輪請直接 <tool>done</tool><input>{\"success\": true, \"summary\": \"...\"}</input>。"
                    "不要再做純驗證或解釋、直接動作。"
                )
                # 上一次注入過就不重複
                _already = (messages and messages[-1].__class__.__name__ == 'HumanMessage'
                            and "系統提醒:目前第" in str(messages[-1].content))
                if not _already:
                    messages.append(HumanMessage(content=_wrap_msg))

            # 冷卻機制：每 SKILL_COOLDOWN_EVERY 次呼叫後暫停
            if iteration > 0 and iteration % SKILL_COOLDOWN_EVERY == 0:
                logger.info(f"[{step_name}] ⏸ 達到 {SKILL_COOLDOWN_EVERY} 次呼叫，冷卻 {SKILL_COOLDOWN_SECONDS} 秒...")
                await asyncio.sleep(SKILL_COOLDOWN_SECONDS)

            # 每次 LLM 呼叫間隔（避免撞 RPM 上限）
            if iteration > 0:
                await asyncio.sleep(SKILL_REQUEST_INTERVAL)

            from llm_factory import invoke_with_streaming
            llm_result = await invoke_with_streaming(
                llm, messages, label=step_name, timeout=600.0, logger=logger,
                return_usage=True,
            )
            reply = (llm_result.get("content") or "").strip()
            _um = llm_result.get("usage_metadata") or {}
            if _um:
                for _k in ("input_tokens", "output_tokens", "total_tokens"):
                    _v = _um.get(_k) or 0
                    if _v: acc_usage[_k] = acc_usage.get(_k, 0) + int(_v)
            if not acc_usage.get("model"):
                acc_usage["model"] = llm_result.get("model") or ""
            # 完整記錄 LLM 回覆（含程式碼），避免 log 截斷讓後續分析誤判
            _reply_preview = reply if len(reply) <= 4000 else reply[:4000] + f"...[已截斷，完整長度 {len(reply)} 字]"
            logger.debug(f"[{step_name}] Agent 回覆：\n{_reply_preview}")

            # 偵錯：如果 reply 包含 done，印出 done 附近的文字
            if 'done' in reply.lower():
                idx = reply.lower().index('done')
                snippet = reply[max(0, idx-80):idx+80]
                logger.info(f"[{step_name}] reply 含 'done'，上下文：…{snippet}…")

            tool_calls = _parse_skill_tool_calls(reply)

            if not tool_calls:
                # 沒有工具呼叫、提示 agent
                # 同時偵測「LLM 口頭說完成但沒打 done tag」連續 2 輪 → 強制 done 收尾
                # 解決 Sonnet / GPT 等強模型常見死循環:寫好檔後一直 verify、不打標準 tag
                if _looks_like_done(reply):
                    done_keyword_streak += 1
                    if done_keyword_streak >= 2:
                        logger.warning(
                            f"[{step_name}] 連續 {done_keyword_streak} 輪 LLM 口頭表示完成但沒下 <tool>done</tool> tag,"
                            f" 強制 done 收尾(避免無謂 iter 燒 token)"
                        )
                        # 條件:最近一次 run_python 成功 OR 沒跑過 run_python(純 read_file 也算)
                        # 否則拒絕強制 done、繼續循環給 LLM 修錯機會
                        if last_run_python_ok is not False:
                            # 模擬 done call、跳出 loop
                            final_summary = reply[:500] if reply else "LLM 連續表示完成、系統強制收尾"
                            logger.info(f"[{step_name}] Skill 執行完成:強制 done — {final_summary[:100]}")
                            all_stdout.append(f"[系統強制 done] {final_summary}")
                            return ExecResult(
                                exit_code=0,
                                stdout="\n".join(all_stdout) if all_stdout else final_summary,
                                stderr="",
                                missing_packages=[],
                            )
                        else:
                            messages.append(HumanMessage(content=reply))
                            messages.append(HumanMessage(content=(
                                "⚠ 你說完成但最近一次 run_python 失敗、必須先修錯。"
                                "請寫程式修正,或如果無法完成就用 done(success=false, reason=...)。"
                            )))
                            continue
                else:
                    done_keyword_streak = 0

                messages.append(HumanMessage(content=reply))
                # 接近 max_iter 時加緊提醒 LLM 收尾
                remaining = SKILL_MAX_ITERATIONS - iteration - 1
                if remaining <= 2:
                    messages.append(HumanMessage(content=(
                        f"⚠ 剩 {remaining} 輪、必須收尾。若任務已完成請立刻打 "
                        f"<tool>done</tool><input>{{...}}</input>;若還沒完成請寫程式繼續。"
                        f"不要再做純文字確認 / 解釋,直接動作。"
                    )))
                else:
                    messages.append(HumanMessage(content="請使用工具來執行任務,或呼叫 done 回報結果。"))
                continue
            else:
                # 有 tool call、清掉 done streak
                done_keyword_streak = 0

            # ── Prose-before-tool 守門(task #160、backport 自 subagent text 版)──
            # 第一個 <tool> 前 prose 過長 = LLM「想很多/寫長篇才動手」(Sonnet 4.6 寫 5 萬字
            # parser 案例)、燒 output token + 拖斷長連線。soft 只 log、enforce 連 _LIMIT 輪中止。
            if _SKILL_PROSE_CAP_MODE != "off":
                _tool_pos = reply.find("<tool>")
                _prose_chars = len(reply[:_tool_pos].strip()) if _tool_pos > 0 else 0
                if _prose_chars > _SKILL_PROSE_BEFORE_TOOL_THRESHOLD:
                    prose_before_tool_violations += 1
                    logger.warning(
                        f"[{step_name}] ⚠ 第 {iteration + 1} 輪 tool 前 prose {_prose_chars:,} 字"
                        f"(上限 {_SKILL_PROSE_BEFORE_TOOL_THRESHOLD})、累計違規 "
                        f"{prose_before_tool_violations}/{_SKILL_PROSE_BEFORE_TOOL_LIMIT} "
                        f"[mode={_SKILL_PROSE_CAP_MODE}]"
                    )
                    if (_SKILL_PROSE_CAP_MODE == "enforce"
                            and prose_before_tool_violations >= _SKILL_PROSE_BEFORE_TOOL_LIMIT):
                        err_msg = (
                            f"連續 {_SKILL_PROSE_BEFORE_TOOL_LIMIT} 輪 tool 前 prose 超過 "
                            f"{_SKILL_PROSE_BEFORE_TOOL_THRESHOLD} 字、強制中止避免燒 token / 拖斷連線。"
                            f"解析請直接寫進 run_python 程式碼、不要在 reply 寫長篇 parser / 解釋。"
                        )
                        logger.error(f"[{step_name}] ✗ {err_msg}")
                        return ExecResult(
                            exit_code=1,
                            stdout="\n".join(all_stdout) if all_stdout else reply[:500],
                            stderr=err_msg,
                            agent_concluded_fail=True,
                        )
                else:
                    prose_before_tool_violations = 0

            # 多工具偵測：LLM 一次塞 run_python + done 是惡習（會把假成功訊息混進 done），
            # 預設只跑第一個 tool（既有行為）、明確告訴 LLM「這次只跑 X、忽略 Y」
            # 用獨立 regex 算 <tool>name</tool> tag 數，比 parser 的多階段 fallback 結果更直接
            _tag_count = len(re.findall(r"<tool>\s*\w+\s*</tool>", reply))
            multi_tool_warn = _tag_count > 1

            call = tool_calls[0]
            tool_name = call["tool"]
            tool_input = call["input"]
            logger.info(f"[{step_name}] 解析結果：tool={tool_name}, input_len={len(tool_input)}"
                        + (f"（⚠ 偵測到 {_tag_count} 個 <tool> 標籤、只跑第一個）" if multi_tool_warn else ""))
            # Trace：push 該輪 tool 呼叫（name + input、result 由後續 dispatch 拿到後 update）
            acc_tool_calls.append({
                "name": tool_name,
                "input_preview": (tool_input or "")[:200],
                "result_preview": "",
            })

            # done → 結束（但先驗證 output 檔案是否存在 + 最近 run_python 必須沒失敗）
            if tool_name == "done":
                try:
                    data = json.loads(tool_input)
                    success = data.get("success", False)
                    summary = data.get("summary", data.get("error", ""))

                    # 守門 1：宣稱成功但「最近一次 run_python 失敗」→ 拒絕 done
                    # 防 LLM 在程式碼炸掉後硬送 done(success=true)、靠殘留檔騙過 output_path 檢查
                    if success and last_run_python_ok is False:
                        logger.warning(f"[{step_name}] Agent 在 run_python 失敗後送 done(success=true) — 拒絕並要求修錯")
                        messages.append(HumanMessage(content=reply))
                        messages.append(HumanMessage(
                            content="[系統] 拒絕 done：你最近一次 run_python 執行失敗（看上面的 stderr / Traceback）。"
                                    "請先用 run_python 修正錯誤、確認真的成功（沒有 [exit code: N] 失敗訊息）後，"
                                    "才能呼叫 done。不要在程式碼失敗後硬送 success=true。"
                        ))
                        continue

                    # 守門 1b:同上,但守的是 run_shell(Cannot find module / no such file / rc!=0 等)
                    # 實測 Sonnet 在 run_shell 失敗後會直接 done(success=true)+ 編造假 summary 騙 validator
                    if success and last_run_shell_ok is False:
                        logger.warning(f"[{step_name}] Agent 在 run_shell 失敗後送 done(success=true) — 拒絕並要求修錯")
                        messages.append(HumanMessage(content=reply))
                        messages.append(HumanMessage(
                            content="[系統] 拒絕 done：你最近一次 run_shell 執行失敗（看上面的 stderr 與 exit code）。"
                                    "請先修正命令（常見:路徑不對、模組沒裝、引號 escape 錯）、確認 run_shell 真的成功"
                                    "(rc=0、沒有 Cannot find / No such file / Error 等)後,才能呼叫 done。"
                                    "**不要在 shell 失敗後硬送 success=true、不要在 summary 編造未發生的執行結果。**"
                        ))
                        continue

                    # 守門 2：宣稱成功但 output 檔案不存在
                    logger.debug(f"[{step_name}] done 檢查：success={success}, output_path={output_path}, exists={Path(output_path).exists() if output_path else 'N/A'}, last_run_python_ok={last_run_python_ok}")
                    if success and output_path and not Path(output_path).exists():
                        logger.warning(f"[{step_name}] Agent 宣稱成功但輸出檔案 {output_path} 不存在，要求重新執行")
                        messages.append(HumanMessage(content=reply))
                        messages.append(HumanMessage(
                            content=f"[系統] 你宣稱成功但輸出檔案 {output_path} 不存在。"
                                    f"你必須使用 run_python 工具實際執行程式碼來產生檔案，"
                                    f"不能只展示程式碼。請使用 <tool>run_python</tool> 執行。"
                        ))
                        continue

                    all_stdout.append(f"[Skill 完成] {summary}")
                    logger.info(f"[{step_name}] Skill 執行完成：{'成功' if success else '失敗'} — {summary}")
                    # 成功 → 儲存 recipe 供下次快速重跑
                    _pending_recipe = None
                    # Recipe 表 workflow_id 有 FK 約束、跑 ad-hoc YAML(無關聯 workflow)時 pipeline_id
                    # 是 config.name 字串、不在 workflows 表 → save 會炸 FOREIGN KEY constraint。
                    # 先驗 pipeline_id 是真實 workflow id 才繼續、否則直接跳過 recipe(這種 run 本來
                    # 也不會被前端用 cache 重播、不存 recipe 不影響)。
                    _is_real_wf = False
                    if success and pipeline_id:
                        try:
                            from db import get_workflow as _gw_check
                            _is_real_wf = _gw_check(pipeline_id) is not None
                        except Exception:
                            _is_real_wf = False
                    if success and pipeline_id and last_successful_code and _is_real_wf:
                        try:
                            import sys as _sys2
                            from pipeline.recipe import _sha1 as _recipe_sha1, _fingerprint_input as _recipe_fp
                            runtime = _time.time() - skill_start_time
                            _fp = {}
                            for p in (input_paths or []):
                                _fp[str(p)] = _recipe_fp(p)
                            recipe_data = {
                                "pipeline_id": pipeline_id,
                                "step_name": _rkey,
                                "task_hash": _recipe_sha1(task_description),
                                "input_fingerprints": _fp,
                                "output_path": output_path,
                                "code": last_successful_code,
                                "python_version": f"{_sys2.version_info.major}.{_sys2.version_info.minor}",
                                "runtime_sec": runtime,
                                "was_interactive": was_interactive,
                            }
                            from db import get_recipe as _get_recipe, save_recipe as _db_save_recipe
                            existing = _get_recipe(pipeline_id, _rkey)
                            if silent_recipe:
                                # 無人值守模式（TG / 排程觸發）:
                                # - 有 recipe → **跳過、不覆寫**（保護用戶手動微調過的版本）
                                # - 無 recipe → 直接建立（首次跑可以 seed）
                                if existing:
                                    logger.info(f"[{step_name}] silent_recipe:Recipe 已存在、跳過(不覆寫)")
                                else:
                                    _db_save_recipe(
                                        pipeline_id, _rkey, recipe_data["task_hash"],
                                        _fp, output_path, last_successful_code,
                                        recipe_data["python_version"], runtime,
                                        was_interactive=was_interactive,
                                    )
                                    logger.info(f"[{step_name}] silent_recipe:首次建立 Recipe")
                            elif no_save_recipe:
                                # 互動延遲模式（桌面手動跑）：
                                if existing:
                                    # 已有 recipe → 延遲儲存等用戶確認(避免覆蓋)
                                    _pending_recipe = recipe_data
                                    logger.info(f"[{step_name}] Recipe 已存在，延遲儲存等待確認")
                                else:
                                    # 無 recipe → 直接儲存(建立新的不算覆蓋)
                                    _db_save_recipe(
                                        pipeline_id, _rkey, recipe_data["task_hash"],
                                        _fp, output_path, last_successful_code,
                                        recipe_data["python_version"], runtime,
                                        was_interactive=was_interactive,
                                    )
                                    logger.info(f"[{step_name}] 首次建立 Recipe")
                            else:
                                # 正常模式（桌面 + 工作流無 skill / 工作流明確要 inline 寫）
                                _db_save_recipe(
                                    pipeline_id, _rkey, recipe_data["task_hash"],
                                    _fp, output_path, last_successful_code,
                                    recipe_data["python_version"], runtime,
                                    was_interactive=was_interactive,
                                )
                        except Exception as e:
                            logger.warning(f"[{step_name}] Recipe 儲存失敗：{e}")
                    pkgs = data.get("missing_packages", []) if not success else []
                    if pkgs:
                        logger.info(f"[{step_name}] LLM 回報缺少套件：{pkgs}")
                    final_stdout = (_build_clean_success_stdout(all_stdout, "[Skill 完成]")
                                    if success else "\n".join(all_stdout))
                    return _attach_trace(ExecResult(
                        exit_code=0 if success else 1,
                        stdout=final_stdout,
                        stderr="" if success else summary,
                        pending_recipe=_pending_recipe,
                        missing_packages=pkgs or None,
                        # 主動 done(success=false) 且非缺套件場景 → 標記為「明確結論」
                        # 讓 runner 跳過 retry。缺套件(pkgs)走既有 missing_dependency 路徑、不重複攔。
                        agent_concluded_fail=(not success and not pkgs),
                    ))
                except json.JSONDecodeError:
                    messages.append(HumanMessage(content=reply))
                    messages.append(HumanMessage(content="[系統] done 的 input 必須是有效 JSON，請重試。"))
                    continue

            # ask_user → 暫停 pipeline，等待使用者回答
            # ask_mode ON 時取消上限（使用者已明確表態想被問、不再防濫用）；OFF 時沿用 ASK_USER_MAX 保護
            if tool_name == "ask_user":
                ask_user_count += 1
                if not ask_mode and ask_user_count > ASK_USER_MAX:
                    tool_result = f"[錯誤] ask_user 已達上限 {ASK_USER_MAX} 次（詢問模式未開啟）。請以預設值完成或呼叫 done(success=false)。"
                    messages.append(HumanMessage(content=reply))
                    messages.append(HumanMessage(content=f"[工具結果 — ask_user]\n{tool_result}"))
                    continue
                try:
                    q_data = json.loads(tool_input)
                    question = (q_data.get("question") or "").strip()
                    options = q_data.get("options") or []
                    context = (q_data.get("context") or "").strip()
                    if not question:
                        raise ValueError("question 不可為空")
                    if not isinstance(options, list):
                        options = []
                except Exception as e:
                    messages.append(HumanMessage(content=reply))
                    messages.append(HumanMessage(
                        content=f"[系統] ask_user input 格式錯誤：{e}。正確格式：{{\"question\":\"...\", \"options\":[...], \"context\":\"...\"}}"
                    ))
                    continue

                answer = await _wait_for_ask_user(run_id, question, options, context, logger, step_name)
                if answer is None:
                    tool_result = "[錯誤] 等待使用者回答逾時或被取消，請以合理預設完成或呼叫 done(success=false)。"
                else:
                    was_interactive = True  # 標記 recipe「首次有人工回答」
                    tool_result = f"使用者回答：{answer}"
                messages.append(HumanMessage(content=reply))
                messages.append(HumanMessage(content=f"[工具結果 — ask_user]\n{tool_result}"))
                continue

            # web_search → 直接呼叫（不走 _execute_skill_tool 的沙盒 pre-flight；它是純 HTTPS API）
            # 這裡單獨處理為了在 call 之前檢查「單一 skill step 上限」
            if tool_name == "web_search":
                web_search_count += 1
                if web_search_count > WEB_SEARCH_MAX_PER_STEP:
                    tool_result = (
                        f"[web_search 錯誤] 本步驟已達搜尋次數上限（{WEB_SEARCH_MAX_PER_STEP} 次）。"
                        "請整合前面搜尋結果回答，或呼叫 done(success=false)。"
                    )
                else:
                    tool_result = await asyncio.get_event_loop().run_in_executor(
                        None,
                        lambda ti=tool_input, cc=web_search_count, lg=logger:
                            _skill_web_search(ti, call_count=cc, logger=lg),
                    )
                all_stdout.append(f"[web_search] {tool_result}")
                messages.append(HumanMessage(content=reply))
                messages.append(HumanMessage(content=f"[工具結果 — web_search]\n{tool_result}"))
                continue

            # view_image → 走多模態：把圖檔讀成 base64 後以 image_url 形式塞進 HumanMessage，
            # 讓視覺模型真的「看到」圖。模型不支援視覺時 LLM 自己會回說看不懂，由 agent 決定下一步。
            if tool_name == "view_image":
                img_data = await asyncio.get_event_loop().run_in_executor(
                    None, _skill_view_image, tool_input
                )
                logger.info(f"[{step_name}] view_image：{img_data['text']}")
                all_stdout.append(f"[view_image] {img_data['text']}")
                messages.append(HumanMessage(content=reply))
                if img_data["image_b64"]:
                    messages.append(HumanMessage(content=[
                        {"type": "text", "text": f"[工具結果 — view_image]\n{img_data['text']}\n請仔細觀察圖片內容後再決定下一步。"},
                        {"type": "image_url", "image_url": {
                            "url": f"data:{img_data['image_mime']};base64,{img_data['image_b64']}"
                        }},
                    ]))
                else:
                    messages.append(HumanMessage(content=f"[工具結果 — view_image]\n{img_data['text']}"))
                continue

            # 執行工具
            logger.info(f"[{step_name}] 工具呼叫：{tool_name}")
            # 若是 run_python / run_shell，先 pre-flight 沙盒狀態：
            #   ask_mode ON：沙盒不可用就問使用者要重試 / 退 host / 中止
            #   ask_mode OFF：維持原本靜默 fallback 行為
            force_host = False
            if tool_name in ("run_python", "run_shell"):
                decision = await _preflight_sandbox(
                    ask_mode=ask_mode,
                    fallback_state=sandbox_fallback_state,
                    run_id=run_id,
                    step_name=step_name,
                    logger=logger,
                )
                if decision == "abort":
                    logger.info(f"[{step_name}] 使用者選擇中止（沙盒不可用）")
                    return ExecResult(
                        exit_code=1,
                        stdout="\n".join(all_stdout),
                        stderr="使用者透過 ask_user 選擇中止（沙盒不可用）",
                        pending_recipe=_pending_recipe,
                        missing_packages=None,
                    )
                force_host = (decision == "host")

            # ── Phase B: 敏感命令攔截 ──
            # 只對 run_python / run_shell 檢查（其他 tool 沒命令執行風險）
            # A1 模式:每次命中都問,不記憶
            # 安裝類命令(pip install / npm install / apt 等)無論 ask_mode 開關都攔
            # — 容器持久化、裝下去就留著、且可能裝奇怪套件、安全不可妥協
            if tool_name in ("run_python", "run_shell"):
                _classification = classify_command(tool_input)
                _always_intercept = _classification and _classification[0] == "install"
                if _classification and (ask_mode or _always_intercept):
                    _cat, _label, _preview = _classification
                    logger.warning(
                        f"[{step_name}] 🛡 ask_mode 偵測到敏感命令（{_cat}）→ 等用戶授權"
                    )
                    _decision = await _wait_for_command_approval(
                        run_id=run_id,
                        category=_cat,
                        label=_label,
                        preview=_preview,
                        tool_name=tool_name,
                        logger=logger,
                        step_name=step_name,
                    )
                    if _decision == "deny":
                        logger.warning(f"[{step_name}] 用戶拒絕命令 → 中止 step")
                        return ExecResult(
                            exit_code=2,
                            stdout="\n".join(all_stdout),
                            stderr=f"使用者拒絕執行敏感命令（{_cat}:{_label}）",
                            pending_recipe=None,
                            missing_packages=None,
                        )
                    if _decision == "hint":
                        # 用戶要改任務 — 走 retry_with_hint 機制（runner 層處理）
                        # 這裡退出 agent loop、stderr 標記讓 runner 跳到 retry_with_hint awaiting
                        logger.info(f"[{step_name}] 用戶選擇改任務、退出 agent loop")
                        return ExecResult(
                            exit_code=2,
                            stdout="\n".join(all_stdout),
                            stderr="使用者選擇改任務（command_approval hint）",
                            pending_recipe=None,
                            missing_packages=None,
                        )
                    if _decision == "timeout":
                        logger.warning(f"[{step_name}] command_approval 逾時、中止 step")
                        return ExecResult(
                            exit_code=2,
                            stdout="\n".join(all_stdout),
                            stderr=f"敏感命令授權等待逾時（{COMMAND_APPROVAL_TIMEOUT}s）",
                            pending_recipe=None,
                            missing_packages=None,
                        )
                    # _decision == "allow" → 繼續執行
                    logger.info(f"[{step_name}] ✓ 用戶授權執行命令（{_cat}）")

            # ── run_shell 的 pip install → 直接轉 missing_dependency 使用者確認(不執行) ──
            # (2026-05-31:sandbox 跑外部腳本(main_CLI.py)撞缺套件時 stderr 常沒完整傳回、
            #  ModuleNotFoundError regex 抓不到;但 skill 想自己裝時命令明擺著有套件名 → 從這抽最可靠。
            #  抽到套件名就直接轉系統『允許安裝』確認,不靠 stderr 格式、也不靠 skill 乖乖 done(missing_packages)。)
            if tool_name == "run_shell":
                # tool_input 可能是純字串、或 native FC 的 JSON {"command":"..."} → 正規化成命令字串
                # (2026-05-31:skill loop 這層的 tool_input 是 JSON、_is_pip_install_cmd 直接判會 False;
                #  _execute_skill_tool 內部才解析成純字串、所以要在這先正規化。)
                _cmd_norm = tool_input
                if isinstance(_cmd_norm, str) and _cmd_norm.lstrip().startswith("{"):
                    try:
                        import json as _j_norm
                        _parsed = _j_norm.loads(tool_input)
                        _cmd_norm = _parsed.get("command") or _parsed.get("input") or tool_input
                    except Exception:
                        _cmd_norm = tool_input
                if _is_pip_install_cmd(_cmd_norm):
                    _after = _cmd_norm.split("install", 1)[1] if "install" in _cmd_norm else ""
                    _pkgs = [p.split("==")[0].split(">")[0].split("<")[0].strip()
                             for p in _after.split()
                             if not p.startswith("-") and not p.endswith((".txt", ".cfg", ".toml"))]
                    if _pkgs:
                        logger.warning(f"[{step_name}] 🛑 攔 run_shell pip install {_pkgs} → 轉 missing_dependency 使用者確認")
                        return ExecResult(
                            exit_code=1,
                            stdout="\n".join(all_stdout),
                            stderr=f"偵測到 skill 嘗試自行安裝套件 {_pkgs}。已改走系統的『允許安裝』使用者確認流程。",
                            pending_recipe=None,
                            missing_packages=_pkgs,
                        )

            tool_result = await asyncio.get_event_loop().run_in_executor(
                None, lambda tn=tool_name, ti=tool_input, lg=logger, fh=force_host, tt=tool_timeout: _execute_skill_tool(tn, ti, cwd=working_dir, run_id=run_id, logger=lg, force_host=fh, tool_timeout=tt)
            )
            # 完整記錄工具結果（錯誤訊息如 ModuleNotFoundError 常超過 300 字）
            _tr_preview = tool_result if len(tool_result) <= 3000 else tool_result[:3000] + f"...[已截斷，完整長度 {len(tool_result)} 字]"
            logger.debug(f"[{step_name}] 工具結果：\n{_tr_preview}")
            # Trace：把 tool result 補回上一輪 push 的 entry
            if acc_tool_calls and acc_tool_calls[-1].get("name") == tool_name and not acc_tool_calls[-1].get("result_preview"):
                acc_tool_calls[-1]["result_preview"] = (str(tool_result) or "")[:300]
            all_stdout.append(f"[{tool_name}] {tool_result}")
            # 追蹤 run_shell 成敗:用 [exit code:] marker(_skill_run_shell 失敗會帶)
            if tool_name == "run_shell":
                if "[exit code:" not in tool_result:
                    last_run_shell_ok = True
                else:
                    last_run_shell_ok = False
                    logger.info(f"[{step_name}] run_shell 失敗 → last_run_shell_ok=False(下次 done 會被守門)")
                    # ── ModuleNotFoundError 早期攔截(同 run_python):run_shell 跑外部腳本(如 main_CLI.py)撞缺套件 ──
                    # 自動抽套件名 → missing_packages → runner 走 missing_dependency 使用者確認。
                    # (2026-05-31:原本只 run_python 有;cli-extractor 用 run_shell 跑 main_CLI.py 撞缺套件時
                    #  沒攔到、靠模型自填 missing_packages 不可靠 → 走 failure。系統層級自動抽才 robust。)
                    import re as _mnf_re2
                    _m_shell = _mnf_re2.search(
                        r"ModuleNotFoundError: No module named ['\"]([\w.]+)['\"]", tool_result)
                    if _m_shell:
                        _pkg_shell = _m_shell.group(1).split(".")[0]
                        logger.warning(
                            f"[{step_name}] 🛑 run_shell 偵測 ModuleNotFoundError: '{_pkg_shell}' "
                            f"→ 中止 agent loop、轉 runner 走 missing_dependency 確認")
                        return ExecResult(
                            exit_code=1,
                            stdout="\n".join(all_stdout),
                            stderr=f"ModuleNotFoundError: 缺少套件 '{_pkg_shell}'。"
                                   f"系統將彈出安裝確認對話框（TG / 前端 modal）。",
                            pending_recipe=None,
                            missing_packages=[_pkg_shell],
                        )
            # 追蹤 run_python 成敗：用 [exit code:] 在 tool_result 是否出現當判斷
            # （sandbox 跟 host 兩條路徑都用同個 marker，見 _skill_run_python / _try_sandbox_exec）
            if tool_name == "run_python":
                if "[exit code:" not in tool_result:
                    last_run_python_ok = True
                    last_successful_code = tool_input  # 成功才記給 recipe

                    # ── Output-driven done 提示 ────────────────────────────────
                    # 偵測:run_python 成功 + output_path 存在 + 檔案有內容 → 提醒早 done
                    # 解決「再驗證一下」「再優化一下」迴圈、Gemma 等弱模型不會主動 done 的問題
                    # 注入 1 次/step、別重複煩 LLM
                    if (
                        output_path
                        and not output_done_hint_injected
                        and Path(output_path).exists()
                        and Path(output_path).stat().st_size > 0
                    ):
                        _size = Path(output_path).stat().st_size
                        if has_external_validator:
                            _hint = (
                                f"[系統提醒] ✅ 輸出檔 `{output_path}` 已存在(**{_size:,} bytes**)且 Python 跑成功。\n"
                                f"此 step 後面有外部 AI validator 會檢查內容、**不要再 read_file 確認**、"
                                f"**不要再優化格式**。如果輸出本身是合理的、現在就 `<tool>done</tool>` 結束。"
                            )
                        else:
                            _hint = (
                                f"[系統提醒] ✅ 輸出檔 `{output_path}` 已存在(**{_size:,} bytes**)且 Python 跑成功。\n"
                                f"此 step **沒有外部 validator** → 用 stdout 已看到的內容快速確認結果合理就 `<tool>done</tool>`,"
                                f"**不要為「驗證檔存在」再 read_file**(浪費一輪)、**不要為「再優化」多寫**(合理就 OK)。"
                            )
                        messages.append(HumanMessage(content=_hint))
                        output_done_hint_injected = True
                        logger.info(
                            f"[{step_name}] 💡 output-driven done 提示已注入"
                            f"(size={_size}, has_validator={has_external_validator})"
                        )
                else:
                    last_run_python_ok = False
                    logger.info(f"[{step_name}] run_python 失敗 → last_run_python_ok=False（下次 done 會被守門）")

                    # ── ModuleNotFoundError 早期攔截 ──
                    # LLM 自己修不了（套件根本沒裝），讓它 iter 5 次只是浪費 token。
                    # 直接 break agent loop、回 ExecResult.missing_packages → runner 接到後
                    # 走新的 awaiting_type=missing_dependency 路徑、彈安裝確認對話框給用戶
                    import re as _mnf_re
                    _mnf_match = _mnf_re.search(
                        r"ModuleNotFoundError: No module named ['\"]([\w.]+)['\"]",
                        tool_result,
                    )
                    if _mnf_match:
                        _missing_pkg = _mnf_match.group(1).split(".")[0]  # pandas.io.x → pandas
                        logger.warning(
                            f"[{step_name}] 🛑 偵測到 ModuleNotFoundError: '{_missing_pkg}' "
                            f"→ 提早中止 agent loop、轉交 runner 處理（不浪費 LLM iter）"
                        )
                        # 注意:不能引用 _pending_recipe — 那個變數只在後面 done(success=true)
                        # 分支才定義,在這個提早 return 路徑會 NameError 被外層 try/except
                        # 抓住、回退成 exit_code=-3 missing_packages=None,讓 runner 看不到攔截結果
                        return ExecResult(
                            exit_code=1,
                            stdout="\n".join(all_stdout),
                            stderr=f"ModuleNotFoundError: 缺少套件 '{_missing_pkg}'。"
                                   f"系統將彈出安裝確認對話框（TG / 前端 modal）。",
                            pending_recipe=None,
                            missing_packages=[_missing_pkg],
                        )
                # 檔案 diff：log 這次 run_python 改了哪些檔（C 優化）
                if _wd_for_diff:
                    cur_mtimes = _snapshot_dir_mtimes(_wd_for_diff)
                    new_files = [n for n in cur_mtimes if n not in prev_mtimes]
                    modified = [n for n in cur_mtimes if n in prev_mtimes and cur_mtimes[n] > prev_mtimes[n]]
                    if new_files or modified:
                        parts = []
                        if new_files: parts.append(f"新增 {len(new_files)}: {', '.join(new_files[:8])}")
                        if modified:  parts.append(f"修改 {len(modified)}: {', '.join(modified[:8])}")
                        logger.info(f"[{step_name}] 📝 檔案變化 — {' / '.join(parts)}")
                    prev_mtimes = cur_mtimes

            messages.append(HumanMessage(content=reply))
            messages.append(HumanMessage(content=f"[工具結果 — {tool_name}]\n{tool_result}"))
            # 多工具警告：reply 含 ≥2 個 <tool>，告訴 LLM 我們只跑第一個（{tool_name}），其他 tag 被忽略
            # 並提醒不要在 <input> 後面接「假的 stdout」當 prediction，那會自我欺騙
            if multi_tool_warn:
                messages.append(HumanMessage(
                    content=f"[系統警告] 你這個 reply 裡有 {_tag_count} 個 <tool> 標籤。系統一次只跑第一個（{tool_name}）"
                            f"— 其他標籤跟它們的 <input> 都被忽略，包括你可能寫在後面的 done。"
                            f"\n規則：每個 reply 只能寫一個 <tool>...</tool><input>...</input>，然後等系統回真實結果。"
                            f"\n禁止在 <input> 之後寫『Successfully executed.』『Recalc Result: ...』這種假裝是執行結果的文字 —"
                            f"那是你在自我欺騙，真實結果上面已經給你了。"
                ))

            # 迴圈偵測：連續多次只執行短程式碼，注入提示打破迴圈
            if tool_name == "run_python" and len(tool_input) < 200:
                short_code_streak += 1
                if short_code_streak >= 3:
                    logger.warning(f"[{step_name}] 偵測到連續 {short_code_streak} 次短程式碼，注入提示打破迴圈")
                    messages.append(HumanMessage(
                        content="[系統警告] 你已經連續多次只執行讀取資料的小段程式碼，但任務尚未完成。"
                                "請立即在一個 <tool>run_python</tool> 呼叫中寫出完整的程式碼來產生輸出檔案。"
                                "不要再分步驟讀取資料，直接把讀取、處理、寫入都放在同一段程式碼中執行。"
                    ))
                    short_code_streak = 0
            else:
                short_code_streak = 0

            # 錯誤重複偵測：連續出現相同錯誤時，注入修正提示
            if tool_name == "run_python" and "[stderr]" in tool_result:
                # 取錯誤的關鍵行作為簽名（最後一行 traceback）
                err_lines = [l for l in tool_result.split("\n") if l.strip() and not l.startswith("[")]
                error_sig = err_lines[-1].strip() if err_lines else ""
                if error_sig and error_sig == last_error_sig:
                    same_error_count += 1
                    if same_error_count >= 2:
                        logger.warning(f"[{step_name}] 相同錯誤連續出現 {same_error_count + 1} 次，注入修正提示")
                        messages.append(HumanMessage(
                            content=f"[系統警告] 你已經連續 {same_error_count + 1} 次遇到相同錯誤：{error_sig}\n"
                                    "你不能重複提交相同的程式碼。請換一個完全不同的方法。\n"
                                    "建議：先用 read_file 或 run_python 讀取輸入檔的前幾行，確認實際的欄位名稱和資料格式，"
                                    "然後根據實際欄位名稱重寫程式碼。"
                        ))
                        same_error_count = 0
                else:
                    last_error_sig = error_sig
                    same_error_count = 1
            else:
                last_error_sig = ""
                same_error_count = 0

            # ── 重複 tool_input 偵測 & 連續失敗早停（不依賴 stderr 有內容）────
            # 這是為了抓這類邊緣情況：LLM 送完全一樣的程式碼，subprocess 吐 exit_code=1
            # 但 stderr 空白，既有的 error_sig 比對因為沒有 [stderr] 而完全失效，
            # iteration 就這樣耗到 cap
            failed_now = ("[exit code:" in tool_result) or ("[錯誤]" in tool_result)
            if failed_now:
                consecutive_failures += 1
            else:
                consecutive_failures = 0

            # 取 tool_input 前 300 字當 signature（避免巨量程式碼比對耗 CPU）
            sig = f"{tool_name}:{tool_input[:300]}"
            last_tool_inputs.append(sig)
            last_tool_inputs = last_tool_inputs[-4:]   # 只保留最近 4 筆

            # 0) 錯誤類型分類:連續 2 次同 kind 錯誤 → 注入對症提示
            #    這比「換策略」空話有用、給 LLM 具體下一步(例如 SyntaxError on triple-quote
            #    → 改用 Path.write_text;AttributeError 'str' has no attribute → 不要 pickle)
            if failed_now:
                _kind = _classify_tool_error(tool_result)
                if _kind:
                    if _kind == last_error_kind:
                        same_error_kind_count += 1
                    else:
                        last_error_kind = _kind
                        same_error_kind_count = 1
                    # 第 2 次同類錯誤 + 還沒注入過這條 hint → 注入
                    if same_error_kind_count >= 2 and _kind not in injected_hint_kinds:
                        _hint = _hint_for_error_kind(_kind)
                        if _hint:
                            logger.warning(f"[{step_name}] 連續 2 次 '{_kind}' 錯誤、注入對症提示")
                            messages.append(HumanMessage(
                                content=f"[系統診斷] 你連續 2 次撞到同類錯誤(`{_kind}`)。針對性建議:\n\n{_hint}"
                            ))
                            injected_hint_kinds.add(_kind)
                else:
                    # 錯誤對不上任何 pattern → reset 計數
                    last_error_kind = None
                    same_error_kind_count = 0
            else:
                last_error_kind = None
                same_error_kind_count = 0

            # 1) 剛剛這次和上次 tool_input 一模一樣且失敗 → 強制打破迴圈
            if failed_now and len(last_tool_inputs) >= 2 and last_tool_inputs[-1] == last_tool_inputs[-2]:
                logger.warning(f"[{step_name}] 偵測到連續送相同 {tool_name}，注入打破迴圈提示")
                messages.append(HumanMessage(
                    content=(
                        "[系統警告] 你剛剛送了**完全一樣的 tool 呼叫**並再次失敗。"
                        "重試同一份程式碼永遠不會有不同結果。立刻改變策略：\n"
                        "1. 先用 read_file 讀輸入檔頭幾行，確認實際格式\n"
                        "2. 或把整段程式用 try/except 包起來，except 裡 `import traceback; traceback.print_exc()` "
                        "然後 `sys.exit(0)` 讓錯誤訊息確實印到 stdout\n"
                        "3. 若仍失敗兩次以上，就呼叫 done(success=false) 並在 error 欄位說明你已窮盡哪些方法"
                    )
                ))

            # 2) 連續 3 次任何形式失敗 → 提早中止（但 ask_mode ON 時改成問使用者）
            if consecutive_failures >= 3:
                # ask_mode ON：使用者表態願意被問 → 不要直接 bail，問一下再決定
                # 這修掉實測痛點：ask_mode 勾了、但 agent 從沒 ask 就被早停中止了
                if ask_mode:
                    logger.info(f"[{step_name}] 連續失敗 {consecutive_failures} 次，詢問模式啟用 → 主動問使用者如何繼續")
                    _err_tail = tool_result[-400:] if tool_result else "（無）"
                    answer = await _wait_for_ask_user(
                        run_id=run_id,
                        question=(
                            f"⚠️ Skill agent 連續失敗 {consecutive_failures} 次。\n\n"
                            f"最後一次錯誤：{_err_tail}\n\n"
                            "該如何繼續？"
                        ),
                        options=["繼續嘗試（換策略）", "放棄此步驟"],
                        context="若選『繼續』可在自由輸入補充策略提示，例如「改用 Selenium」「先試 RSS feed」。",
                        logger=logger, step_name=step_name,
                    )
                    if answer is None or "放棄" in answer:
                        logger.info(f"[{step_name}] 使用者選擇放棄（或 ask_user 逾時）")
                        return ExecResult(
                            exit_code=1,
                            stdout="\n".join(all_stdout),
                            stderr=f"使用者選擇放棄此步驟（連續失敗 {consecutive_failures} 次後）",
                            pending_recipe=_pending_recipe,
                            missing_packages=None,
                        )
                    # 使用者選擇繼續：把 answer 當額外提示注入對話
                    consecutive_failures = 0  # 重置計數器讓 agent 繼續
                    was_interactive = True
                    messages.append(HumanMessage(
                        content=f"[使用者補充指示] {answer}\n\n"
                                "請根據以上指示調整策略、不要重複之前失敗的做法。"
                    ))
                    logger.info(f"[{step_name}] 使用者同意繼續，指示：{answer[:100]}")
                    continue  # 回到迭代頂端，不要走下面的 consecutive_failures bail-out
                # ask_mode OFF：照舊行為，直接中止
                logger.error(f"[{step_name}] ⛔ 連續失敗 {consecutive_failures} 次，提早中止避免浪費 token")
                return _attach_trace(ExecResult(
                    exit_code=1,
                    stdout="\n".join(all_stdout),
                    stderr=(
                        f"Skill 連續失敗 {consecutive_failures} 次(累計 {iteration + 1} 次迭代)，提早中止。"
                        f"最後一次錯誤：{tool_result[-500:]}"
                    ),
                ))

        # 超過最大迭代
        logger.warning(f"[{step_name}] Skill agent 達到最大迭代次數")
        return _attach_trace(ExecResult(
            exit_code=1,
            stdout="\n".join(all_stdout),
            stderr=f"Skill agent 在 {SKILL_MAX_ITERATIONS} 次迭代內未完成任務",
        ))

    except Exception as e:
        # 429 / RESOURCE_EXHAUSTED：用獨立 exit code -429 標記，runner 看到後不重試
        # 避免「skill 撞 quota → 步驟 retry → 再撞 quota → 燒光配額」的連環錯
        _err_str = str(e)
        is_quota = ("429" in _err_str or "RESOURCE_EXHAUSTED" in _err_str
                    or "quota" in _err_str.lower() or "rate limit" in _err_str.lower())
        if is_quota:
            logger.error(f"[{step_name}] Skill 執行 LLM 配額/速率受限（429）— 不重試，避免燒光配額：{_err_str[:200]}")
            return _attach_trace(ExecResult(
                exit_code=-429,
                stdout="\n".join(all_stdout),
                stderr=f"LLM provider 配額用盡或速率受限（429）：{_err_str}",
            ))
        logger.error(f"[{step_name}] Skill 執行異常：{e}")
        return _attach_trace(ExecResult(
            exit_code=-3,
            stdout="\n".join(all_stdout),
            stderr=f"Skill 執行異常：{e}",
        ))


# ============================================================
# Skill recipe 儲存(native + text loop 共用、避免 drift)
# ============================================================
def _recipe_code_from(shell_cmd: Optional[str], py_code: Optional[str]) -> Optional[str]:
    """recipe 要存的 code:**優先**用「跑某支 .py 的 run_shell 指令」(= cli-extractor /
    既有 CLI 的實際執行,已含使用者選擇),包成 subprocess 可重播;否則退回最後一段成功的
    run_python code。優先 shell 是因為 agent 常在 run_shell 跑完 CLI 後又補一個只 print
    output 存在的驗證 run_python,若取「最近成功動作」會被驗證步驟蓋掉(E2E 抓到的破綻)。"""
    if shell_cmd:
        return (
            "import subprocess, sys\n"
            f"_p = subprocess.run({shell_cmd!r}, shell=True)\n"
            "sys.exit(_p.returncode)"
        )
    return py_code


def _save_skill_recipe(
    *, pipeline_id, rkey, task_description, input_paths, output_path,
    code, was_interactive, runtime, silent_recipe, no_save_recipe,
    logger, step_name,
):
    """done(success=true) 後存 skill recipe。回傳 _pending_recipe(no_save_recipe 且已有舊 recipe 時、
    需等使用者確認再寫)、否則 None。與 text loop 行為一致:
    - 只在 pipeline_id 是真實 workflow 時存(ad-hoc YAML 的 pipeline_id 非 workflow → FK 約束會炸)
    - silent_recipe:有舊的跳過、無則 seed;no_save_recipe:有舊的延遲、無則建;正常:直接存
    """
    if not (pipeline_id and code):
        return None
    try:
        from db import get_workflow as _gw
        if _gw(pipeline_id) is None:
            return None  # 非真實 workflow → 跳過(本來也不會被前端 cache 重播)
    except Exception:
        return None
    try:
        import sys as _sys2
        from pipeline.recipe import _sha1 as _recipe_sha1, _fingerprint_input as _recipe_fp
        from db import get_recipe as _get_recipe, save_recipe as _db_save_recipe
        _fp = {str(p): _recipe_fp(p) for p in (input_paths or [])}
        task_hash = _recipe_sha1(task_description)
        pyver = f"{_sys2.version_info.major}.{_sys2.version_info.minor}"
        existing = _get_recipe(pipeline_id, rkey)

        def _save():
            _db_save_recipe(pipeline_id, rkey, task_hash, _fp, output_path, code,
                            pyver, runtime, was_interactive=was_interactive)

        if silent_recipe:
            if existing:
                logger.info(f"[{step_name}] silent_recipe:Recipe 已存在、跳過(不覆寫)")
            else:
                _save(); logger.info(f"[{step_name}] silent_recipe:首次建立 Recipe")
        elif no_save_recipe:
            if existing:
                logger.info(f"[{step_name}] Recipe 已存在、延遲儲存等待確認")
                return {
                    "pipeline_id": pipeline_id, "step_name": rkey, "task_hash": task_hash,
                    "input_fingerprints": _fp, "output_path": output_path, "code": code,
                    "python_version": pyver, "runtime_sec": runtime,
                    "was_interactive": was_interactive,
                }
            _save(); logger.info(f"[{step_name}] 首次建立 Recipe")
        else:
            _save(); logger.info(f"[{step_name}] ✅ Recipe 已儲存(native skill loop)")
    except Exception as e:
        logger.warning(f"[{step_name}] Recipe 儲存失敗:{e}")
    return None


# ============================================================
# Phase A.2 — SKILL loop native function calling
# ============================================================
async def _execute_skill_native_loop(
    *,
    llm,
    messages: list,
    output_path: Optional[str],
    working_dir: Optional[str],
    run_id: str,
    step_name: str,
    logger: logging.Logger,
    tool_timeout: int,
    ask_mode: bool,
    acc_usage: dict,
    acc_tool_calls: list,
    all_stdout: list,
    pipeline_id: Optional[str] = None,
    task_description: str = "",
    input_paths: Optional[list] = None,
    no_save_recipe: bool = False,
    silent_recipe: bool = False,
    recipe_step_key: Optional[str] = None,
) -> ExecResult:
    """Phase A.2 — SKILL loop 用 LangChain bind_tools() native function calling 版本。

    跟 text 版同 ExecResult 回傳格式、execute_step_with_skill 看 SUBAGENT_LOOP_MODE=native 就 dispatch 過來。
    保留守門:done preflight(output 存在 / run_python 成功 / run_shell 成功)、consecutive_no_tool。
    省略(MVP):recipe save(下次同任務不快取、LLM 重跑)、visual_validation 由外層處理。
    """
    from langchain_core.messages import AIMessage, ToolMessage, HumanMessage as _HM
    from pipeline.sandbox_tools import build_subagent_tools
    import time as _time  # recipe runtime 計算用(native loop 區域、與 text loop 對齊)

    # SKILL 完整工具集 = 跟 _execute_skill_tool dispatch + system prompt 教的工具對齊。
    # (含 write_file/edit_file/grep/glob/view_image —— 這些 prompt 有教、dispatch 接得住,
    #  之前 native 遷移漏註冊導致一呼叫就被擋,屬與 export_var 同類的 native FC 副作用。)
    # view_image 由 native loop 特判 tc_name 走多模態注入(見下方 view_image 處理)。
    allowed_tools = {"run_python", "run_shell", "read_file", "write_file", "edit_file",
                     "grep", "glob", "view_image", "ask_user", "export_var", "done"}
    # web_search 與 system prompt 的「條件揭露」對齊:沒開 / 沒 tavily key 就不註冊,
    # 避免 native 模式下工具 schema 仍餵給 LLM、被呼叫卻無效(gate 失效)。
    try:
        from settings import get_settings as _gs_ws
        _ws_cfg = _gs_ws()
        if _ws_cfg.get("web_search_enabled", True) and _ws_cfg.get("tavily_api_key"):
            allowed_tools.add("web_search")
    except Exception:
        allowed_tools.add("web_search")  # 取不到設定 → 保守保留(維持舊行為)
    web_counter = {"count": 0}

    tools = build_subagent_tools(
        cwd=working_dir, run_id=run_id, logger=logger,
        tool_timeout=tool_timeout, allowed_tool_names=allowed_tools,
        step_name=step_name, web_search_counter=web_counter,
    )
    name_to_tool = {t.name: t for t in tools}

    try:
        llm_with_tools = llm.bind_tools(tools)
    except Exception as e:
        logger.error(f"[{step_name}] bind_tools 失敗:{e}")
        return ExecResult(
            exit_code=-1, stdout="\n".join(all_stdout) or "",
            stderr=f"SKILL bind_tools 失敗: {e}",
        )

    # Native FC override — SKILL system_prompt(execute_step_with_skill line 2181 起)教 LLM
    # 用 <tool>name</tool> 文字協議,跟 bind_tools 的 native function_declarations 衝突。
    # Gemma 4 等模型會選文字協議、回 content 含 <tool>...</tool> 純文字、tool_calls=[]。
    # 在 system_prompt 結尾加 override 強制走 native。
    if messages and isinstance(messages[0].content, str):
        _native_override = (
            "\n\n## ⛔ 最高優先級 — Tool 呼叫協議(本對話使用)\n"
            "本對話**使用 native function calling API**、tool 已透過 function_declarations 註冊。\n"
            "你**必須**透過 API 的 function_call 機制呼叫工具、**禁止**寫 `<tool>name</tool>` 文字格式。\n"
            "✓ 正確:直接 emit tool_calls(API 結構化欄位)、不要在 reply 文字內寫 `<tool>` tag。\n"
            "✗ 錯誤:寫 `<tool>run_python</tool>\\n```python\\n...` 純文字 — orchestrator 不會解析、會被視為沒呼叫 tool。\n"
            "上面 system prompt 內若有 `<tool>...</tool>` 範例、那是文字協議的舊範例、本次 native 模式請忽略格式、保留語意(該用哪個 tool / 何時 done)。\n"
            "\n## 🔁 失敗重試規則(重要)\n"
            "若某次 run_python / run_shell 失敗(回傳含 `[exit code: N]` / Error / traceback)、"
            "**禁止把相同的程式碼或命令原封不動再送一次** — 那只會得到一模一樣的錯誤、純粹浪費。\n"
            "你必須:① 先讀錯誤訊息找真因 → ② **實質改寫**(換 API 用法 / 修語法 / 補缺套件 / 換做法)再試;"
            "③ 若連續嘗試都修不動、坦白呼叫 done(success=false)說明卡在哪、不要硬撐重送相同內容。\n"
        )
        messages[0] = SystemMessage(
            content=messages[0].content + _native_override,
            additional_kwargs=getattr(messages[0], "additional_kwargs", {}) or {},
        )

    last_run_python_ok: Optional[bool] = None
    last_run_shell_ok: Optional[bool] = None
    last_successful_code: Optional[str] = None  # 最後一段成功 run_python code → 供 recipe 儲存
    last_successful_shell: Optional[str] = None  # 最後一條成功「跑 .py 腳本」的 run_shell(CLI 執行)→ recipe 優先用它
    was_interactive = False                      # 過程用過 ask_user → recipe 標記(replay 時提醒可能要再問)
    _recipe_start = _time.time()                 # 算 recipe runtime_sec 用
    consecutive_no_tool = 0
    prose_before_tool_violations = 0   # task #160:有 tool_calls 但 content 過長累計
    fake_done_count = 0
    _FAKE_DONE_LIMIT_SKILL = 3

    # ── 不收斂守門(task #187、2026-05-29):偵測「相同 tool input 重複失敗」死循環 ──
    # 真實案例:Gemma 寫的 docx-js code 有 API bug、node rc=1,它每輪重寫一模一樣的
    # 9230 字 code、又 rc=1、鬼打牆到 max_iter(空轉 6 分鐘燒 quota)。
    # 同一份失敗 input hash:第 2 次→注入「別重送」強提示、第 3 次→提早中止 + 診斷。
    repeat_fail_hashes: dict[str, int] = {}
    _REPEAT_FAIL_REMIND_AT = 2
    _REPEAT_FAIL_ABORT_AT = 3

    # Output-ready 自動 done 偵測:LLM 寫成 output 檔後常持續 verify / thumbnail / preview
    # 不主動 done(2026-05-26 PPT workflow 真實案例:iter 11 寫成、iter 12-20 全空轉燒 ~770K tokens)。
    #
    # 誤殺防護(mtime-based):
    # - 檔案 mtime 變動 = LLM 真的在改檔 / polish → reset counter(允許繼續優化)
    # - 檔案 mtime 沒變 = LLM 在做別的(verify / thumbnail / ls)→ counter++
    # 這樣只會擋「真正空轉」、不會誤殺「正在 polish」的合理場景。
    output_ready_no_done_count = 0
    last_output_mtime: Optional[float] = None
    _OUTPUT_READY_REMINDER_AT = 1   # 第幾輪 ready+mtime未變後注入 reminder
    _OUTPUT_READY_FORCE_AT = 4       # 第幾輪 ready+mtime未變後強制收尾(從 3 調 4、更保守)
    _OFFICE_EXTS_SKILL = {".pptx", ".docx", ".xlsx"}
    _OFFICE_MIN_BYTES_SKILL = 5000
    _GENERIC_MIN_BYTES_SKILL = 100

    def _output_is_ready() -> tuple[bool, int, Optional[float]]:
        """檢查 output_path 是否存在且 size 達門檻。回 (ready, size, mtime)。"""
        if not output_path:
            return False, 0, None
        p = Path(output_path)
        if not p.exists():
            return False, 0, None
        try:
            st = p.stat()
        except OSError:
            return False, 0, None
        floor = (
            _OFFICE_MIN_BYTES_SKILL
            if p.suffix.lower() in _OFFICE_EXTS_SKILL
            else _GENERIC_MIN_BYTES_SKILL
        )
        return st.st_size >= floor, st.st_size, st.st_mtime

    _RETRIABLE = (
        "503", "429", "unavailable", "rate limit", "rate_limit",
        "overloaded", "internal error", "500", "deadline", "resource_exhausted",
        # 網路瞬斷 / DNS 抖動 — 退避重連、不直接判失敗
        "apiconnectionerror", "connection error", "connection refused",
        "connection reset", "connection aborted", "getaddrinfo",
        "timed out", "timeout", "temporarily unavailable", "econnreset",
        "remotedisconnected", "max retries",
    )

    for iteration in range(SKILL_MAX_ITERATIONS):
        logger.info(f"[{step_name}] Skill 執行迭代 {iteration + 1}/{SKILL_MAX_ITERATIONS} [NATIVE]")

        # Context 雪崩防護(跟 text 版同邏輯)
        if iteration >= SKILL_CONTEXT_KEEP_RECENT_FULL + 1:
            _compacted = _compact_old_tool_results(messages)
            if _compacted:
                logger.debug(f"[{step_name}] 🪶 壓縮 {_compacted} 條舊 tool 結果(防 context 膨脹)")

        # 第一輪後 strip FIRST_ITER 區塊(省 prompt token)
        if iteration == 1 and isinstance(messages[0].content, str) and "<!--FIRST_ITER_BEGIN-->" in messages[0].content:
            _stripped = re.sub(
                r"<!--FIRST_ITER_BEGIN-->.*?<!--FIRST_ITER_END-->\s*",
                "", messages[0].content, flags=re.DOTALL,
            )
            messages[0] = SystemMessage(
                content=_stripped,
                additional_kwargs=getattr(messages[0], "additional_kwargs", {}) or {},
            )
            logger.debug(f"[{step_name}] 🪶 第二輪起 system prompt 已 strip first-iter blocks")

        if iteration > 0:
            await asyncio.sleep(SKILL_REQUEST_INTERVAL)

        # LLM call with retry
        response: Optional[AIMessage] = None
        last_err: Optional[Exception] = None
        _input_chars = sum(len(str(getattr(m, "content", "") or "")) for m in messages)
        for _attempt in range(3):
            try:
                logger.info(f"[{step_name}] 🤖 LLM 開始處理(input {_input_chars:,} 字)…")
                _t0 = asyncio.get_event_loop().time()
                response = await asyncio.wait_for(
                    llm_with_tools.ainvoke(messages), timeout=600.0,
                )
                _el = asyncio.get_event_loop().time() - _t0
                _content_str = _extract_text(response.content)
                _tc_count = len(getattr(response, 'tool_calls', []) or [])
                # Cache stats(Anthropic / OpenAI 有,Gemini / Groq / Ollama 沒)
                _um_this = getattr(response, "usage_metadata", None) or {}
                _itd_this = _um_this.get("input_token_details") if isinstance(_um_this, dict) else None
                _cache_read = (_itd_this or {}).get("cache_read", 0) or 0
                _cache_create = (_itd_this or {}).get("cache_creation", 0) or 0
                _in_tok_this = _um_this.get("input_tokens", 0) or 0
                _cache_str = ""
                if _cache_read or _cache_create:
                    _total_prompt = _in_tok_this + _cache_read + _cache_create
                    _hit_pct = (_cache_read / _total_prompt * 100) if _total_prompt > 0 else 0
                    _cache_str = f", cache_read {_cache_read:,} ({_hit_pct:.0f}%), cache_write {_cache_create:,}"
                logger.info(
                    f"[{step_name}] ✅ LLM 完成({_el:.0f}s, content {len(_content_str)} 字, "
                    f"tool_calls={_tc_count}{_cache_str})"
                )
                # 診斷:LLM 回 content + tool_calls 都空 → dump 完整 response 看哪個 field 真有輸出
                if len(_content_str) == 0 and _tc_count == 0:
                    logger.warning(
                        f"[{step_name}] ⚠ 空回應診斷:\n"
                        f"  response.content = {response.content!r}\n"
                        f"  response.additional_kwargs = {dict(getattr(response, 'additional_kwargs', None) or {})}\n"
                        f"  response.response_metadata = {dict(getattr(response, 'response_metadata', None) or {})}\n"
                        f"  response.usage_metadata = {dict(getattr(response, 'usage_metadata', None) or {})}\n"
                        f"  response.tool_calls (raw) = {getattr(response, 'tool_calls', None)!r}\n"
                        f"  response.invalid_tool_calls = {getattr(response, 'invalid_tool_calls', None)!r}"
                    )
                # MALFORMED_FUNCTION_CALL = 壞生成、可重試(不是模型不肯呼叫工具)。
                # 用既有 3 次預算重試;耗盡才落回下游 no-tool 處理(不比修前差)。
                if _is_malformed_empty(response) and _attempt < 2:
                    _w = 2 ** _attempt
                    logger.warning(
                        f"[{step_name}] finish_reason=MALFORMED_FUNCTION_CALL(壞生成、非不呼叫工具)"
                        f"→ {_w}s retry({_attempt + 1}/2)"
                    )
                    response = None
                    await asyncio.sleep(_w)
                    continue
                last_err = None
                break
            except asyncio.TimeoutError as e:
                last_err = e
                if _attempt < 2:
                    _w = 2 ** _attempt
                    logger.warning(f"[{step_name}] LLM 逾時、{_w}s retry({_attempt + 1}/2)")
                    await asyncio.sleep(_w)
                    continue
                break
            except Exception as e:
                last_err = e
                if any(k in str(e).lower() for k in _RETRIABLE):
                    if _attempt < 2:
                        _w = 2 ** _attempt
                        logger.warning(
                            f"[{step_name}] LLM 暫時錯誤、{_w}s retry({_attempt + 1}/2):"
                            f"{type(e).__name__}: {str(e)[:200]}"
                        )
                        await asyncio.sleep(_w)
                        continue
                break

        if last_err is not None or response is None:
            err = f"LLM 失敗: {type(last_err).__name__ if last_err else 'None'}: {last_err}"
            logger.error(f"[{step_name}] {err}")
            return ExecResult(
                exit_code=-2, stdout="\n".join(all_stdout), stderr=err,
            )

        # 累計 token usage
        um = getattr(response, "usage_metadata", None) or {}
        if isinstance(um, dict) and um:
            for _k in ("input_tokens", "output_tokens", "total_tokens"):
                _v = um.get(_k) or 0
                if _v:
                    acc_usage[_k] = acc_usage.get(_k, 0) + int(_v)
            itd = um.get("input_token_details") or {}
            if isinstance(itd, dict):
                acc_usage["cache_read_tokens"] = (
                    acc_usage.get("cache_read_tokens", 0) + (itd.get("cache_read", 0) or 0)
                )
                acc_usage["cache_creation_tokens"] = (
                    acc_usage.get("cache_creation_tokens", 0) + (itd.get("cache_creation", 0) or 0)
                )
        if not acc_usage.get("model"):
            _rm = getattr(response, "response_metadata", None) or {}
            acc_usage["model"] = _rm.get("model_name") or _rm.get("model") or ""

        # AIMessage 加進 messages(含 tool_calls)
        messages.append(response)

        tool_calls = list(getattr(response, "tool_calls", []) or [])
        content_str = _extract_text(response.content)

        # 沒 tool_calls → Claude 原生 end_turn(想結束)。對齊 Anthropic 官方:
        # 完成時就不呼叫工具、回純文字,這是天生的結束信號、不是異常。
        # 若 output 檔已 ready → 視為正常完成,別當 consecutive_no_tool 懲罰逼它空轉。
        if not tool_calls:
            _ready_et, _sz_et, _ = _output_is_ready()
            if _ready_et:
                logger.info(
                    f"[{step_name}] ✅ LLM 回 end_turn(純文字結束)且輸出檔已 ready"
                    f"({_sz_et:,} bytes)→ 視為完成(對齊原生 end_turn 結束)"
                )
                return ExecResult(
                    exit_code=0,
                    stdout=_build_clean_success_stdout(all_stdout, "[Skill 完成]"),
                    stderr="",
                )
            consecutive_no_tool += 1
            logger.warning(
                f"[{step_name}] ⚠ 第 {iteration + 1} 輪沒 tool_calls"
                f"(content {len(content_str)} 字)、累計 {consecutive_no_tool}/2"
            )
            if consecutive_no_tool >= 2:
                err_msg = (
                    f"連 2 輪沒呼叫任何 tool、強制中止。"
                    f"LLM 把分析寫 content 而非 run_python 寫檔。"
                )
                logger.error(f"[{step_name}] ✗ {err_msg}")
                return ExecResult(
                    exit_code=1,
                    stdout="\n".join(all_stdout) or content_str[:500],
                    stderr=err_msg,
                    agent_concluded_fail=True,
                )
            messages.append(_HM(content=(
                "請呼叫一個 tool(run_python / read_file / 等)繼續、或 done 結束。"
                "分析請寫進 run_python 程式碼、不要寫 content。"
            )))
            continue
        consecutive_no_tool = 0

        # 把 LLM 這一輪的決策旁白(content)記進 log。native FC 下這段原本只計字數、
        # 看不到內容,長任務無法回溯「它在想什麼 / 為什麼這樣做」。實測 tool 輪的 content
        # 一定 < 3000 字(prose-before-tool 守門在 1000+ log 從沒觸發),8000 字寬鬆上限
        # 只擋病態案例、並保留誠實截斷標記。注意:深層 reasoning(thinking token)由
        # llm_factory 消耗後丟棄、不在 content 內,這裡顯示的是「簡短決策旁白」。
        _prose = (content_str or "").strip()
        if _prose:
            _pl = len(_prose)
            _shown = _prose if _pl <= 8000 else _prose[:8000] + f"\n…[截斷、完整 {_pl} 字]"
            logger.info(f"[{step_name}] 💬 {_shown}")

        # ── Prose-before-tool 守門(task #160、backport 自 subagent)──
        # 有 tool_calls 但伴隨超長 content(Sonnet 4.6 寫 5 萬字 parser 案例)→ 燒 output token、
        # 拖斷長連線。soft 只 log、enforce 連 _LIMIT 輪違規才中止。
        if _SKILL_PROSE_CAP_MODE != "off":
            _prose_chars = len((content_str or "").strip())
            if _prose_chars > _SKILL_PROSE_BEFORE_TOOL_THRESHOLD:
                prose_before_tool_violations += 1
                logger.warning(
                    f"[{step_name}] ⚠ 第 {iteration + 1} 輪 tool 前 content {_prose_chars:,} 字"
                    f"(上限 {_SKILL_PROSE_BEFORE_TOOL_THRESHOLD})、累計違規 "
                    f"{prose_before_tool_violations}/{_SKILL_PROSE_BEFORE_TOOL_LIMIT} "
                    f"[mode={_SKILL_PROSE_CAP_MODE}]"
                )
                if (_SKILL_PROSE_CAP_MODE == "enforce"
                        and prose_before_tool_violations >= _SKILL_PROSE_BEFORE_TOOL_LIMIT):
                    err_msg = (
                        f"連續 {_SKILL_PROSE_BEFORE_TOOL_LIMIT} 輪 tool 前 content 超過 "
                        f"{_SKILL_PROSE_BEFORE_TOOL_THRESHOLD} 字、強制中止避免燒 token / 拖斷連線。"
                        f"解析請直接寫進 run_python 程式碼、不要在 content 寫長篇 parser / 解釋。"
                    )
                    logger.error(f"[{step_name}] ✗ {err_msg}")
                    return ExecResult(
                        exit_code=1,
                        stdout="\n".join(all_stdout) or content_str[:500],
                        stderr=err_msg,
                        agent_concluded_fail=True,
                    )
            else:
                # 沒違規這輪 → 重置(只擋「連續」冗長、跟 subagent 一致)
                prose_before_tool_violations = 0

        # 處理 tool_calls — 先掃出 done(若有)、其他先跑
        done_call: Optional[dict] = None
        regular_calls = []
        for tc in tool_calls:
            if tc.get("name") == "done":
                done_call = tc  # 多個 done 取最後
            else:
                regular_calls.append(tc)

        _repeat_fail_reminder: Optional[str] = None
        for tc in regular_calls:
            tc_name = tc.get("name", "")
            tc_id = tc.get("id") or ""
            tc_args = tc.get("args", {}) or {}
            # native FC 的「意圖」在結構化 tool_call 裡(content 常為空)→ 必須印出實際參數,
            # 否則 log 只剩 args_keys、看不到 LLM 到底要跑什麼指令 / 程式碼(改 native 前的可見度)。
            _pv = (tc_args.get("command") or tc_args.get("code") or tc_args.get("question")
                   or tc_args.get("path") or tc_args.get("query") or "")
            _pv = str(_pv).rstrip()
            _full_len = len(_pv)
            if _full_len > 1500:
                _pv = _pv[:1500] + f"\n…[截斷、完整 {_full_len} 字]"
            if _pv:
                # 保留原始換行 → log 檢視器多行整齊顯示(對齊舊 text 模式的可讀性),
                # 不要壓成單行;延續行沒有 [INFO] 等級標記 → 前端上中性灰、不會誤判紅字。
                logger.info(f"[{step_name}] 🛠 tool={tc_name}（{_full_len} 字）:\n{_pv}")
            else:
                logger.info(f"[{step_name}] 🛠 tool={tc_name} (keys={list(tc_args.keys())})")

            # ── native loop:pip install 企圖(run_shell command + run_python code)→ 轉 missing_dependency ──
            # (共用 helper detect_pip_install:涵蓋 run_shell 的 pip 命令 + run_python code 內的
            #  subprocess/os.system/pip.main/['pip','install'] 繞法。skill 是 pipeline step → 直接轉確認。)
            _pip_pkgs = detect_pip_install(tc_name, tc_args)
            if _pip_pkgs:
                logger.warning(f"[{step_name}] 🛑 攔 native pip install {_pip_pkgs} → 轉 missing_dependency 使用者確認")
                _real_pkgs = [p for p in _pip_pkgs if p != "(未知套件)"]
                return ExecResult(
                    exit_code=1,
                    stdout="\n".join(all_stdout),
                    stderr=f"偵測到嘗試自行安裝套件 {_pip_pkgs}。已改走系統的『允許安裝』使用者確認流程。",
                    missing_packages=_real_pkgs or None,
                )

            # run_shell 用:跑工具前先記 output 檔 mtime,跑完比對 → 只認「真的產出/更新 output」的
            # .py 指令當 recipe(排除 --help / ls / 探查 / 驗證,那些不動 output 檔)。
            _out_mtime_before = None
            if tc_name == "run_shell" and output_path:
                try:
                    _out_mtime_before = Path(output_path).stat().st_mtime
                except Exception:
                    _out_mtime_before = None

            # view_image 特判:走多模態注入(讀圖→base64→image_url),不走泛用 dispatch
            # (泛用 dispatch 只會回 "__VIEW_IMAGE__" sentinel、對 LLM 無意義)。
            if tc_name == "view_image":
                _vi_path = tc_args.get("path", "") if isinstance(tc_args, dict) else str(tc_args)
                _vi = await asyncio.get_event_loop().run_in_executor(None, _skill_view_image, _vi_path)
                _vi_text = _vi.get("text", "") if isinstance(_vi, dict) else str(_vi)
                logger.info(f"[{step_name}] view_image:{_vi_text}")
                all_stdout.append(f"[view_image] {_vi_text}")
                acc_tool_calls.append({
                    "name": "view_image",
                    "input_preview": json.dumps(tc_args, ensure_ascii=False)[:200],
                    "result_preview": _vi_text[:300],
                })
                # 先回 ToolMessage 滿足 tool_call 配對,再 append 帶圖的 HumanMessage(視覺模型才看得到圖)
                messages.append(ToolMessage(content=f"[view_image] {_vi_text}", tool_call_id=tc_id))
                if isinstance(_vi, dict) and _vi.get("image_b64"):
                    messages.append(_HM(content=[
                        {"type": "text", "text": f"[view_image 圖片]\n{_vi_text}\n請觀察圖片內容後再決定下一步。"},
                        {"type": "image_url", "image_url": {
                            "url": f"data:{_vi.get('image_mime', 'image/png')};base64,{_vi['image_b64']}"
                        }},
                    ]))
                continue

            tool_fn = name_to_tool.get(tc_name)
            try:
                if tool_fn is None:
                    result = f"[錯誤] tool '{tc_name}' 不在白名單。可用:{sorted(name_to_tool.keys())}"
                else:
                    raw = await tool_fn.ainvoke(tc_args)
                    result = str(raw) if raw is not None else ""
            except Exception as e:
                result = f"[執行失敗] {type(e).__name__}: {e}"

            # 追蹤 last_run_python_ok / last_run_shell_ok(守門用)+ was_interactive(recipe 用)
            if tc_name == "run_python":
                last_run_python_ok = "[exit code:" not in result
                if last_run_python_ok:
                    last_successful_code = tc_args.get("code", "")
            elif tc_name == "run_shell":
                last_run_shell_ok = "[exit code:" not in result
                if last_run_shell_ok:
                    # cli-extractor / 既有 CLI 這種「靠 run_shell 跑 <某>.py <args>」的 skill:
                    # 只記「執行某支 .py 腳本、且真的寫出/更新了 output 檔」的 run_shell
                    # (= CLI 實際執行、已含使用者選擇)。用 mtime 比對排除 --help / ls / 探查 /
                    # 驗證類(那些不動 output);recipe save 時**優先**用它,避免被後面補的
                    # 驗證 run_python(只 print output 存在)蓋掉(E2E 抓到的破綻)。
                    _cmd = str(tc_args.get("command", "") or "")
                    if ".py" in _cmd:
                        if output_path:
                            try:
                                _mt_after = Path(output_path).stat().st_mtime
                            except Exception:
                                _mt_after = None
                            if _mt_after is not None and (
                                _out_mtime_before is None or _mt_after > _out_mtime_before
                            ):
                                last_successful_shell = _cmd  # 這條真的產出了 output
                        else:
                            # 無 output_path 可比對 → 退回寬鬆(仍優於完全不抓)
                            last_successful_shell = _cmd
            elif tc_name == "ask_user":
                was_interactive = True  # 用過人工問答 → recipe 標記(replay 時提醒可能要再問)

            # ── ModuleNotFoundError 早期攔截(共用 helper detect_missing_module)──
            # 傳 task_description:讓 import 名→pip 名 能用任務文字當線索(如任務寫 openai-whisper)
            _miss = detect_missing_module(result, task_description)
            if _miss:
                logger.warning(f"[{step_name}] 🛑 native 偵測 ModuleNotFoundError: {_miss} → 轉 missing_dependency 使用者確認")
                return ExecResult(
                    exit_code=1,
                    stdout="\n".join(all_stdout),
                    stderr=f"ModuleNotFoundError: 缺少套件 {_miss}。系統將彈出安裝確認對話框（TG / 前端 modal）。",
                    missing_packages=_miss,
                )

            # ── 不收斂守門:相同失敗 input 重複偵測 ──
            _call_failed = ("[exit code:" in result) or result.startswith("[執行失敗]")
            if _call_failed and tc_name in ("run_python", "run_shell"):
                _fh = hashlib.md5(
                    (tc_name + "\x00"
                     + json.dumps(tc_args, ensure_ascii=False, sort_keys=True))
                    .encode("utf-8", "replace")
                ).hexdigest()
                repeat_fail_hashes[_fh] = repeat_fail_hashes.get(_fh, 0) + 1
                _rep = repeat_fail_hashes[_fh]
                if _rep >= _REPEAT_FAIL_ABORT_AT:
                    err_msg = (
                        f"相同的 {tc_name} input 連續失敗 {_rep} 次(內容完全相同、錯誤相同)、"
                        f"判定不收斂、提早中止避免空轉燒 token。"
                        f"可能此任務超出當前模型能力(常見:docx-js 等嚴格 API)、"
                        f"建議改用強模型或換做法。"
                    )
                    logger.error(f"[{step_name}] ✗ 不收斂守門觸發:{err_msg}")
                    return ExecResult(
                        exit_code=1,
                        stdout="\n".join(all_stdout),
                        stderr=err_msg,
                        agent_concluded_fail=True,
                    )
                if _rep >= _REPEAT_FAIL_REMIND_AT and _repeat_fail_reminder is None:
                    _repeat_fail_reminder = (
                        f"[系統] ⚠ 你已用**完全相同**的 {tc_name} 內容失敗 {_rep} 次。"
                        f"再送一次必然得到同樣的錯。請先讀上面錯誤訊息找真因、"
                        f"**實質改寫**(換 API 用法 / 修語法 / 補缺套件)再試;"
                        f"若真的修不動就呼叫 done(success=false)說明卡點、別再重送相同內容。"
                    )
                    logger.warning(
                        f"[{step_name}] 🔁 相同 {tc_name} 失敗第 {_rep} 次、注入「別重送」提示"
                    )

            all_stdout.append(f"[{tc_name}] {result[:500]}")
            acc_tool_calls.append({
                "name": tc_name,
                "input_preview": json.dumps(tc_args, ensure_ascii=False)[:200],
                "result_preview": result[:300],
            })

            # 截斷大 result 防 context 雪崩
            _MAX = 3000 if tc_name == "read_file" else 5000
            _HEAD = 2000 if tc_name == "read_file" else 4000
            _TAIL = 1000
            if len(result) > _MAX:
                _h = result[:_HEAD]
                _t = result[-_TAIL:]
                result = (
                    f"{_h}\n…[中間省略 {len(result) - _HEAD - _TAIL} 字、"
                    f"完整長度 {len(result)}]…\n{_t}"
                )

            messages.append(ToolMessage(content=result, tool_call_id=tc_id))

        # 不收斂守門:本輪偵測到相同失敗重送 → 注入一次「別重送」強提示
        # (放在所有 tool 結果後、維持 tool_call/tool_result 配對、再加 human turn)
        if _repeat_fail_reminder is not None:
            messages.append(HumanMessage(content=_repeat_fail_reminder))

        # Output-ready 自動 done 偵測 — mtime-based、區分「真 polish」vs「純空轉」
        if done_call is None:
            _ready, _size, _mtime = _output_is_ready()
            if _ready:
                # mtime 有變 = LLM 真的在改檔(可能是 polish)→ 重置 counter、允許繼續
                if last_output_mtime is None or (_mtime is not None and _mtime > last_output_mtime):
                    if last_output_mtime is not None:
                        logger.info(
                            f"[{step_name}] 📝 Output 檔有更新(mtime 變動、{_size:,} bytes)、"
                            f"LLM 仍在 polish → counter 重置"
                        )
                    last_output_mtime = _mtime
                    output_ready_no_done_count = 0
                else:
                    # mtime 沒變 = LLM 在做別的事(verify / thumbnail / ls)→ 累計
                    output_ready_no_done_count += 1
                    logger.info(
                        f"[{step_name}] 📦 Output ready 但 mtime 未變({_size:,} bytes)、"
                        f"LLM 未動檔也未 done、累計 {output_ready_no_done_count}/{_OUTPUT_READY_FORCE_AT}"
                    )
                    if output_ready_no_done_count >= _OUTPUT_READY_FORCE_AT:
                        logger.warning(
                            f"[{step_name}] ⚠ Output 已 ready 且 mtime 未變 {output_ready_no_done_count} 輪、"
                            f"LLM 仍未 done — 強制 success 收尾。檔案:{output_path}({_size:,} bytes)"
                        )
                        all_stdout.append(
                            f"[Skill 完成] 系統強制收尾:輸出檔 {Path(output_path).name} "
                            f"({_size:,} bytes)已 ready 且 {output_ready_no_done_count} 輪未動、LLM 未 done"
                        )
                        _pending_fs = None
                        _rc_code = _recipe_code_from(last_successful_shell, last_successful_code)
                        if _rc_code:
                            _pending_fs = _save_skill_recipe(
                                pipeline_id=pipeline_id, rkey=(recipe_step_key or step_name),
                                task_description=task_description, input_paths=input_paths,
                                output_path=output_path, code=_rc_code,
                                was_interactive=was_interactive, runtime=_time.time() - _recipe_start,
                                silent_recipe=silent_recipe, no_save_recipe=no_save_recipe,
                                logger=logger, step_name=step_name,
                            )
                        return ExecResult(
                            exit_code=0,
                            stdout=_build_clean_success_stdout(all_stdout, "[Skill 完成]"),
                            stderr="",
                            pending_recipe=_pending_fs,
                        )
                    if output_ready_no_done_count >= _OUTPUT_READY_REMINDER_AT:
                        messages.append(HumanMessage(content=(
                            f"[系統] ✅ 目標檔案 {output_path} 已存在({_size:,} bytes)且本輪未改檔。"
                            f"如果你還在 polish — 直接寫檔(run_python 改檔)、系統會偵測 mtime 變動允許繼續。"
                            f"如果你只是想 verify / 看 thumbnail / 確認大小 — 那是多餘的、**請立刻呼叫 done** 結束 step。"
                            f"連續 {_OUTPUT_READY_FORCE_AT} 輪不動檔也不 done 會被強制收尾。"
                        )))
            else:
                # 沒 ready → 重置 counter + mtime tracking(允許 LLM 繼續寫檔嘗試)
                output_ready_no_done_count = 0
                last_output_mtime = None

        # 處理 done(在 regular_calls 全跑完後)
        if done_call is not None:
            tc_id = done_call.get("id") or ""
            tc_args = done_call.get("args", {}) or {}
            _success = bool(tc_args.get("success", True))
            _summary = (
                tc_args.get("summary")
                or tc_args.get("error")
                or "(空 summary)"
            )

            # 守門 1:run_python 失敗 → 拒 done(防 LLM 在程式炸掉後硬送 success=true)
            if _success and last_run_python_ok is False:
                logger.warning(f"[{step_name}] done 被拒:最近 run_python 失敗")
                messages.append(ToolMessage(
                    content=(
                        "[拒收] 最近一次 run_python 執行失敗、看上面 stderr / traceback。"
                        "請先用 run_python 修正、確認沒 [exit code: N] 才能 done。"
                    ),
                    tool_call_id=tc_id,
                ))
                continue

            # 守門 1b:run_shell 失敗 → 拒 done
            if _success and last_run_shell_ok is False:
                logger.warning(f"[{step_name}] done 被拒:最近 run_shell 失敗")
                messages.append(ToolMessage(
                    content=(
                        "[拒收] 最近一次 run_shell 失敗。請修正命令(常見:路徑、模組缺裝、quote)、"
                        "確認 rc=0 沒 Error 後才能 done。"
                    ),
                    tool_call_id=tc_id,
                ))
                continue

            # 守門 2:output_path 不存在 → 拒 done + surgical retry
            if (
                _success and output_path
                and not Path(output_path).exists()
                and fake_done_count < _FAKE_DONE_LIMIT_SKILL
            ):
                fake_done_count += 1
                logger.warning(
                    f"[{step_name}] done 被拒:output 檔 {output_path} 不存在"
                    f"({fake_done_count}/{_FAKE_DONE_LIMIT_SKILL})"
                )
                messages.append(ToolMessage(
                    content=(
                        f"[拒收] 你宣稱成功但 output 檔 {output_path} 不存在!"
                        f"請用 run_python 真實寫檔到那個路徑(Path(...).write_text() / "
                        f"df.to_excel() / fig.savefig() 等)、跑完 print Path('{output_path}').exists() "
                        f"確認 True 才能 done。不准只展示 code、必須真跑。"
                    ),
                    tool_call_id=tc_id,
                ))
                continue

            # 通過 → 結束 loop
            messages.append(ToolMessage(content="__DONE_ACCEPTED__", tool_call_id=tc_id))
            acc_tool_calls.append({
                "name": "done",
                "input_preview": json.dumps(tc_args, ensure_ascii=False)[:200],
                "result_preview": "",
            })
            all_stdout.append(f"[Skill 完成] {_summary}")
            logger.info(
                f"[{step_name}] ✅ Skill 完成(success={_success}, "
                f"summary 前 80 字={_summary[:80]})"
            )

            # 成功 → 存 recipe(下次同任務 + 同輸入指紋直接 replay、省 LLM token)
            # native loop 原本省略了這段(MVP)→ native 成預設後 skill recipe 整個沒在存,這裡補回。
            _pending = None
            _rc_code = _recipe_code_from(last_successful_shell, last_successful_code)
            if _success and _rc_code:
                _pending = _save_skill_recipe(
                    pipeline_id=pipeline_id, rkey=(recipe_step_key or step_name),
                    task_description=task_description, input_paths=input_paths,
                    output_path=output_path, code=_rc_code,
                    was_interactive=was_interactive, runtime=_time.time() - _recipe_start,
                    silent_recipe=silent_recipe, no_save_recipe=no_save_recipe,
                    logger=logger, step_name=step_name,
                )

            _final_stdout = (
                _build_clean_success_stdout(all_stdout, "[Skill 完成]")
                if _success else "\n".join(all_stdout)
            )

            return ExecResult(
                exit_code=0 if _success else 1,
                stdout=_final_stdout,
                stderr="" if _success else _summary,
                agent_concluded_fail=(not _success),
                pending_recipe=_pending,
            )

    # max_iter 用完未 done
    err_msg = f"max_iter ({SKILL_MAX_ITERATIONS}) 用完未 done"
    logger.error(f"[{step_name}] {err_msg}")
    return ExecResult(
        exit_code=1,
        stdout="\n".join(all_stdout),
        stderr=err_msg,
        agent_concluded_fail=True,
    )


# ── Outlook 自動化節點專屬 agent ──────────────────────────────────────────────
# 跟 skill 模式類似（LLM 寫 code、解析 tool calls、迴圈），但有以下差異：
#   1. 強制 host 執行（pywin32 在 sandbox / Linux 容器跑不了）
#   2. AST allowlist 檢查（每次 run_python 前用 win32_agent_config.check_imports 過濾）
#   3. 系統提示限定 win32_helpers + 允許的套件，做不到的需求要 agent 直接 done(success=false)
#   4. 不接 recipe 快取（Outlook 環境每次狀態不同 — 收件匣會新增/刪信，沒辦法穩定 replay）
#   5. 較簡單的 tool 集（run_python / done / ask_user，不要 web_search / view_image）

OUTLOOK_AGENT_MAX_ITERATIONS = 12
OUTLOOK_AGENT_REQUEST_INTERVAL = 2.0


_OUTLOOK_SYSTEM_PROMPT = """你是 Outlook 自動化專家。透過 pywin32 + Outlook COM 處理寄信、收信、行事曆、附件等需求。

## 環境限制（嚴格）

你只能使用以下套件：
- **win32_helpers.outlook**：本專案 wrapper（最推薦，用這個就好）
- **win32com.client / pywintypes / pythoncom**：原始 COM（罕見場景才用）
- **pandas / numpy / openpyxl**：資料整理 / 寫 xlsx
- **python-docx / python-pptx**：產生 docx / pptx 報告
- **bs4 / jinja2 / markdown**：HTML / 模板 / md 渲染
- **PIL（Pillow）**：圖片處理
- **標準庫**：re / datetime / pathlib / json / csv / html / email 等

**禁止 import**：requests / httpx / urllib / selenium / playwright / subprocess / smtplib / imaplib / sklearn / torch 等。
做不到的需求（例如要連 Web API、操作 Slack、Teams、瀏覽器）→ 直接 `done(success=false, error="此需求需要 X，不在 Outlook 自動化節點範圍。建議使用一般 Skill 節點。")`。

## 推薦的 wrapper API（import 路徑：`from pipeline.win32_helpers.outlook import ...`）

注意：本節點的 sys.path 會自動注入 backend dir，所以你直接 `from pipeline.win32_helpers.outlook import search_mail` 就能 import。**不要寫 `from win32_helpers.outlook import ...`**（少了 `pipeline.` 前綴會 ModuleNotFoundError）。

```python
# 讀信（回 DataFrame）
search_mail(*, subject=None, sender=None, body_keyword=None,
            since=None, until=None, folder="inbox",
            unread_only=False, has_attachment=None,
            exact_match=False, limit=500) -> pd.DataFrame
# 欄位：entry_id / received / sender_name / sender_email / subject /
#       body_preview / body_text / has_attachments / attachment_names /
#       is_unread / importance / folder_name

get_mail_by_id(entry_id) -> dict          # 單封信完整資料（含 body_html）
download_attachments(*, entry_ids, out_dir, name_template="...") -> list[Path]

# 寄信
send_mail(*, to, subject, body, body_format="html",
          cc=None, bcc=None, attachments=None,
          importance=1, save_to_drafts=False) -> str  # EntryID
reply_mail(*, entry_id, body, body_format="html", reply_all=False) -> str
forward_mail(*, entry_id, to, body="", body_format="html") -> str

# 行事曆
calendar_list(*, since=None, until=None, folder="calendar",
              include_recurring=True, limit=200) -> pd.DataFrame
create_meeting(*, subject, start, end, location="", body="",
               required_attendees=None, optional_attendees=None,
               reminder_minutes=15, send_invitation=True) -> str
```

## 工具（原生 function calling）

本對話**使用原生 function calling API**、工具已透過 function_declarations 註冊。你**必須**透過 API 的 function_call 機制呼叫工具、**禁止**寫 `<tool>name</tool>` / `<input>...</input>` 這類文字格式 — 那種純文字 orchestrator 不會解析、會被視為沒呼叫任何工具。直接 emit 結構化 tool_calls 即可。

可用工具：
- **run_python(code)**：執行 Python 程式碼處理 Outlook 需求。
  例：
  ```python
  import pandas as pd
  from pipeline.win32_helpers.outlook import search_mail
  df = search_mail(subject="報告", since="2026-04-25")
  print(df.head())
  ```
- **done(success, summary)**：結束任務。success=true/false、summary 用中文說明結果（失敗時說明卡在哪）。
- **ask_user(question, options, context)**：問使用者（極少用，例如「找到 50 封信，要不要全部處理？」）。

每次 reply 呼叫**一個**工具即可。

## 規則

1. 寫 Python 程式碼前先思考一下整體流程（不要一次寫太短的 read_csv 然後再 print）
2. 整理結果如果指定了 output_path，**一定要把結果存到那個路徑**（xlsx / md / json 等格式由情境決定）
3. 出錯後不要直接 done(success=true)！先用 run_python 修錯、確認 stdout 沒有 traceback 才 done
4. 永遠不要寫 fake stdout（不要在 <input> 後面寫『Successfully sent.』『DataFrame: ...』之類字串）— 真實結果系統會回給你
5. 如果使用者描述太模糊、缺關鍵資訊（例如要寄給誰、日期區間），用 ask_user 問；不要瞎猜亂寄信

## 進階：建立資料夾 / 收件規則（wrapper 沒有、用原始 win32com COM）

- **新增資料夾**（可重複跑、已存在就沿用、不要報錯）：
  ```python
  import win32com.client as w
  ns = w.Dispatch("Outlook.Application").GetNamespace("MAPI")
  inbox = ns.GetDefaultFolder(6)  # olFolderInbox
  name = "我的資料夾"
  existing = {inbox.Folders.Item(i + 1).Name for i in range(inbox.Folders.Count)}
  folder = inbox.Folders[name] if name in existing else inbox.Folders.Add(name)
  ```
- **建立收件規則（自動分流「未來」的信）** — ⚠️ 規則 API **需要 Exchange 在線**：
  ```python
  try:
      rules = ns.DefaultStore.GetRules()      # 離線 / 未連 Exchange 會丟例外
  except Exception as e:
      # ↓ 不要假裝成功！老實回報需要上線
      raise RuntimeError(f"Outlook 離線或未連 Exchange、無法建規則：{e}")
  r = rules.Create("規則名稱", 0)              # 0 = olRuleReceive（收件時觸發）
  c = r.Conditions.SenderAddress; c.Enabled = True; c.Address = ["someone@example.com"]
  a = r.Actions.MoveToFolder;     a.Enabled = True; a.Folder = folder   # 上面建/取到的資料夾物件
  rules.Save()
  ```
- **建完規則務必「重新讀回」驗證**：再 `ns.DefaultStore.GetRules()`、確認找得到該規則名才算成功；找不到 → `done(success=false)`。**別憑「沒丟例外」就說成功。**

6. ⚠️ **離線 / 連不上 Exchange 的硬規則（最重要、防假成功）**：凡是需要「伺服器端生效」的操作（建/改收件規則、伺服器端搬大量信、寄信）一旦遇到「離線工作 / 必須連線 Microsoft Exchange」這類 COM 例外 → **立刻 `done(success=false)` 說明「Outlook 目前離線、此操作需先連上 Exchange」**，**絕對不要回報成功**。建資料夾通常本機就能成功、不受此限。"""


def _build_outlook_prompt(
    *,
    template: str,
    template_params: dict,
    free_text: str,
    output_path: Optional[str],
    prev_outputs: Optional[list],
    prefetched_data: str = "",
    prefetch_error: str = "",
) -> str:
    """根據 template / params / 自由輸入文字組出給 LLM 的 user prompt。

    若 prefetched_data 有值，prompt 會明確告訴 LLM「資料已預抓」並把它嵌進來，
    LLM 的工作就只剩「整理 + 寫檔」。"""
    parts: list[str] = []

    if template:
        parts.append(f"## 任務模板：{template}")
        parts.append("使用者透過選單選了這個模板，需求參數如下：")
        if template_params:
            for k, v in template_params.items():
                if v == "" or v is None or (isinstance(v, list) and not v):
                    continue
                parts.append(f"  - **{k}**：{v}")
        else:
            parts.append("  （未填參數，請依模板預設行為執行）")
    elif free_text:
        parts.append("## 自由輸入需求")
        parts.append(free_text)
    else:
        parts.append("## 任務未明確設定")
        parts.append("使用者既沒選模板也沒打字。請呼叫 done(success=false, error=...) "
                     "回報需要更明確的指示。")

    # 預抓資料區塊 — 若後端已預抓成功，LLM 不應再呼叫 search_mail
    if prefetched_data:
        parts.append("")
        parts.append("## 已預抓的資料（後端已從 Outlook 抓好）")
        parts.append("**這段是真實 Outlook 內容，已經完成抓取。你不需要也不應該重新呼叫 "
                     "search_mail/calendar_list 等函式抓資料。**")
        parts.append("你的工作：根據下方資料 + 模板參數，整理成適合的格式（通常是 markdown）"
                     "寫到 output_path。")
        parts.append("")
        parts.append("---")
        parts.append(prefetched_data)
        parts.append("---")
    elif prefetch_error:
        # prefetch 失敗了 → 告訴 LLM 自己想辦法
        parts.append("")
        parts.append(f"## 注意：預抓資料失敗（{prefetch_error}）")
        parts.append("後端嘗試自動抓資料但失敗了。請你自己用 win32_helpers / pywin32 "
                     "嘗試把資料抓出來、處理。")

    if output_path:
        parts.append("")
        parts.append(f"## 輸出檔案路徑")
        parts.append(f"請把整理 / 摘要結果寫到：`{output_path}`")
        # 從副檔名告訴 LLM 該怎麼寫
        _suffix = Path(output_path).suffix.lower()
        if _suffix == ".md":
            parts.append("**副檔名 .md → 寫 markdown 純文字**：用 `Path(...).write_text(report_str, encoding='utf-8')`")
        elif _suffix == ".docx":
            parts.append("**副檔名 .docx → 用 python-docx 寫 Word 檔**：")
            parts.append("```python")
            parts.append("from docx import Document")
            parts.append("doc = Document()")
            parts.append("doc.add_heading('標題', 0)")
            parts.append("doc.add_paragraph('摘要內容...')")
            parts.append(f"doc.save(r'{output_path}')")
            parts.append("```")
        elif _suffix == ".xlsx":
            parts.append("**副檔名 .xlsx → 用 openpyxl 或 pandas.to_excel 寫 Excel 檔**")
        elif _suffix == ".pdf":
            parts.append("**副檔名 .pdf**：先寫 docx 再用 docx2pdf 轉、或用 reportlab")
        elif _suffix == ".json":
            parts.append("**副檔名 .json → 用 json.dump 寫結構化資料**")
        else:
            parts.append("（父資料夾已建好。請依副檔名決定寫法 — md / docx / xlsx / json 等）")
        if prefetched_data and _suffix == ".md":
            parts.append("典型寫法（資料已備齊、不用再抓）：")
            parts.append("```python")
            parts.append("from pathlib import Path")
            parts.append("report = '''# 標題\\n\\n（你的摘要內容）\\n'''")
            parts.append(f"Path(r'{output_path}').write_text(report, encoding='utf-8')")
            parts.append("```")

    if prev_outputs:
        parts.append("")
        parts.append("## 前一步驟的輸出檔案（可讀取）")
        for o in prev_outputs:
            p = o.get("path", "")
            schema = o.get("schema", "")
            parts.append(f"  - `{p}`" + (f"：{schema}" if schema else ""))

    parts.append("")
    parts.append("請開始執行。完成後呼叫 done(success=true, summary='...')。")
    return "\n".join(parts)


async def execute_step_with_outlook(
    *,
    template: str,
    template_params: dict,
    free_text: str,
    timeout: int,
    logger: logging.Logger,
    step_name: str,
    output_path: Optional[str] = None,
    working_dir: Optional[str] = None,
    prev_outputs: Optional[list] = None,
    run_id: str = "",
    ask_mode: bool = False,
    llm_role: str = "primary",
) -> ExecResult:
    """Outlook 自動化節點。

    執行路徑分流：
      A. Direct path（template 在 outlook_templates.DIRECT_HANDLERS 裡）→ 不走 LLM、
         直接呼叫對應 handler（快、確定性、零 token 成本）
      B. LLM path（template 不在 direct 清單 / 為空 / 自由輸入）→ 進 agent loop、
         讓 LLM 寫 code 處理（吃 token、可能多輪 retry）
    """
    # 處理 output_path（兩條路徑共用）
    if output_path:
        output_path = str(Path(output_path).expanduser())
        # 若 template_params 帶 output_format，把路徑副檔名同步調整
        # 例：runner default 給 .md 但使用者選 docx → 改成 .docx
        # 這樣 LLM / direct handler 都直接拿到正確副檔名的路徑、不用自己判斷
        _fmt = (template_params or {}).get("output_format") or ""
        _fmt = str(_fmt).strip().lower().lstrip(".")
        _ext_map = {"md": ".md", "markdown": ".md", "xlsx": ".xlsx", "excel": ".xlsx",
                    "txt": ".txt", "docx": ".docx", "word": ".docx",
                    "pdf": ".pdf", "json": ".json", "csv": ".csv"}
        if _fmt in _ext_map:
            desired = _ext_map[_fmt]
            if not output_path.lower().endswith(desired):
                old = output_path
                output_path = str(Path(output_path).with_suffix(desired))
                logger.info(f"[{step_name}] 依 output_format={_fmt} 調整 output_path：{old} → {output_path}")
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        _out = Path(output_path)
        if _out.exists():
            try:
                _out.unlink()
                logger.info(f"[{step_name}] 刪除舊輸出檔：{output_path}")
            except Exception as e:
                logger.warning(f"[{step_name}] 舊輸出檔刪除失敗（可能被開啟中）：{e}")
    if working_dir:
        Path(working_dir).mkdir(parents=True, exist_ok=True)

    # ── 路徑 A：Direct template handler（不需 LLM）────────────────────
    from pipeline.outlook_templates import is_direct_template, run_direct_template
    if template and is_direct_template(template):
        logger.info(f"[{step_name}] 走 direct 模板路徑：{template}（不進 LLM）")
        # 重要：run_direct_template 內部呼叫 search_mail / calendar_list 都是
        # 同步阻塞 COM call，可能跑 4-5 分鐘。直接在 asyncio loop 跑會 block
        # 整個 FastAPI、frontend 1.5s 的 log polling 全部排隊到 run 結束才回應，
        # 看起來像「跑完才一次印出 log」。所以丟去 thread pool。
        ok, summary = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: run_direct_template(
                template=template,
                params=template_params or {},
                output_path=output_path or "",
                prev_outputs=prev_outputs,
                step_name=step_name,
                logger_obj=logger,
            ),
        )
        return ExecResult(
            exit_code=0 if ok else 1,
            stdout=f"[Outlook 完成] {summary}" if ok else "",
            stderr="" if ok else summary,
        )

    # ── 路徑 B：LLM agent loop（自由輸入 / 需要摘要的模板）─────────────
    # AST gate
    backend_dir = str(Path(__file__).resolve().parent.parent)
    if backend_dir not in sys.path:
        sys.path.insert(0, backend_dir)
    from pipeline.win32_agent_config import check_imports, format_errors_for_agent

    # 給下方 run_python 的 sys.path 注入用 — agent code 跑在 subprocess、不繼承我們這邊的 sys.path
    _backend_dir_for_outlook = backend_dir

    # 預抓資料（mid-path）：若 template 有 prefetch handler，後端先把資料抓好給 LLM。
    # LLM 的工作從「寫 search_mail 程式碼 + 摘要」變成「讀資料 + 摘要」 — 大幅縮短迭代、
    # 避開 LLM 寫 search_mail 程式碼可能踩的 timezone / pandas / 套件坑。
    prefetched_data = ""
    prefetch_error = ""
    from pipeline.outlook_templates import has_prefetch, run_prefetch
    if template and has_prefetch(template):
        # 同 direct path：run_prefetch 是同步 COM call，丟 thread pool 避免 block event loop。
        # 期間 frontend 的 log polling 才能即時收到 search_mail 的進度訊息。
        ok_pf, md_pf, err_pf = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: run_prefetch(
                template=template, params=template_params or {},
                prev_outputs=prev_outputs, logger_obj=logger,
            ),
        )
        if ok_pf:
            prefetched_data = md_pf
            logger.info(f"[{step_name}] ✓ 預抓資料完成（{len(md_pf)} 字），LLM 只需摘要")
        else:
            prefetch_error = err_pf
            logger.warning(f"[{step_name}] ⚠ 預抓失敗：{err_pf}（LLM 將自己抓）")

    # LLM
    from llm_factory import build_llm
    from langchain_core.messages import HumanMessage, SystemMessage, AIMessage, ToolMessage
    from langchain_core.tools import tool as _lc_tool
    llm = build_llm(role=llm_role)

    user_prompt = _build_outlook_prompt(
        template=template,
        template_params=template_params or {},
        free_text=free_text,
        output_path=output_path,
        prev_outputs=prev_outputs,
        prefetched_data=prefetched_data,
        prefetch_error=prefetch_error,
    )
    logger.info(f"[{step_name}] Outlook agent 啟動，template={template or '(自由輸入)'}"
                + ("（資料已預抓）" if prefetched_data else ""))
    logger.debug(f"[{step_name}] user prompt:\n{user_prompt}")

    messages = [SystemMessage(content=_OUTLOOK_SYSTEM_PROMPT), HumanMessage(content=user_prompt)]
    all_stdout: list[str] = []
    last_run_python_ok: Optional[bool] = None
    import time as _time
    start_time = _time.time()

    # ── 3 個工具（native function calling）─────────────────────────────
    # 這些 @tool wrapper 只負責「對 LLM 暴露乾淨 schema」、實際執行（AST 檢查 /
    # sys.path 注入 / host 執行 / done 守門 / ask_user 等待）仍在下方 loop 處理，
    # loop 看 tc["name"] 自己 dispatch（跟 validator 的 view_image/done 同模式）。
    @_lc_tool
    def run_python(code: str) -> str:
        """執行 Python 程式碼來處理 Outlook 需求（寄信 / 讀信 / 行事曆 / 寫檔）。

        Args:
            code: 要執行的 Python 程式碼（完整可獨立執行；用 from pipeline.win32_helpers.outlook import ...）
        Returns:
            stdout + stderr（如有）
        """
        return "__RUN_PYTHON__"  # 不在此真執行、loop 看 tool name 做 AST 檢查 + host 執行

    @_lc_tool
    def done(success: bool, summary: str = "") -> str:
        """結束任務、回傳最終結果。

        Args:
            success: 是否成功完成
            summary: 結果說明（中文；失敗時說明卡在哪）
        Returns:
            結束標記
        """
        return "__DONE__"

    @_lc_tool
    def ask_user(question: str, context: str = "") -> str:
        """需求模糊 / 缺關鍵資訊時問使用者（極少用）。

        Args:
            question: 要問的問題(若有多選請直接寫在問題內,例「A/B/C 哪個?」)
            context: 補充情境（可省略）
        Returns:
            使用者的回答
        """
        return "__ASK_USER__"  # 不在此真執行、loop 看 tool name 等待使用者回答
        # 註:不暴露 options 參數 — Gemini function declaration 對 Optional[list](anyOf/無 items)
        # 會回 400 INVALID_ARGUMENT(E2E 抓到)。對齊 sandbox_tools.ask_user 的 gemma-safe schema。

    outlook_tools = [run_python, done, ask_user]
    name_to_tool = {t.name: t for t in outlook_tools}

    try:
        llm_with_tools = llm.bind_tools(outlook_tools)
    except Exception as _be:
        logger.error(f"[{step_name}] Outlook bind_tools 失敗：{_be}")
        return ExecResult(
            exit_code=-1, stdout="\n".join(all_stdout) or "",
            stderr=f"Outlook bind_tools 失敗：{_be}",
        )

    consecutive_empty = 0  # 連續空回應守門：gemma native-FC config 下可能 tool_calls + content 皆空 → 早退、別空轉

    try:
        for iteration in range(OUTLOOK_AGENT_MAX_ITERATIONS):
            logger.info(f"[{step_name}] Outlook 迭代 {iteration + 1}/{OUTLOOK_AGENT_MAX_ITERATIONS} [NATIVE]")
            if _time.time() - start_time > timeout:
                logger.warning(f"[{step_name}] Outlook agent 整體 timeout（{timeout}s）達到")
                return ExecResult(
                    exit_code=-2, stdout="\n".join(all_stdout),
                    stderr=f"Outlook agent timeout（{timeout}s）",
                )

            if iteration > 0:
                await asyncio.sleep(OUTLOOK_AGENT_REQUEST_INTERVAL)

            response: AIMessage = await asyncio.wait_for(
                llm_with_tools.ainvoke(messages), timeout=180.0
            )
            tool_calls = list(getattr(response, "tool_calls", []) or [])
            content_str = _extract_text(response.content)
            logger.info(
                f"[{step_name}] LLM 回覆（content {len(content_str)} 字, tool_calls={len(tool_calls)}）"
            )
            if content_str:
                logger.debug(f"[{step_name}] LLM content：{content_str[:1500]}")

            # AIMessage 進歷史（含 tool_calls、維持 tool_call/tool_result 配對）
            messages.append(response)

            if not tool_calls:
                # 連續空回應守門：tool_calls + content 皆空（content<5 字）→ streak++、連 2 次中止。
                # gemma 在 native-FC config 下若不吐 tool 也不吐字會卡住，別空轉到上限。
                if len(content_str.strip()) < 5:
                    consecutive_empty += 1
                    logger.warning(
                        f"[{step_name}] ⚠ 第 {iteration + 1} 輪空回應"
                        f"（tool_calls + content 皆空）、累計 {consecutive_empty}/2"
                    )
                    if consecutive_empty >= 2:
                        err_msg = (
                            "Outlook agent 連續 2 輪空回應（tool_calls + content 皆空、"
                            "疑似模型不相容 native function calling）、判定卡住、中止。"
                        )
                        logger.error(f"[{step_name}] ✗ {err_msg}")
                        return ExecResult(
                            exit_code=1, stdout="\n".join(all_stdout), stderr=err_msg,
                        )
                else:
                    consecutive_empty = 0
                messages.append(HumanMessage(content=
                    "請使用工具（原生 function calling）呼叫 run_python / done / ask_user。"
                    "不要用文字描述工具呼叫。"))
                continue

            consecutive_empty = 0

            # 掃出 done（若有）、其餘工具先執行；done 在最後處理
            done_call: Optional[dict] = None
            regular_calls = []
            for tc in tool_calls:
                if tc.get("name") == "done":
                    done_call = tc  # 多個 done 取最後一個
                else:
                    regular_calls.append(tc)

            for tc in regular_calls:
                tool_name = tc.get("name", "")
                tc_id = tc.get("id") or ""
                tc_args = tc.get("args", {}) or {}
                logger.info(f"[{step_name}] tool={tool_name} [NATIVE]")

                if tool_name == "run_python":
                    tool_input = str(tc_args.get("code", "") or "")
                    logger.info(f"[{step_name}] run_python input_len={len(tool_input)}")
                    # AST 檢查 — disallowed import 直接擋
                    try:
                        errs = check_imports(tool_input)
                    except SyntaxError as e:
                        messages.append(ToolMessage(
                            content=f"[系統] 你提交的 Python 有語法錯誤：{e}", tool_call_id=tc_id,
                        ))
                        last_run_python_ok = False
                        continue
                    if errs:
                        err_msg = format_errors_for_agent(errs)
                        logger.warning(f"[{step_name}] 偵測到 disallowed imports：{[e.module for e in errs]}")
                        messages.append(ToolMessage(content=err_msg, tool_call_id=tc_id))
                        last_run_python_ok = False
                        continue

                    # 執行（強制 host）
                    # 注入 sys.path 讓 subprocess 找得到 backend dir 裡的 win32_helpers / pipeline 套件
                    # （_skill_run_python 寫到 Windows temp 目錄後 spawn subprocess，預設 sys.path
                    # 沒有 backend dir → import win32_helpers / from pipeline.X 會 ModuleNotFoundError）
                    _injected_code = (
                        f"import sys\n"
                        f"sys.path.insert(0, r{repr(_backend_dir_for_outlook)})\n"
                        + tool_input
                    )
                    tool_result = await asyncio.get_event_loop().run_in_executor(
                        None,
                        lambda ti=_injected_code, lg=logger, tt=_compute_tool_timeout(timeout): _execute_skill_tool(
                            "run_python", ti, cwd=working_dir, run_id=run_id,
                            logger=lg, force_host=True, tool_timeout=tt,
                        ),
                    )
                    logger.debug(f"[{step_name}] 執行結果：{tool_result[:1500]}")
                    all_stdout.append(f"[run_python] {tool_result}")
                    last_run_python_ok = "[exit code:" not in tool_result
                    messages.append(ToolMessage(content=tool_result, tool_call_id=tc_id))
                    continue

                if tool_name == "ask_user":
                    question = str(tc_args.get("question", "") or "")
                    options = tc_args.get("options") or []
                    context = str(tc_args.get("context", "") or "")
                    answer = await _wait_for_ask_user(
                        run_id, question, options, context, logger, step_name,
                    )
                    all_stdout.append(f"[ask_user] {question} → {answer}")
                    messages.append(ToolMessage(content=f"[ask_user 答案] {answer}", tool_call_id=tc_id))
                    continue

                # 不支援的工具（理論上 bind_tools 後不會發生）
                messages.append(ToolMessage(
                    content=(f"[系統] 工具 {tool_name} 不在 Outlook 節點允許清單。"
                             f"可用：run_python / done / ask_user。"),
                    tool_call_id=tc_id,
                ))

            # done 工具（在 regular_calls 全跑完後處理）
            if done_call is not None:
                tc_id = done_call.get("id") or ""
                tc_args = done_call.get("args", {}) or {}
                success = bool(tc_args.get("success", False))
                summary = tc_args.get("summary") or tc_args.get("error") or ""

                # 守門：宣稱成功但最近 run_python 失敗 → 拒絕
                if success and last_run_python_ok is False:
                    logger.warning(f"[{step_name}] 在 run_python 失敗後送 done(success=true)，拒絕")
                    messages.append(ToolMessage(
                        content="[系統] 拒絕 done：上一次 run_python 失敗。先修錯再 done。",
                        tool_call_id=tc_id,
                    ))
                    continue
                # 守門：宣稱成功但 output 檔不存在
                if success and output_path and not Path(output_path).exists():
                    logger.warning(f"[{step_name}] done 宣稱成功但 {output_path} 不存在，拒絕")
                    messages.append(ToolMessage(
                        content=(f"[系統] 你宣稱成功但輸出檔 {output_path} 不存在。"
                                 f"請用 run_python 實際寫入後再 done。"),
                        tool_call_id=tc_id,
                    ))
                    continue

                all_stdout.append(f"[Outlook 完成] {summary}")
                logger.info(f"[{step_name}] {'成功' if success else '失敗'}：{summary}")
                final_stdout = (_build_clean_success_stdout(all_stdout, "[Outlook 完成]")
                                if success else "\n".join(all_stdout))
                return ExecResult(
                    exit_code=0 if success else 1,
                    stdout=final_stdout,
                    stderr="" if success else summary,
                )

        # 達到迭代上限
        logger.warning(f"[{step_name}] Outlook agent 達迭代上限 {OUTLOOK_AGENT_MAX_ITERATIONS}")
        return ExecResult(
            exit_code=-2, stdout="\n".join(all_stdout),
            stderr=f"Outlook agent 在 {OUTLOOK_AGENT_MAX_ITERATIONS} 次內未完成",
        )

    except Exception as e:
        _err_str = str(e)
        is_quota = ("429" in _err_str or "RESOURCE_EXHAUSTED" in _err_str
                    or "quota" in _err_str.lower() or "rate limit" in _err_str.lower())
        if is_quota:
            logger.error(f"[{step_name}] Outlook agent LLM 配額：{_err_str[:200]}")
            return ExecResult(
                exit_code=-429, stdout="\n".join(all_stdout),
                stderr=f"LLM provider 配額用盡或速率受限（429）：{_err_str}",
            )
        logger.error(f"[{step_name}] Outlook agent 例外：{e}", exc_info=True)
        return ExecResult(
            exit_code=-3, stdout="\n".join(all_stdout),
            stderr=f"Outlook agent 例外：{e}",
        )


# ── 網頁爬蟲節點（V5 新增）──────────────────────────────────────────

async def execute_step_with_web_crawler(
    *,
    mode: str,
    # web 模式
    url: str,
    urls: list[str],            # 多 URL 列表；非空時優先用、走 crawl_urls
    js_render: bool,
    wait_for_selector: str,
    cloudflare_fallback: bool,
    cookies: str,
    interactions: list[dict],
    download_assets: bool,
    scroll_count: int = 0,          # 0=自動偵測 / >0=固定滾 N 次
    target_post_count: int = 0,     # 0=不設目標 / >0=滾到至少 N 個貼文連結
    # 論壇 / 列表模式
    with_children: bool = False,           # True = 列表頁 + 抽前 N 子頁全部抓回合併
    child_link_pattern: str = "",          # 子頁 URL pattern（空 = auto 內建 12 種）
    max_children: int = 10,                # 最多抓幾個子頁
    # 影片模式
    video_url: str,
    video_quality: str,
    video_max_filesize_mb: int,
    video_max_duration_min: int,
    video_subs: bool,
    video_subs_langs: str,
    video_save_info_json: bool,
    # 共用
    output_path: str,
    timeout: int,
    logger: logging.Logger,
    step_name: str,
) -> ExecResult:
    """網頁爬蟲節點。mode 決定走哪條路徑：
      "web"   → Crawl4AI（網頁 → markdown）；CF 偵測到時 fallback FlareSolverr
      "video" → yt-dlp（影音站 → mp4 + 字幕 + .info.json + 摘要 .md）
    """
    if not output_path:
        return ExecResult(exit_code=2, stdout="",
                          stderr="web_crawler 節點未填 output.path")

    if mode == "video":
        if not video_url:
            return ExecResult(exit_code=2, stdout="",
                              stderr="web_crawler 影片模式未填 YouTube/Vimeo URL")
        from pipeline.web_crawler import crawl_video
        try:
            result = await crawl_video(
                url=video_url,
                output_path=output_path,
                quality=video_quality or "720p",
                max_filesize_mb=video_max_filesize_mb or 500,
                max_duration_min=video_max_duration_min if video_max_duration_min is not None else 30,
                subs=video_subs,
                subs_langs=video_subs_langs or "",
                save_info_json=video_save_info_json,
                cookies=cookies,
                timeout=timeout,
                logger=logger,
                step_name=step_name,
            )
        except Exception as e:
            logger.error(f"[{step_name}] 爬蟲例外：{e}", exc_info=True)
            return ExecResult(exit_code=-3, stdout="", stderr=f"爬蟲例外：{e}")
    else:
        # ── 論壇 / 列表模式（with_children）：自動抓子頁、合併單一 markdown ──
        # 取代之前「使用者拉 skill 節點讓 LLM 寫 crawl4ai code 抓 N 篇」的脆弱方案。
        # 必須是單 URL 才適用（列表頁本身）；多 URL 模式跟這個互斥
        if with_children:
            single_url = url or (
                [u.strip() for u in (urls or []) if u and u.strip()][0]
                if (urls and any(u.strip() for u in urls)) else ""
            )
            if not single_url:
                return ExecResult(exit_code=2, stdout="",
                                  stderr="web_crawler 論壇模式（with_children）未填列表頁 URL")
            from pipeline.web_crawler import crawl_list_with_children
            try:
                result = await crawl_list_with_children(
                    list_url=single_url,
                    output_path=output_path,
                    js_render=js_render,
                    wait_for_selector=wait_for_selector,
                    cloudflare_fallback=cloudflare_fallback,
                    cookies=cookies,
                    interactions=interactions or [],
                    download_assets=download_assets,
                    scroll_count=scroll_count,
                    target_post_count=target_post_count,
                    child_link_pattern=child_link_pattern,
                    max_children=max_children,
                    timeout=timeout,
                    logger=logger,
                    step_name=step_name,
                )
            except Exception as e:
                logger.error(f"[{step_name}] 論壇模式爬蟲例外：{e}", exc_info=True)
                return ExecResult(exit_code=-3, stdout="", stderr=f"爬蟲例外：{e}")
            if not result.ok:
                return ExecResult(exit_code=1, stdout="",
                                  stderr=f"列表頁爬取失敗（tier={result.tier}）：{result.error}")
            summary = (
                f"[爬蟲完成] 論壇模式 tier={result.tier} 列表 + 子頁、"
                f"合併 {len(result.markdown.split())} 字 耗時={result.duration_ms}ms → {output_path}"
            )
            return ExecResult(exit_code=0, stdout=summary, stderr="")

        # 預設 web 模式 — 多 URL 還是單 URL？
        # 過濾空行 / 註解後 list 多於 1 → 多 URL；其他都單 URL（含只填 wc_url）
        cleaned_urls = [u.strip() for u in (urls or []) if u and u.strip() and not u.strip().startswith("#")]
        if cleaned_urls and len(cleaned_urls) > 1:
            # ── 多 URL 路徑 ──
            from pipeline.web_crawler import crawl_urls
            from pathlib import Path as _Path
            # output_path 解讀為「資料夾」(去掉副檔名 / 拿 parent)
            outp = _Path(output_path)
            output_dir = str(outp.parent if outp.suffix else outp)
            try:
                manifest = await crawl_urls(
                    urls=cleaned_urls,
                    output_dir=output_dir,
                    js_render=js_render,
                    wait_for_selector=wait_for_selector,
                    cloudflare_fallback=cloudflare_fallback,
                    cookies=cookies,
                    interactions=interactions or [],
                    download_assets=download_assets,
                    scroll_count=scroll_count,
                    target_post_count=target_post_count,
                    timeout=timeout,
                    logger=logger,
                    step_name=step_name,
                )
            except Exception as e:
                logger.error(f"[{step_name}] 多 URL 爬蟲例外：{e}", exc_info=True)
                return ExecResult(exit_code=-3, stdout="", stderr=f"爬蟲例外：{e}")
            if not manifest.get("ok"):
                return ExecResult(
                    exit_code=1, stdout="",
                    stderr=f"多 URL 爬取失敗：{manifest.get('successful')}/{manifest.get('total')} 成功；"
                           f"{manifest.get('error') or '見 index.json'}",
                )
            summary = (
                f"[爬蟲完成] 多 URL：{manifest['successful']}/{manifest['total']} 成功 "
                f"耗時={manifest.get('duration_ms', 0)}ms → {output_dir}/index.json"
            )
            return ExecResult(exit_code=0, stdout=summary, stderr="")

        # ── 單 URL 路徑（向後相容：用 wc_url 或 wc_urls 只有一個）──
        single_url = cleaned_urls[0] if cleaned_urls else url
        if not single_url:
            return ExecResult(exit_code=2, stdout="",
                              stderr="web_crawler 網頁模式未填 URL")
        from pipeline.web_crawler import crawl_single_url
        try:
            result = await crawl_single_url(
                url=single_url,
                output_path=output_path,
                js_render=js_render,
                wait_for_selector=wait_for_selector,
                cloudflare_fallback=cloudflare_fallback,
                cookies=cookies,
                interactions=interactions or [],
                download_assets=download_assets,
                scroll_count=scroll_count,
                target_post_count=target_post_count,
                timeout=timeout,
                logger=logger,
                step_name=step_name,
            )
        except Exception as e:
            logger.error(f"[{step_name}] 爬蟲例外：{e}", exc_info=True)
            return ExecResult(exit_code=-3, stdout="", stderr=f"爬蟲例外：{e}")

    if not result.ok:
        return ExecResult(exit_code=1, stdout="",
                          stderr=f"爬取失敗（tier={result.tier}）：{result.error}")

    # 404 fail-fast:crawl 本身「成功」回傳、但抓到的是 HTTP 4xx/5xx 錯誤頁 或 內容過短的空殼
    # (JS 未渲染 / 反爬擋頁)→ 直接判 fail,別把錯誤頁當成功往下傳、害下游 parser 解析空殼才爆。
    # (對齊 crawl_empty_diagnosis:先看 status + word_count;404 頁/空殼不是有效內容)
    _wc = len(result.markdown.split())
    _sc = getattr(result, "status_code", None)
    if isinstance(_sc, int) and _sc >= 400:
        return ExecResult(exit_code=1, stdout="",
            stderr=(f"爬到 HTTP {_sc} 錯誤頁（tier={result.tier}、字數={_wc}）→ "
                    f"URL 可能失效 / 改版 / 反爬,請確認網址正確。不把錯誤頁當成功往下傳。"))
    if _wc < 30:
        return ExecResult(exit_code=1, stdout="",
            stderr=(f"爬取內容過少（僅 {_wc} 字、status={_sc}、tier={result.tier}）→ "
                    f"可能是空殼 / JS 未渲染 / 反爬擋頁,非有效內容,無法供下游解析。"))

    summary = (
        f"[爬蟲完成] tier={result.tier} status={result.status_code} "
        f"title={result.title!r} 字數={_wc} "
        f"耗時={result.duration_ms}ms → {output_path}"
    )
    return ExecResult(exit_code=0, stdout=summary, stderr="")
